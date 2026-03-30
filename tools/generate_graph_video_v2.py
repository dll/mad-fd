from __future__ import annotations

import math
import os
import re
import shutil
import struct
import subprocess
import sys
import textwrap
import wave
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional, Tuple

from PIL import Image, ImageDraw, ImageFont

try:
    import pyttsx3
except Exception:
    pyttsx3 = None

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
DIAGRAMS_DIR = DOCS_DIR / "diagrams"
VIDEO_DIR = DOCS_DIR / "video"
GENERATED_DIR = VIDEO_DIR / "generated"
SLIDES_DIR = VIDEO_DIR / "slides_v2"
AUDIO_DIR = GENERATED_DIR / "audio_v2"
CLIPS_DIR = GENERATED_DIR / "clips_v2"
VIDEO_OUTPUT_DIR = ROOT / "video_output"

SCRIPT_MD_PATH = VIDEO_DIR / "video_script_v2.md"
PPTX_PATH = VIDEO_OUTPUT_DIR / "知识图谱核心功能_图谱功能教程_v2.pptx"
VIDEO_PATH = VIDEO_OUTPUT_DIR / "知识图谱核心功能_图谱功能教程_v2.mp4"

ARCH_PUML = DIAGRAMS_DIR / "knowledge_graph_feature_architecture.puml"
FLOW_PUML = DIAGRAMS_DIR / "knowledge_graph_feature_flow.puml"
ARCH_PNG = DIAGRAMS_DIR / "knowledge_graph_feature_architecture.png"
FLOW_PNG = DIAGRAMS_DIR / "knowledge_graph_feature_flow.png"
DATA_PNG = DIAGRAMS_DIR / "knowledge_graph_feature_data_model.png"

MAIN_PATH = ROOT / "lib" / "main.dart"
HOME_PATH = ROOT / "lib" / "presentation" / "pages" / "home" / "home_page.dart"
GRAPH_LIST_PATH = (
    ROOT / "lib" / "presentation" / "pages" / "graph" / "graph_list_page.dart"
)
GRAPH_DETAIL_PATH = (
    ROOT / "lib" / "presentation" / "pages" / "graph" / "graph_detail_page.dart"
)
DATABASE_HELPER_PATH = ROOT / "lib" / "data" / "local" / "database_helper.dart"
GRAPH_DAO_PATH = ROOT / "lib" / "data" / "local" / "graph_dao.dart"
FAVORITE_DAO_PATH = ROOT / "lib" / "data" / "local" / "favorite_dao.dart"
LEARNING_RECORD_DAO_PATH = ROOT / "lib" / "data" / "local" / "learning_record_dao.dart"
QUIZ_PAGE_PATH = ROOT / "lib" / "presentation" / "pages" / "quiz" / "quiz_page.dart"
DOCUMENT_PAGE_PATH = (
    ROOT / "lib" / "presentation" / "pages" / "learning" / "document_page.dart"
)
VIDEO_PAGE_PATH = (
    ROOT / "lib" / "presentation" / "pages" / "learning" / "video_page.dart"
)
PROGRESS_PAGE_PATH = (
    ROOT / "lib" / "presentation" / "pages" / "learning" / "progress_page.dart"
)

WIDTH = 1920
HEIGHT = 1080

BG_TOP = (244, 248, 255)
BG_BOTTOM = (230, 238, 252)
WHITE = (255, 255, 255)
TEXT = (34, 44, 66)
TEXT_LIGHT = (92, 104, 128)
PRIMARY = (76, 110, 245)
PRIMARY_DARK = (48, 77, 201)
ACCENT = (18, 184, 134)
ORANGE = (255, 146, 43)
RED = (235, 87, 87)
PURPLE = (125, 92, 255)
BORDER = (210, 220, 240)
BOX_BG = (250, 252, 255)
TAG_BG = (233, 239, 255)
CODE_BG = (246, 248, 252)

TITLE_FONT_SIZE = 52
SUBTITLE_FONT_SIZE = 28
BODY_FONT_SIZE = 26
SMALL_FONT_SIZE = 20
CODE_FONT_SIZE = 21
SECTION_FONT_SIZE = 30


@dataclass
class TutorialSection:
    section_no: str
    title: str
    duration_hint: str
    scene: str
    narration: str
    points: List[str]
    image_mode: str = "arch"
    code_ref: Optional[Tuple[Path, str]] = None


def ensure_dirs() -> None:
    for path in [
        DIAGRAMS_DIR,
        VIDEO_DIR,
        GENERATED_DIR,
        SLIDES_DIR,
        AUDIO_DIR,
        CLIPS_DIR,
        VIDEO_OUTPUT_DIR,
    ]:
        path.mkdir(parents=True, exist_ok=True)


def load_font(size: int, bold: bool = False):
    candidates = []
    if os.name == "nt":
        if bold:
            candidates += [
                r"C:\Windows\Fonts\msyhbd.ttc",
                r"C:\Windows\Fonts\simhei.ttf",
                r"C:\Windows\Fonts\arialbd.ttf",
                r"C:\Windows\Fonts\calibrib.ttf",
            ]
        else:
            candidates += [
                r"C:\Windows\Fonts\msyh.ttc",
                r"C:\Windows\Fonts\simsun.ttc",
                r"C:\Windows\Fonts\arial.ttf",
                r"C:\Windows\Fonts\calibri.ttf",
            ]
    else:
        candidates += [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        ]

    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size=size)
        except Exception:
            continue
    return ImageFont.load_default()


FONT_TITLE = load_font(TITLE_FONT_SIZE, True)
FONT_SUBTITLE = load_font(SUBTITLE_FONT_SIZE, False)
FONT_BODY = load_font(BODY_FONT_SIZE, False)
FONT_BODY_BOLD = load_font(BODY_FONT_SIZE, True)
FONT_SMALL = load_font(SMALL_FONT_SIZE, False)
FONT_SMALL_BOLD = load_font(SMALL_FONT_SIZE, True)
FONT_CODE = load_font(CODE_FONT_SIZE, False)
FONT_SECTION = load_font(SECTION_FONT_SIZE, True)


def draw_vertical_gradient(
    img: Image.Image, top: Tuple[int, int, int], bottom: Tuple[int, int, int]
) -> None:
    draw = ImageDraw.Draw(img)
    for y in range(img.height):
        ratio = y / max(1, img.height - 1)
        color = tuple(int(top[i] * (1 - ratio) + bottom[i] * ratio) for i in range(3))
        draw.line([(0, y), (img.width, y)], fill=color)


def rounded_rect(draw: ImageDraw.ImageDraw, xy, fill, outline=None, radius=24, width=2):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def text_size(draw: ImageDraw.ImageDraw, text: str, font) -> Tuple[int, int]:
    bbox = draw.textbbox((0, 0), text, font=font)
    return bbox[2] - bbox[0], bbox[3] - bbox[1]


def center_text(
    draw: ImageDraw.ImageDraw, x: int, y: int, text: str, font, fill
) -> None:
    w, h = text_size(draw, text, font)
    draw.text((x - w / 2, y - h / 2), text, font=font, fill=fill)


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font, max_width: int) -> List[str]:
    lines: List[str] = []
    for paragraph in text.splitlines():
        p = paragraph.strip()
        if not p:
            lines.append("")
            continue
        current = ""
        for ch in p:
            candidate = current + ch
            w, _ = text_size(draw, candidate, font)
            if w <= max_width or not current:
                current = candidate
            else:
                lines.append(current)
                current = ch
        if current:
            lines.append(current)
    return lines


def draw_wrapped_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    box: Tuple[int, int, int, int],
    font,
    fill,
    line_spacing: int = 8,
) -> int:
    x1, y1, x2, y2 = box
    max_width = x2 - x1
    y = y1
    for line in wrap_text(draw, text, font, max_width):
        if y > y2:
            break
        draw.text((x1, y), line, font=font, fill=fill)
        _, h = text_size(draw, line or "中", font)
        y += h + line_spacing
    return y


def safe_name(text: str) -> str:
    return re.sub(r"[\\/:*?\"<>| ]+", "_", text)


def clean_text_for_tts(text: str) -> str:
    text = text.replace("SQLite", "S Q Lite")
    text = text.replace("Flutter", "Flutter")
    text = text.replace("GraphListPage", "Graph List Page")
    text = text.replace("GraphDetailPage", "Graph Detail Page")
    text = text.replace("HomePage", "Home Page")
    text = text.replace("InteractiveViewer", "Interactive Viewer")
    text = text.replace("GraphDao", "Graph Dao")
    text = text.replace("DatabaseHelper", "Database Helper")
    text = text.replace("learning_records", "learning records")
    text = text.replace("favorites", "favorites")
    return text


def write_puml_reference_files() -> None:
    arch = """@startuml
skinparam backgroundColor #F8FAFF
skinparam componentStyle rectangle
skinparam shadowing false
skinparam defaultFontName Microsoft YaHei
skinparam rectangle {
  BorderColor #4C6EF5
  RoundCorner 18
  BackgroundColor #FFFFFF
}
title 知识图谱核心功能架构图

actor 学生
actor 教师
actor 管理员

rectangle "Flutter UI" {
  rectangle "LoginPage\\n登录入口"
  rectangle "HomePage\\n导航中心"
  rectangle "GraphListPage\\n图谱列表"
  rectangle "GraphDetailPage\\n图谱详情与节点交互"
  rectangle "DocumentListPage\\n资料 / PPT"
  rectangle "VideoListPage\\n视频资源"
  rectangle "QuizPage\\n章节测验"
  rectangle "ProgressPage\\n学习进度"
  rectangle "FavoritesPage\\n收藏内容"
}

rectangle "Service / DAO" {
  rectangle "AuthService"
  rectangle "GraphDao"
  rectangle "QuizDao"
  rectangle "FavoriteDao"
  rectangle "LearningRecordDao"
  rectangle "UserDao"
  rectangle "DatabaseHelper"
}

database "SQLite / learning_data.db" as DB

学生 --> LoginPage
教师 --> LoginPage
管理员 --> LoginPage
LoginPage --> HomePage

HomePage --> GraphListPage
GraphListPage --> GraphDetailPage
GraphDetailPage --> FavoriteDao
GraphDetailPage --> LearningRecordDao
GraphListPage --> GraphDao
GraphDetailPage --> GraphDao
QuizPage --> QuizDao
ProgressPage --> QuizDao
ProgressPage --> LearningRecordDao
FavoritesPage --> FavoriteDao

GraphDao --> DatabaseHelper
QuizDao --> DatabaseHelper
FavoriteDao --> DatabaseHelper
LearningRecordDao --> DatabaseHelper
UserDao --> DatabaseHelper
DatabaseHelper --> DB
@enduml
"""
    flow = """@startuml
skinparam backgroundColor #F8FAFF
skinparam defaultFontName Microsoft YaHei
skinparam shadowing false
title 图谱功能演示流程图

start
:进入 HomePage;
:点击底部导航“图谱”;
:GraphListPage 加载图谱列表;
if (图谱数据存在?) then (是)
  :进入 GraphDetailPage;
  :读取 nodes 和 edges;
  :InteractiveViewer 缩放与拖拽浏览;
  :点击节点;
  :显示节点详情卡片;
  if (开始学习?) then (是)
    :写入 learning_records;
  else (收藏)
    :写入 favorites;
  endif
  :继续图谱浏览;
else (否)
  :显示暂无图谱数据;
endif
stop
@enduml
"""
    ARCH_PUML.write_text(arch, encoding="utf-8")
    FLOW_PUML.write_text(flow, encoding="utf-8")


def render_arch_png() -> None:
    img = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw_vertical_gradient(img, BG_TOP, BG_BOTTOM)
    draw = ImageDraw.Draw(img)

    draw.text((82, 54), "知识图谱核心功能架构图", font=FONT_TITLE, fill=TEXT)
    draw.text(
        (84, 126),
        "对照参考脚本的“系统架构图”风格，展示页面层、DAO层与数据层关系",
        font=FONT_SUBTITLE,
        fill=TEXT_LIGHT,
    )

    def actor(x: int, y: int, name: str):
        draw.ellipse((x, y, x + 70, y + 70), outline=PRIMARY_DARK, width=5, fill=WHITE)
        draw.line((x + 35, y + 70, x + 35, y + 145), fill=PRIMARY_DARK, width=5)
        draw.line((x - 5, y + 96, x + 75, y + 96), fill=PRIMARY_DARK, width=5)
        draw.line((x + 35, y + 145, x, y + 205), fill=PRIMARY_DARK, width=5)
        draw.line((x + 35, y + 145, x + 70, y + 205), fill=PRIMARY_DARK, width=5)
        center_text(draw, x + 35, y + 238, name, FONT_BODY, TEXT)

    actor(110, 300, "学生")
    actor(110, 480, "教师")
    actor(110, 660, "管理员")

    rounded_rect(
        draw, (280, 210, 1190, 860), fill=WHITE, outline=BORDER, radius=32, width=3
    )
    draw.text((315, 230), "Flutter 页面层", font=FONT_SECTION, fill=PRIMARY_DARK)

    ui_boxes = [
        ("LoginPage", "登录入口", (320, 300, 560, 390), PRIMARY),
        ("HomePage", "导航中心", (610, 300, 850, 390), PRIMARY),
        ("GraphListPage", "图谱列表", (900, 300, 1140, 390), ACCENT),
        ("GraphDetailPage", "图谱详情 / 节点交互", (320, 450, 610, 560), ACCENT),
        ("DocumentListPage", "资料 / PPT", (650, 450, 900, 560), ORANGE),
        ("VideoListPage", "视频资源", (940, 450, 1140, 560), RED),
        ("QuizPage", "章节测验", (320, 620, 560, 730), PURPLE),
        ("ProgressPage", "学习进度", (610, 620, 850, 730), ACCENT),
        ("FavoritesPage", "收藏内容", (900, 620, 1140, 730), ORANGE),
    ]
    for title, subtitle, rect, color in ui_boxes:
        rounded_rect(draw, rect, fill=BOX_BG, outline=color, radius=22, width=4)
        center_text(
            draw, (rect[0] + rect[2]) // 2, rect[1] + 26, title, FONT_SMALL_BOLD, color
        )
        draw_wrapped_text(
            draw,
            subtitle,
            (rect[0] + 18, rect[1] + 48, rect[2] - 18, rect[3] - 14),
            FONT_SMALL,
            TEXT,
        )

    rounded_rect(
        draw, (1230, 210, 1770, 760), fill=WHITE, outline=BORDER, radius=32, width=3
    )
    draw.text((1265, 230), "DAO / Service 层", font=FONT_SECTION, fill=PRIMARY_DARK)

    dao_boxes = [
        ("AuthService", (1260, 300, 1500, 376), PRIMARY),
        ("GraphDao", (1520, 300, 1720, 376), ACCENT),
        ("QuizDao", (1260, 420, 1450, 496), PURPLE),
        ("FavoriteDao", (1470, 420, 1720, 496), ORANGE),
        ("LearningRecordDao", (1260, 540, 1515, 616), RED),
        ("UserDao", (1535, 540, 1720, 616), PRIMARY),
        ("DatabaseHelper", (1345, 660, 1650, 736), PRIMARY_DARK),
    ]
    for label, rect, color in dao_boxes:
        rounded_rect(draw, rect, fill=BOX_BG, outline=color, radius=18, width=4)
        center_text(
            draw,
            (rect[0] + rect[2]) // 2,
            (rect[1] + rect[3]) // 2,
            label,
            FONT_SMALL_BOLD,
            color,
        )

    rounded_rect(
        draw,
        (1330, 830, 1660, 970),
        fill=WHITE,
        outline=PRIMARY_DARK,
        radius=30,
        width=4,
    )
    draw.ellipse((1405, 820, 1585, 868), fill=TAG_BG, outline=PRIMARY_DARK, width=4)
    draw.rectangle((1405, 844, 1585, 940), fill=TAG_BG, outline=PRIMARY_DARK, width=4)
    draw.ellipse((1405, 916, 1585, 964), fill=TAG_BG, outline=PRIMARY_DARK, width=4)
    center_text(
        draw, 1495, 898, "SQLite / learning_data.db", FONT_SMALL_BOLD, PRIMARY_DARK
    )

    def arrow(x1, y1, x2, y2, color=PRIMARY_DARK, width=4):
        draw.line((x1, y1, x2, y2), fill=color, width=width)
        angle = math.atan2(y2 - y1, x2 - x1)
        a = 14
        left = (
            x2 - a * math.cos(angle - math.pi / 6),
            y2 - a * math.sin(angle - math.pi / 6),
        )
        right = (
            x2 - a * math.cos(angle + math.pi / 6),
            y2 - a * math.sin(angle + math.pi / 6),
        )
        draw.polygon([(x2, y2), left, right], fill=color)

    arrow(225, 335, 320, 335)
    arrow(225, 515, 320, 345)
    arrow(225, 695, 320, 355)
    arrow(560, 335, 610, 335, PRIMARY)
    arrow(850, 345, 900, 345, ACCENT)
    arrow(455, 560, 1620, 338, ACCENT)
    arrow(455, 560, 1595, 456, ORANGE)
    arrow(455, 560, 1385, 578, RED)
    arrow(438, 674, 1355, 456, PURPLE)
    arrow(730, 730, 1385, 578, RED)
    arrow(1010, 676, 1595, 456, ORANGE)
    arrow(1380, 376, 1498, 660)
    arrow(1622, 376, 1498, 660)
    arrow(1355, 496, 1498, 660)
    arrow(1592, 496, 1498, 660)
    arrow(1390, 616, 1498, 660)
    arrow(1622, 616, 1498, 660)
    arrow(1498, 736, 1498, 830)

    draw.text(
        (82, 1006),
        "核心结论：图谱模块位于系统中心，向下连接本地数据库，向外联动收藏、学习记录、测验与进度。",
        font=FONT_SMALL,
        fill=TEXT_LIGHT,
    )
    img.save(ARCH_PNG)


def render_flow_png() -> None:
    img = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw_vertical_gradient(img, (248, 251, 255), (236, 244, 255))
    draw = ImageDraw.Draw(img)

    draw.text((82, 54), "图谱功能演示流程图", font=FONT_TITLE, fill=TEXT)
    draw.text(
        (84, 126),
        "对照参考脚本的“流程图”风格，展示从入口到学习与收藏联动的完整流程",
        font=FONT_SUBTITLE,
        fill=TEXT_LIGHT,
    )

    boxes = [
        ("1", "进入 HomePage\n首页导航", (130, 260, 400, 420), PRIMARY),
        ("2", "点击“图谱”\n进入 GraphListPage", (455, 260, 785, 420), ACCENT),
        ("3", "读取图谱列表\n本地数据库 graphs", (840, 260, 1170, 420), ORANGE),
        (
            "4",
            "进入 GraphDetailPage\n读取 nodes / edges",
            (1225, 260, 1595, 420),
            PURPLE,
        ),
        ("5", "缩放拖拽浏览\nInteractiveViewer", (260, 610, 620, 770), PRIMARY),
        ("6", "点击节点\n显示详情卡片", (720, 610, 1040, 770), RED),
        ("7", "开始学习或收藏\n沉淀学习行为", (1140, 610, 1585, 770), ACCENT),
    ]

    for tag, title, rect, color in boxes:
        rounded_rect(draw, rect, fill=WHITE, outline=color, radius=24, width=5)
        rounded_rect(
            draw,
            (rect[0] + 18, rect[1] + 18, rect[0] + 80, rect[1] + 80),
            fill=color,
            outline=color,
            radius=18,
            width=1,
        )
        center_text(draw, rect[0] + 49, rect[1] + 48, tag, FONT_SMALL_BOLD, WHITE)
        draw_wrapped_text(
            draw,
            title,
            (rect[0] + 105, rect[1] + 32, rect[2] - 20, rect[3] - 20),
            FONT_BODY,
            TEXT,
        )

    def arrow(x1, y1, x2, y2, color=PRIMARY_DARK, width=5):
        draw.line((x1, y1, x2, y2), fill=color, width=width)
        angle = math.atan2(y2 - y1, x2 - x1)
        a = 16
        left = (
            x2 - a * math.cos(angle - math.pi / 6),
            y2 - a * math.sin(angle - math.pi / 6),
        )
        right = (
            x2 - a * math.cos(angle + math.pi / 6),
            y2 - a * math.sin(angle + math.pi / 6),
        )
        draw.polygon([(x2, y2), left, right], fill=color)

    arrow(400, 340, 455, 340, PRIMARY)
    arrow(785, 340, 840, 340, ACCENT)
    arrow(1170, 340, 1225, 340, ORANGE)
    arrow(1410, 420, 1410, 505, PURPLE)
    arrow(1410, 505, 440, 610, PRIMARY_DARK)
    arrow(620, 690, 720, 690, PRIMARY)
    arrow(1040, 690, 1140, 690, RED)

    draw.text((130, 860), "演示结论：", font=FONT_BODY_BOLD, fill=TEXT)
    draw.text(
        (260, 860),
        "图谱不是静态展示页，而是课程知识的组织中心，也是学习记录、收藏、测验与进度的上游入口。",
        font=FONT_BODY,
        fill=TEXT_LIGHT,
    )
    draw.text((130, 920), "适合视频讲解重点：", font=FONT_BODY_BOLD, fill=TEXT)
    draw.text(
        (370, 920),
        "从“怎么进入图谱”讲到“节点交互后如何沉淀学习行为”。",
        font=FONT_BODY,
        fill=TEXT_LIGHT,
    )

    img.save(FLOW_PNG)


def render_data_model_png() -> None:
    img = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw_vertical_gradient(img, BG_TOP, BG_BOTTOM)
    draw = ImageDraw.Draw(img)

    draw.text((82, 54), "知识图谱核心数据模型图", font=FONT_TITLE, fill=TEXT)
    draw.text(
        (84, 126),
        "用数据库结构解释图谱、节点、边与学习行为如何互相关联",
        font=FONT_SUBTITLE,
        fill=TEXT_LIGHT,
    )

    def table(rect, title, fields, color):
        rounded_rect(draw, rect, fill=WHITE, outline=color, radius=24, width=4)
        rounded_rect(
            draw,
            (rect[0], rect[1], rect[2], rect[1] + 54),
            fill=color,
            outline=color,
            radius=22,
            width=1,
        )
        draw.text((rect[0] + 18, rect[1] + 13), title, font=FONT_SMALL_BOLD, fill=WHITE)
        y = rect[1] + 74
        for field in fields:
            draw.text((rect[0] + 18, y), field, font=FONT_SMALL, fill=TEXT)
            y += 34

    table(
        (130, 260, 470, 620),
        "graphs",
        ["id (PK)", "title", "graph_type", "layout"],
        PRIMARY,
    )
    table(
        (560, 220, 930, 720),
        "nodes",
        [
            "id (PK)",
            "graph_id (FK)",
            "title",
            "content",
            "node_type",
            "level",
            "x / y",
            "parent_id",
            "visible",
        ],
        ACCENT,
    )
    table(
        (1010, 260, 1370, 620),
        "edges",
        [
            "id (PK)",
            "graph_id (FK)",
            "source_id",
            "target_id",
            "edge_type",
            "label",
            "weight",
            "visible",
        ],
        ORANGE,
    )
    table(
        (1450, 220, 1800, 560),
        "favorites",
        ["id (PK)", "user_id", "node_id", "node_title", "favorite_time"],
        RED,
    )
    table(
        (1450, 620, 1800, 930),
        "learning_records",
        ["id (PK)", "user_id", "node_id", "node_title", "study_time", "completed_at"],
        PURPLE,
    )

    def arrow(x1, y1, x2, y2, color=PRIMARY_DARK, width=5):
        draw.line((x1, y1, x2, y2), fill=color, width=width)
        angle = math.atan2(y2 - y1, x2 - x1)
        a = 16
        left = (
            x2 - a * math.cos(angle - math.pi / 6),
            y2 - a * math.sin(angle - math.pi / 6),
        )
        right = (
            x2 - a * math.cos(angle + math.pi / 6),
            y2 - a * math.sin(angle + math.pi / 6),
        )
        draw.polygon([(x2, y2), left, right], fill=color)

    arrow(470, 420, 560, 420, PRIMARY)
    arrow(930, 420, 1010, 420, ACCENT)
    arrow(930, 490, 1450, 380, RED)
    arrow(930, 560, 1450, 760, PURPLE)

    draw.text(
        (86, 978),
        "说明：graphs 组织图谱；nodes 与 edges 组成图结构；favorites 与 learning_records 把图谱浏览转化为可追踪的学习行为。",
        font=FONT_SMALL,
        fill=TEXT_LIGHT,
    )
    img.save(DATA_PNG)


def choose_diagram(mode: str) -> Path:
    if mode == "arch":
        return ARCH_PNG
    if mode == "flow":
        return FLOW_PNG
    if mode == "data":
        return DATA_PNG
    return ARCH_PNG


def extract_block(path: Path, marker: str, max_lines: int = 32) -> str:
    if not path.exists():
        return f"// missing file: {path.name}"
    text = path.read_text(encoding="utf-8")
    lines = text.splitlines()
    start = 0
    for i, line in enumerate(lines):
        if marker in line:
            start = i
            break
    end = min(len(lines), start + max_lines)
    snippet = "\n".join(lines[start:end]).strip()
    return snippet or f"// marker not found: {marker}"


def draw_code_block(
    draw: ImageDraw.ImageDraw, code: str, box: Tuple[int, int, int, int]
) -> None:
    rounded_rect(draw, box, fill=CODE_BG, outline=(224, 230, 242), radius=18, width=2)
    x1, y1, x2, y2 = box
    y = y1 + 18
    for line in code.splitlines():
        if y > y2 - 28:
            break
        draw.text((x1 + 18, y), line.rstrip(), font=FONT_CODE, fill=TEXT)
        _, h = text_size(draw, line or "中", FONT_CODE)
        y += h + 6


def paste_image(
    base: Image.Image, source: Path, box: Tuple[int, int, int, int]
) -> None:
    if not source.exists():
        return
    img = Image.open(source).convert("RGB")
    target_w = box[2] - box[0]
    target_h = box[3] - box[1]
    ratio = min(target_w / img.width, target_h / img.height)
    new_size = (max(1, int(img.width * ratio)), max(1, int(img.height * ratio)))
    resized = img.resize(new_size)
    x = box[0] + (target_w - new_size[0]) // 2
    y = box[1] + (target_h - new_size[1]) // 2
    base.paste(resized, (x, y))


def build_sections() -> List[TutorialSection]:
    return [
        TutorialSection(
            section_no="第一部分",
            title="项目概述",
            duration_hint="约2分钟",
            scene="显示项目名称、课程定位与核心图谱架构图",
            narration=(
                "大家好，欢迎来到移动应用开发知识图谱学习系统教程。"
                "这次视频不再讲登录注册，而是聚焦系统真正的核心能力，也就是知识图谱。"
                "这个项目基于 Flutter 构建，数据落在本地 SQLite 数据库中。"
                "图谱模块负责把课程章节、知识点、能力要求和学习行为组织成结构化关系，"
                "并进一步联动资料、视频、测验和进度分析。"
            ),
            points=[
                "项目定位：移动应用开发知识图谱学习系统",
                "实现方式：Flutter 页面层 + DAO 数据访问层 + SQLite 本地数据库",
                "核心价值：把课程知识结构化，把学习行为可视化",
            ],
            image_mode="arch",
        ),
        TutorialSection(
            section_no="第一部分",
            title="功能需求分析",
            duration_hint="约2分钟",
            scene="显示知识图谱需求点与系统闭环关系",
            narration=(
                "从业务角度看，知识图谱模块至少要满足四个需求。"
                "第一，能够展示图谱总览，让用户快速进入不同主题图谱。"
                "第二，能够展示节点与边，形成完整知识结构。"
                "第三，点击节点后要进入学习动作，比如开始学习或加入收藏。"
                "第四，这些行为还要反馈到学习进度、测验和复习流程中，构成学习闭环。"
            ),
            points=[
                "图谱总览：快速进入课程图谱、技术栈图谱、学习图谱等主题",
                "结构浏览：基于节点与边展示知识关系",
                "行为沉淀：学习记录与收藏必须可追踪",
                "闭环联动：进度、资料、视频、测验围绕图谱展开",
            ],
            image_mode="flow",
        ),
        TutorialSection(
            section_no="第二部分",
            title="整体架构设计",
            duration_hint="约3分钟",
            scene="展示页面层、DAO层和数据库层结构图",
            narration=(
                "接下来我们看整体架构。"
                "在页面层，HomePage 负责统一导航，GraphListPage 负责图谱列表，"
                "GraphDetailPage 负责图谱详情和节点交互。"
                "在数据访问层，GraphDao、FavoriteDao、LearningRecordDao 等类"
                "分别封装图谱读取、收藏写入和学习记录写入逻辑。"
                "最底层由 DatabaseHelper 统一管理本地数据库连接和建表逻辑。"
            ),
            points=[
                "页面层：首页导航、图谱列表、图谱详情、进度、资料、视频、测验",
                "DAO层：GraphDao、FavoriteDao、LearningRecordDao、QuizDao",
                "数据层：DatabaseHelper + learning_data.db",
            ],
            image_mode="arch",
            code_ref=(HOME_PATH, "class HomePage"),
        ),
        TutorialSection(
            section_no="第二部分",
            title="数据库设计",
            duration_hint="约3分钟",
            scene="显示数据模型图与关键表关系",
            narration=(
                "知识图谱能否真正工作，关键在于底层数据结构。"
                "graphs 表负责图谱总体信息，nodes 表记录具体知识节点，"
                "edges 表负责节点之间的连接关系。"
                "除此之外，favorites 表和 learning_records 表"
                "把图谱浏览转化为后续可分析的学习行为。"
                "这也是知识图谱系统和普通静态目录页最大的区别。"
            ),
            points=[
                "graphs：图谱元数据",
                "nodes：知识点、章节、文件等节点",
                "edges：节点之间的包含或关联关系",
                "favorites / learning_records：记录用户行为",
            ],
            image_mode="data",
            code_ref=(DATABASE_HELPER_PATH, "CREATE TABLE IF NOT EXISTS graphs"),
        ),
        TutorialSection(
            section_no="第三部分",
            title="图谱列表页实现",
            duration_hint="约2分钟",
            scene="展示 GraphListPage 与 GraphDao 的协作逻辑",
            narration=(
                "图谱功能的入口是 GraphListPage。"
                "这个页面启动后，会调用 GraphDao 读取全部图谱，"
                "并以列表卡片的形式展示图谱标题和类型。"
                "用户点击任意图谱后，就会跳转到 GraphDetailPage。"
                "所以这里承担的是图谱发现与入口分发的职责。"
            ),
            points=[
                "页面初始化后调用 _loadGraphs",
                "GraphDao 负责从 graphs 表查询图谱记录",
                "列表卡片点击后跳转到 GraphDetailPage",
            ],
            image_mode="arch",
            code_ref=(GRAPH_LIST_PATH, "Future<void> _loadGraphs()"),
        ),
        TutorialSection(
            section_no="第三部分",
            title="图谱详情页实现",
            duration_hint="约3分钟",
            scene="显示 GraphDetailPage 的节点加载与详情卡片结构",
            narration=(
                "进入图谱详情页之后，系统会继续读取当前图谱对应的 nodes 和 edges。"
                "如果节点没有预设坐标，页面还会自动计算树状布局。"
                "用户可以通过 InteractiveViewer 进行缩放、拖拽和浏览，"
                "点击某个节点之后，下方展示节点详情卡片，"
                "包括标题、类型、内容，以及开始学习和收藏两个操作入口。"
            ),
            points=[
                "加载 nodes / edges 并构建图谱视图",
                "支持自动布局、缩放和拖拽",
                "点击节点后展示详情卡片与操作按钮",
            ],
            image_mode="flow",
            code_ref=(GRAPH_DETAIL_PATH, "Future<void> _loadGraphData()"),
        ),
        TutorialSection(
            section_no="第三部分",
            title="节点交互与学习行为",
            duration_hint="约2分钟",
            scene="显示节点详情中的操作按钮与学习行为流向",
            narration=(
                "知识图谱真正的价值，不在于把节点画出来，"
                "而在于用户点击节点后可以产生后续行为。"
                "在节点详情里，开始学习按钮对应学习记录的沉淀，"
                "收藏按钮对应重要知识点的归档。"
                "这两个动作使得图谱成为整个学习系统的上游入口。"
            ),
            points=[
                "开始学习：用于记录学习节点与完成时间",
                "收藏：用于构建个人重点知识点清单",
                "图谱从展示层升级为行为入口层",
            ],
            image_mode="flow",
            code_ref=(GRAPH_DETAIL_PATH, "Widget _buildNodeDetail()"),
        ),
        TutorialSection(
            section_no="第四部分",
            title="DAO 与数据访问实现",
            duration_hint="约2分钟",
            scene="展示 GraphDao、FavoriteDao、LearningRecordDao 的职责拆分",
            narration=(
                "从代码组织上看，项目没有把所有逻辑堆到页面里，"
                "而是通过 DAO 进行职责拆分。"
                "GraphDao 只负责图谱读取，FavoriteDao 只负责收藏读写，"
                "LearningRecordDao 只负责学习记录统计与保存。"
                "这种方式让页面层关注交互，数据层关注存取，结构更清晰。"
            ),
            points=[
                "GraphDao：图谱、节点、边查询",
                "FavoriteDao：收藏新增、删除、统计",
                "LearningRecordDao：学习记录新增、统计、完成情况",
            ],
            image_mode="arch",
            code_ref=(GRAPH_DAO_PATH, "Future<List<GraphModel>> getAllGraphs()"),
        ),
        TutorialSection(
            section_no="第四部分",
            title="与资料、视频、测验的联动",
            duration_hint="约3分钟",
            scene="展示图谱作为上游入口如何联动其它模块",
            narration=(
                "知识图谱并不是孤立存在。"
                "在这个项目里，资料页面管理 PDF 和 PPT，视频页面管理课程视频，"
                "测验页面提供章节练习，进度页面汇总学习记录和测验结果。"
                "图谱模块可以看作课程知识的地图，而资料、视频和测验"
                "则是围绕这张地图展开的学习资源和验证工具。"
            ),
            points=[
                "资料模块：提供 PDF 与 PPT 课件访问",
                "视频模块：提供课程视频资源入口",
                "测验模块：按章节组织练习题并沉淀成绩",
                "进度模块：对学习记录和测验结果做汇总",
            ],
            image_mode="arch",
            code_ref=(PROGRESS_PAGE_PATH, "class ProgressPage"),
        ),
        TutorialSection(
            section_no="第五部分",
            title="演示流程总结",
            duration_hint="约2分钟",
            scene="完整回顾从入口到数据沉淀的流程图",
            narration=(
                "现在我们把整个演示流程再串起来。"
                "用户从首页进入图谱，列表页读取图谱，详情页读取节点和边，"
                "然后通过缩放与拖拽浏览知识结构。"
                "点击节点后，系统展示详情，并允许开始学习或加入收藏。"
                "最终这些行为被写入数据库，为后续的进度分析和学习建议提供依据。"
            ),
            points=[
                "入口：HomePage 到 GraphListPage",
                "浏览：GraphDetailPage + InteractiveViewer",
                "沉淀：favorites 与 learning_records",
                "结果：形成可追踪、可分析的学习闭环",
            ],
            image_mode="flow",
        ),
        TutorialSection(
            section_no="第六部分",
            title="测试与运行结果",
            duration_hint="约2分钟",
            scene="展示测试通过、APK 构建成功与视频生成产物",
            narration=(
                "最后总结一下当前工程状态。"
                "项目的测试已经能够通过，Android 调试包和发布包也都能成功构建。"
                "在此基础上，我们进一步生成了图谱功能教程脚本、"
                "结构模型图、PPT 和最终视频。"
                "这说明项目不仅具备功能实现，也具备了教学展示与成果汇报的条件。"
            ),
            points=[
                "测试：Widget 测试已通过",
                "构建：Debug APK 与 Release APK 均可输出",
                "成果：脚本、PPT、模型图、视频已成套生成",
            ],
            image_mode="data",
        ),
        TutorialSection(
            section_no="第六部分",
            title="结束语与后续扩展",
            duration_hint="约1分钟",
            scene="收束画面，说明后续还可扩展路径、资料、视频和测验四支视频",
            narration=(
                "以上就是知识图谱核心功能的完整讲解。"
                "这一版重点解释了为什么图谱是项目核心，"
                "以及它如何连接学习记录、收藏、资料、视频和测验。"
                "在此基础上，还可以继续扩展生成学习路径视频、"
                "资料播放视频、视频资源播放视频，以及章节测验视频。"
                "这样就能形成完整的五支系列教学视频。"
            ),
            points=[
                "本视频定位：图谱功能教程",
                "后续可扩展：路径、资料、视频、测验四支配套视频",
                "适合场景：课程汇报、项目答辩、教学展示",
            ],
            image_mode="arch",
        ),
    ]


def build_markdown_script(sections: List[TutorialSection]) -> None:
    lines: List[str] = []
    lines += [
        "# 知识图谱核心功能教程脚本 v2",
        "",
        "## 教程信息",
        "",
        "- **教程标题**：移动应用开发知识图谱核心功能教程 - 图谱模块完整讲解",
        "- **时长**：约18-22分钟（可根据语速与字幕停留时间调整）",
        "- **目标受众**：课程学习者、项目答辩评审、教学演示观众",
        "- **学习目标**：理解知识图谱在项目中的定位、结构、数据模型和交互流程",
        "",
        "***",
        "",
    ]
    current = None
    for idx, sec in enumerate(sections, start=1):
        if current != sec.section_no:
            current = sec.section_no
            lines += [
                f"## {sec.section_no}：{sec.title if '项目概述' in sec.title else '教程内容'}",
                "",
            ]
        lines += [
            f"### {idx}. {sec.title}（{sec.duration_hint}）",
            "",
            f"**画面**：{sec.scene}",
            "",
            "**旁白**：",
            f"“{sec.narration}”",
            "",
            "**要点**：",
            "",
        ]
        lines += [f"- {p}" for p in sec.points]
        if sec.code_ref:
            lines += [
                "",
                f"**代码来源**：`{sec.code_ref[0].relative_to(ROOT)}`",
            ]
        lines += ["", ""]
    lines += [
        "***",
        "",
        "## 生成产物",
        "",
        "- `docs/diagrams/knowledge_graph_feature_architecture.puml`",
        "- `docs/diagrams/knowledge_graph_feature_flow.puml`",
        "- `docs/diagrams/knowledge_graph_feature_architecture.png`",
        "- `docs/diagrams/knowledge_graph_feature_flow.png`",
        "- `docs/diagrams/knowledge_graph_feature_data_model.png`",
        "- `docs/video/video_script_v2.md`",
        "- `video_output/知识图谱核心功能_图谱功能教程_v2.pptx`",
        "- `video_output/知识图谱核心功能_图谱功能教程_v2.mp4`",
    ]
    SCRIPT_MD_PATH.write_text("\n".join(lines), encoding="utf-8")


def render_slide(section: TutorialSection, index: int, total: int) -> Path:
    img = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw_vertical_gradient(img, BG_TOP, BG_BOTTOM)
    draw = ImageDraw.Draw(img)

    rounded_rect(
        draw, (50, 36, 1870, 160), fill=WHITE, outline=BORDER, radius=28, width=2
    )
    draw.text((85, 58), "移动应用开发知识图谱核心功能教程", font=FONT_TITLE, fill=TEXT)
    draw.text(
        (88, 118),
        f"{section.section_no} · {section.title}",
        font=FONT_SUBTITLE,
        fill=TEXT_LIGHT,
    )
    rounded_rect(
        draw, (1660, 62, 1820, 118), fill=PRIMARY, outline=PRIMARY, radius=20, width=1
    )
    center_text(draw, 1740, 90, f"{index}/{total}", FONT_SMALL_BOLD, WHITE)

    left_box = (55, 190, 910, 1005)
    right_top = (940, 190, 1865, 620)
    right_bottom = (940, 650, 1865, 1005)

    rounded_rect(draw, left_box, fill=WHITE, outline=BORDER, radius=30, width=3)
    rounded_rect(draw, right_top, fill=WHITE, outline=BORDER, radius=30, width=3)
    rounded_rect(draw, right_bottom, fill=WHITE, outline=BORDER, radius=30, width=3)

    draw.text((90, 220), "教程讲解", font=FONT_SECTION, fill=PRIMARY_DARK)
    rounded_rect(
        draw, (90, 272, 250, 320), fill=TAG_BG, outline=TAG_BG, radius=16, width=1
    )
    draw.text(
        (112, 286),
        f"建议时长：{section.duration_hint}",
        font=FONT_SMALL_BOLD,
        fill=PRIMARY_DARK,
    )

    y = 348
    draw.text((90, y), "画面：", font=FONT_BODY_BOLD, fill=TEXT)
    y = draw_wrapped_text(draw, section.scene, (170, y, 850, y + 120), FONT_BODY, TEXT)
    y += 14

    draw.text((90, y), "旁白：", font=FONT_BODY_BOLD, fill=TEXT)
    y += 42
    y = draw_wrapped_text(
        draw, section.narration, (90, y, 850, 730), FONT_BODY, TEXT, line_spacing=10
    )
    y += 20

    draw.text((90, y), "要点：", font=FONT_BODY_BOLD, fill=TEXT)
    y += 44
    for point in section.points:
        rounded_rect(
            draw,
            (98, y + 6, 118, y + 26),
            fill=PRIMARY,
            outline=PRIMARY,
            radius=7,
            width=1,
        )
        y = (
            draw_wrapped_text(
                draw, point, (132, y, 845, y + 80), FONT_BODY, TEXT, line_spacing=8
            )
            + 12
        )

    draw.text((975, 220), "模型图 / 结构图", font=FONT_SECTION, fill=PRIMARY_DARK)
    paste_image(img, choose_diagram(section.image_mode), (975, 270, 1830, 590))

    draw.text((975, 680), "代码要点", font=FONT_SECTION, fill=PRIMARY_DARK)
    if section.code_ref:
        code = extract_block(section.code_ref[0], section.code_ref[1], 16)
    else:
        code = "// 本页以结构和流程讲解为主，无单独代码片段"
    draw_code_block(draw, code, (975, 730, 1830, 965))

    out = SLIDES_DIR / f"{index:02d}_{safe_name(section.title)}.png"
    img.save(out)
    return out


def generate_slides(sections: List[TutorialSection]) -> List[Path]:
    return [
        render_slide(section, i, len(sections))
        for i, section in enumerate(sections, start=1)
    ]


def choose_voice(engine, preferred_keywords: Iterable[str]) -> Optional[str]:
    voices = engine.getProperty("voices") or []
    preferred_lower = [k.lower() for k in preferred_keywords]
    for voice in voices:
        content = f"{getattr(voice, 'id', '')} {getattr(voice, 'name', '')}".lower()
        if any(k in content for k in preferred_lower):
            return voice.id
    return voices[0].id if voices else None


def save_tts_audio(text: str, out_path: Path, rate: int = 165) -> bool:
    if pyttsx3 is None:
        return False
    try:
        engine = pyttsx3.init()
        voice_id = choose_voice(
            engine, ["zh", "chinese", "huihui", "xiaoxiao", "yaoyao", "hanhan"]
        )
        if voice_id:
            engine.setProperty("voice", voice_id)
        engine.setProperty("rate", rate)
        engine.save_to_file(clean_text_for_tts(text), str(out_path))
        engine.runAndWait()
        engine.stop()
        return out_path.exists() and out_path.stat().st_size > 0
    except Exception:
        return False


def create_silent_wav(duration_seconds: float, out_path: Path) -> None:
    sample_rate = 22050
    channels = 1
    sampwidth = 2
    total_frames = max(1, int(duration_seconds * sample_rate))
    silence = struct.pack("<h", 0)
    with wave.open(str(out_path), "wb") as wav_file:
        wav_file.setnchannels(channels)
        wav_file.setsampwidth(sampwidth)
        wav_file.setframerate(sample_rate)
        for _ in range(total_frames):
            wav_file.writeframesraw(silence)


def wav_duration(path: Path) -> float:
    with wave.open(str(path), "rb") as wav_file:
        frames = wav_file.getnframes()
        rate = wav_file.getframerate()
        return frames / float(rate)


def which(programs: List[str]) -> Optional[str]:
    for program in programs:
        found = shutil.which(program)
        if found:
            return found
    return None


def ffmpeg_path() -> Optional[str]:
    return which(
        ["ffmpeg", r"C:\Users\ldl\AppData\Local\Microsoft\WinGet\Links\ffmpeg.exe"]
    )


def build_clip(image_path: Path, audio_path: Path, output_path: Path) -> bool:
    ffmpeg = ffmpeg_path()
    if not ffmpeg:
        return False
    cmd = [
        ffmpeg,
        "-y",
        "-loop",
        "1",
        "-i",
        str(image_path),
        "-i",
        str(audio_path),
        "-c:v",
        "libx264",
        "-tune",
        "stillimage",
        "-c:a",
        "aac",
        "-b:a",
        "192k",
        "-pix_fmt",
        "yuv420p",
        "-shortest",
        "-vf",
        "scale=1920:1080,fps=30",
        str(output_path),
    ]
    result = subprocess.run(
        cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, check=False
    )
    return result.returncode == 0 and output_path.exists()


def concat_clips(clips: List[Path], output_path: Path) -> bool:
    ffmpeg = ffmpeg_path()
    if not ffmpeg:
        return False
    manifest = GENERATED_DIR / "clips_manifest_v2.txt"
    lines = [f"file '{clip.as_posix()}'" for clip in clips]
    manifest.write_text("\n".join(lines), encoding="utf-8")
    cmd = [
        ffmpeg,
        "-y",
        "-f",
        "concat",
        "-safe",
        "0",
        "-i",
        str(manifest),
        "-c",
        "copy",
        str(output_path),
    ]
    result = subprocess.run(
        cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, check=False
    )
    return result.returncode == 0 and output_path.exists()


def generate_audio_and_video(
    slides: List[Path], sections: List[TutorialSection]
) -> bool:
    clips: List[Path] = []
    for idx, (slide, section) in enumerate(zip(slides, sections), start=1):
        audio_path = AUDIO_DIR / f"{idx:02d}_{safe_name(section.title)}.wav"
        clip_path = CLIPS_DIR / f"{idx:02d}_{safe_name(section.title)}.mp4"

        text = f"{section.title}。{section.narration}。要点包括：{'；'.join(section.points)}。"
        ok = save_tts_audio(text, audio_path)
        if not ok:
            create_silent_wav(8.0, audio_path)

        duration = max(6.0, wav_duration(audio_path))
        if duration < 6.5:
            create_silent_wav(7.0, audio_path)
        if not build_clip(slide, audio_path, clip_path):
            return False
        clips.append(clip_path)

    return concat_clips(clips, VIDEO_PATH)


def generate_pptx(slides: List[Path], sections: List[TutorialSection]) -> None:
    if Presentation is None:
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_image, section in zip(slides, sections):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title = slide.shapes.add_textbox(
            Inches(0.45), Inches(0.2), Inches(8.3), Inches(0.5)
        )
        p = title.text_frame.paragraphs[0]
        p.text = f"{section.section_no} · {section.title}"
        p.font.size = Pt(28)
        p.font.bold = True

        subtitle = slide.shapes.add_textbox(
            Inches(0.47), Inches(0.72), Inches(8.8), Inches(0.35)
        )
        p2 = subtitle.text_frame.paragraphs[0]
        p2.text = f"画面：{section.scene}"
        p2.font.size = Pt(14)

        slide.shapes.add_picture(
            str(slide_image),
            Inches(0.35),
            Inches(1.15),
            width=Inches(12.55),
            height=Inches(6.0),
        )

    prs.save(PPTX_PATH)


def print_summary(slides: List[Path], video_ok: bool) -> None:
    print("已生成教程式图谱视频 v2：")
    print(f"- {ARCH_PUML}")
    print(f"- {FLOW_PUML}")
    print(f"- {ARCH_PNG}")
    print(f"- {FLOW_PNG}")
    print(f"- {DATA_PNG}")
    print(f"- {SCRIPT_MD_PATH}")
    print(f"- {PPTX_PATH}")
    for slide in slides:
        print(f"- {slide}")
    if video_ok:
        print(f"- {VIDEO_PATH}")
    else:
        print("- 视频未成功生成，请检查 ffmpeg 或音频环境。")


def main() -> None:
    ensure_dirs()
    write_puml_reference_files()
    render_arch_png()
    render_flow_png()
    render_data_model_png()

    sections = build_sections()
    build_markdown_script(sections)
    slides = generate_slides(sections)
    generate_pptx(slides, sections)
    video_ok = generate_audio_and_video(slides, sections)
    print_summary(slides, video_ok)


if __name__ == "__main__":
    main()
