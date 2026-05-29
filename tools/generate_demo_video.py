#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
MAD-KGDT Demo Video Generator (6 segments, moviepy + edge_tts)

Generates the 6 demo segments from demo_script.md as slide-based videos
with TTS narration, subtitle burn-in, and SRT output.

Segments:
  #0 Opening         15s — PPT Slide 1-3
  #1 Knowledge Graph 30s — PPT Slide 6
  #2 AI 24 Agents    30s — PPT Slide 7
  #3 Digital Twin    45s — PPT Slide 8-10
  #4 4-Platform      30s — PPT Slide 11-13
  #5 Closing         15s — PPT Slide 14-16

Usage:
  python tools/generate_demo_video.py

Requirements:
  pip install moviepy==1.0.3 edge-tts Pillow
"""

from __future__ import annotations

import asyncio
import os
import sys
import time
import wave
from dataclasses import dataclass
from pathlib import Path

# Fix Windows GBK console encoding
if sys.stdout.encoding in ('gbk', 'cp936'):
    sys.stdout.reconfigure(encoding='utf-8')  # type: ignore[attr-defined]

from PIL import Image, ImageDraw, ImageFont

# ── edge_tts ──────────────────────────────────────────────────────────────
try:
    import edge_tts

    _EDGE_OK = True
except ImportError:
    _EDGE_OK = False
    print("Warning: edge-tts not found.  pip install edge-tts")

# ── moviepy ───────────────────────────────────────────────────────────────
try:
    from moviepy.editor import AudioFileClip, ImageClip, concatenate_videoclips

    _MOVIEPY_OK = True
except ImportError:
    _MOVIEPY_OK = False
    print("Error: moviepy not found.  pip install moviepy==1.0.3")

# ── pyttsx3 fallback ──────────────────────────────────────────────────────
try:
    import pyttsx3 as _pyttsx3_mod  # type: ignore[import]

    _PYTTSX3_OK = True
except Exception:
    _pyttsx3_mod = None
    _PYTTSX3_OK = False

# ═══════════════════════════════════════════════════════════════════════════
# PATHS
# ═══════════════════════════════════════════════════════════════════════════
ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "video_output"
DEMO_DIR = ROOT / "docs" / "video" / "demo"
SLIDES_DIR = DEMO_DIR / "slides"
SENT_DIR = DEMO_DIR / "sent"
AUDIO_DIR = DEMO_DIR / "audio"
TEMP_DIR = DEMO_DIR / "temp"
SCREENSHOTS_DIR = DEMO_DIR / "screenshots"

def _ss(name: str) -> Path | None:
    p = SCREENSHOTS_DIR / name
    return p if p.exists() else None

# ═══════════════════════════════════════════════════════════════════════════
# VIDEO PARAMETERS
# ═══════════════════════════════════════════════════════════════════════════
W, H = 1920, 1080
FPS = 30
TTS_VOICE = "zh-CN-XiaoxiaoNeural"
TTS_RATE = "+15%"
CLIP_TAIL = 0.15
SLIDE_TAIL = 0.40
SUB_H = 90
CONTENT_H = H - SUB_H
TITLE_H = 84

# ── colours ───────────────────────────────────────────────────────────────
BG_TOP = (20, 36, 68)
BG_BOT = (32, 54, 96)
TITLE_BG = (13, 24, 50)
ACCENT = (255, 198, 55)
TEXT = (236, 241, 252)
TEXT_MUTE = (165, 185, 220)
PRIMARY = (92, 172, 255)
GREEN = (72, 208, 126)
ORANGE = (255, 152, 55)
PURPLE = (168, 110, 255)
RED_C = (235, 85, 85)
BORDER = (52, 84, 140)
CODE_BG = (24, 30, 44)
PANEL = (26, 44, 76)
SUB_BG = (6, 6, 12)
WHITE = (255, 255, 255)

# ── fon ts ────────────────────────────────────────────────────────────────
def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    paths: list[str] = (
        [r"C:\Windows\Fonts\msyhbd.ttc", r"C:\Windows\Fonts\simhei.ttf"]
        if bold
        else [r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\simsun.ttc"]
    ) + [r"C:\Windows\Fonts\arial.ttf"]
    for p in paths:
        try:
            return ImageFont.truetype(p, size)
        except Exception:
            continue
    return ImageFont.load_default()


FT = _font(48, True)
FS = _font(26, False)
FB = _font(24, True)
FN = _font(22, False)
FK = _font(19, False)
FKB = _font(19, True)
FSU = _font(27, False)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE SPEC
# ═══════════════════════════════════════════════════════════════════════════
@dataclass
class SlideSpec:
    seg_id: str
    title: str
    subtitle: str
    bullets: list[str]
    narration: str
    voice_segments: list[str] | None = None
    image_path: Path | None = None
    image_caption: str | None = None


# ═══════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════
def _tsz(draw: ImageDraw.ImageDraw, text: str, font) -> tuple[int, int]:
    bb = draw.textbbox((0, 0), text, font=font)
    return int(bb[2] - bb[0]), int(bb[3] - bb[1])


def _wrap(draw: ImageDraw.ImageDraw, text: str, font, max_w: int) -> list[str]:
    lines: list[str] = []
    for para in text.splitlines():
        cur = ""
        for ch in para.strip():
            cand = cur + ch
            if _tsz(draw, cand, font)[0] > max_w and cur:
                lines.append(cur)
                cur = ch
            else:
                cur = cand
        if cur:
            lines.append(cur)
    return lines


def _draw_wrapped(
    draw: ImageDraw.ImageDraw,
    text: str,
    x: int,
    y: int,
    font,
    fill,
    max_w: int,
    gap: int = 5,
) -> int:
    for line in _wrap(draw, text, font, max_w):
        draw.text((x, y), line, font=font, fill=fill)
        y += _tsz(draw, line or "M", font)[1] + gap
    return y


def _gradient(img: Image.Image) -> None:
    draw = ImageDraw.Draw(img)
    for y in range(img.height):
        r = BG_TOP[0] + int((BG_BOT[0] - BG_TOP[0]) * y / img.height)
        g = BG_TOP[1] + int((BG_BOT[1] - BG_TOP[1]) * y / img.height)
        b = BG_TOP[2] + int((BG_BOT[2] - BG_TOP[2]) * y / img.height)
        draw.line([(0, y), (img.width, y)], fill=(r, g, b))


def _paste_fit(
    base: Image.Image, src: Path | None, box: tuple[int, int, int, int]
) -> None:
    if not src or not src.exists():
        return
    try:
        im = Image.open(src).convert("RGB")
    except Exception:
        return
    bw, bh = box[2] - box[0], box[3] - box[1]
    ratio = min(bw / im.width, bh / im.height)
    nw, nh = max(1, int(im.width * ratio)), max(1, int(im.height * ratio))
    im = im.resize((nw, nh), Image.LANCZOS)
    base.paste(im, (box[0] + (bw - nw) // 2, box[1] + (bh - nh) // 2))


def _safe(text: str) -> str:
    keep = []
    for ch in text:
        if ch.isalnum() or ch in ("-", "_") or "\u4e00" <= ch <= "\u9fff":
            keep.append(ch)
        else:
            keep.append("_")
    return "".join(keep).strip("_") or "slide"


# ═══════════════════════════════════════════════════════════════════════════
# DRAWING: Mock Page, Data Card, Badge, Feature Grid
# ═══════════════════════════════════════════════════════════════════════════
def draw_feature_card(
    draw: ImageDraw.ImageDraw,
    x: int,
    y: int,
    w: int,
    h: int,
    icon: str,
    title: str,
    desc: str,
    color: tuple[int, int, int],
) -> None:
    """Draw a single feature card with icon, title, description."""
    draw.rounded_rectangle((x, y, x + w, y + h), radius=12, fill=PANEL, outline=BORDER, width=2)
    # icon circle
    cx, cy = x + 30, y + 28
    draw.ellipse((cx - 14, cy - 14, cx + 14, cy + 14), fill=color)
    draw.text((cx - 6, cy - 10), icon, font=_font(20, True), fill=(0, 0, 0))
    # title
    draw.text((x + 56, y + 14), title, font=_font(22, True), fill=TEXT)
    # description
    _draw_wrapped(draw, desc, x + 56, y + 46, _font(17), TEXT_MUTE, w - 70, 4)


def draw_badge(
    draw: ImageDraw.ImageDraw, label: str, x: int, y: int, color: tuple[int, int, int]
) -> tuple[int, int]:
    """Draw a small badge/label, return (x2, y2)."""
    tw, th = _tsz(draw, label, FK)
    bw, bh = tw + 24, th + 12
    bx, by = x, y
    draw.rounded_rectangle((bx, by, bx + bw, by + bh), radius=8, fill=color, outline=color, width=1)
    draw.text((bx + 12, by + 6), label, font=FK, fill=(8, 16, 36))
    return bx + bw, by + bh


def draw_data_card(
    draw: ImageDraw.ImageDraw,
    x: int,
    y: int,
    w: int,
    h: int,
    value: str,
    label: str,
    color: tuple[int, int, int],
) -> None:
    """Draw a KPI data card (big number + label)."""
    draw.rounded_rectangle((x, y, x + w, y + h), radius=14, fill=PANEL, outline=BORDER, width=2)
    # big number
    draw.text((x + 20, y + 16), value, font=_font(40, True), fill=color)
    # label
    draw.text((x + 20, y + 62), label, font=FK, fill=TEXT_MUTE)


# ═══════════════════════════════════════════════════════════════════════════
# TTS + DURATION
# ═══════════════════════════════════════════════════════════════════════════
async def _edge_async(text: str, path: Path) -> None:
    comm = edge_tts.Communicate(text, TTS_VOICE, rate=TTS_RATE)
    await comm.save(str(path))


_pyttsx3_engine = None


def _init_pyttsx3() -> None:
    global _pyttsx3_engine
    if _pyttsx3_engine is not None or not _PYTTSX3_OK:
        return
    try:
        eng = _pyttsx3_mod.init()
        eng.setProperty("rate", 155)
        voices = eng.getProperty("voices") or []
        for v in voices:
            combined = f"{getattr(v, 'name', '')} {getattr(v, 'id', '')}".lower()
            if any(k in combined for k in ["zh", "chinese", "huihui", "xiaoxiao"]):
                eng.setProperty("voice", v.id)
                break
        _pyttsx3_engine = eng
    except Exception as e:
        print(f"  pyttsx3 init error: {e}")


def _pyttsx3_gen(text: str, wav_path: Path) -> bool:
    _init_pyttsx3()
    if _pyttsx3_engine is None:
        return False
    try:
        _pyttsx3_engine.save_to_file(text, str(wav_path))
        _pyttsx3_engine.runAndWait()
        time.sleep(0.3)
        return wav_path.exists() and wav_path.stat().st_size > 200
    except Exception:
        return False


def generate_tts(text: str, stem: Path) -> Path | None:
    mp3 = stem.with_suffix(".mp3")
    wav = stem.with_suffix(".wav")
    if mp3.exists() and mp3.stat().st_size > 500:
        return mp3
    if wav.exists() and wav.stat().st_size > 500:
        return wav
    if _EDGE_OK:
        for attempt in range(3):
            try:
                asyncio.run(_edge_async(text, mp3))
                if mp3.exists() and mp3.stat().st_size > 500:
                    return mp3
                if mp3.exists():
                    mp3.unlink(missing_ok=True)
            except Exception as e:
                print(f"    edge_tts attempt {attempt + 1}: {e}")
                mp3.unlink(missing_ok=True)
                if attempt < 2:
                    time.sleep(1.5 * (attempt + 1))
        print("    edge_tts failed, switching to pyttsx3 fallback")

    if _pyttsx3_gen(text, wav):
        return wav

    return None


def get_duration(audio_path: Path) -> float:
    try:
        clip = AudioFileClip(str(audio_path))
        dur = clip.duration
        clip.close()
        return dur
    except Exception:
        return 0.0


def make_silent_wav(duration: float, path: Path) -> None:
    sr = 22050
    n = max(1, int(duration * sr))
    with wave.open(str(path), "wb") as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(sr)
        wf.writeframes(b"\x00\x00" * n)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE RENDERING
# ═══════════════════════════════════════════════════════════════════════════
def render_demo_slide(slide: SlideSpec, idx: int, total: int) -> Image.Image:
    """Render a custom demo slide (different layout per segment type)."""
    img = Image.new("RGB", (W, H), BG_TOP)
    _gradient(img)
    draw = ImageDraw.Draw(img)

    # ── title bar ────────────────────────────────────────────────────────
    draw.rectangle((0, 0, W, TITLE_H), fill=TITLE_BG)
    draw.line((0, TITLE_H, W, TITLE_H), fill=ACCENT, width=3)
    draw.text((48, 12), slide.title, font=FT, fill=ACCENT)
    draw.text((50, 58), slide.subtitle, font=FS, fill=TEXT_MUTE)

    badge = f"{idx} / {total}"
    bw, _ = _tsz(draw, badge, FKB)
    bx = W - bw - 50
    draw.rounded_rectangle((bx - 12, 20, W - 20, 64), radius=10, fill=PRIMARY, outline=PRIMARY, width=1)
    draw.text((bx, 28), badge, font=FKB, fill=(8, 16, 36))

    # ── content area ─────────────────────────────────────────────────────
    MARGIN = 42
    CTOP = TITLE_H + 42
    CBOT = CONTENT_H - 8
    LEFT_W = 600
    RSTART = LEFT_W + MARGIN * 2

    # left: bullets
    draw.rounded_rectangle(
        (MARGIN, CTOP, LEFT_W + MARGIN, CBOT), radius=16, fill=PANEL, outline=BORDER, width=2
    )
    y = CTOP + 18
    for bullet in slide.bullets:
        if y + 28 > CBOT:
            break
        draw.ellipse((MARGIN + 18, y + 7, MARGIN + 31, y + 20), fill=ACCENT)
        y = _draw_wrapped(draw, bullet, MARGIN + 42, y, FN, TEXT, LEFT_W - 62, 5) + 10

    # right: image or feature area
    IMG_W = W - RSTART - MARGIN
    IMG_H = CBOT - CTOP
    draw.rounded_rectangle(
        (RSTART, CTOP, W - MARGIN, CBOT), radius=16, fill=PANEL, outline=BORDER, width=2
    )

    # if image path is provided, paste it
    if slide.image_path and slide.image_path.exists():
        _paste_fit(img, slide.image_path, (RSTART + 8, CTOP + 8, W - MARGIN - 8, CBOT - 28))
        if slide.image_caption:
            draw.text(
                (RSTART + 12, CBOT - 24), slide.image_caption[:90], font=FK, fill=TEXT_MUTE,
            )
    else:
        # draw a mock feature display inside the right panel
        mid_x = RSTART + (IMG_W // 2)
        draw.text((mid_x - 60, CTOP + 30), slide.seg_id, font=_font(72, True), fill=ACCENT)
        if slide.seg_id == "#0":
            # Four pillar cards
            pillars = [
                ("KG", "知识图谱", "6 图谱 · 200+ 节点 · 11 布局", PRIMARY),
                ("AI", "AI Agent", "24 Agent · 4 大类别 · 秒级响应", GREEN),
                ("DT", "数字孪生", "学生镜像 · OBE 达成 · 5× 效率", ORANGE),
                ("4P", "四端同步", "Android/Windows/Web/HarmonyOS", PURPLE),
            ]
            for i, (icon, title, desc, color) in enumerate(pillars):
                px = RSTART + 20 + (i % 2) * 290
                py = CTOP + 110 + (i // 2) * 175
                draw_feature_card(draw, px, py, 270, 160, icon, title, desc, color)
        elif slide.seg_id == "#1":
            # Graph + layout icons
            graphs = [
                "01-课程图谱", "02-技术栈", "03-实验图谱",
                "04-项目图谱", "05-教学图谱", "06-学习图谱",
            ]
            for i, g in enumerate(graphs):
                gx = RSTART + 20 + (i % 3) * 210
                gy = CTOP + 100 + (i // 3) * 80
                draw.rounded_rectangle((gx, gy, gx + 190, gy + 60), radius=10, fill=(36, 58, 102), outline=ACCENT, width=1)
                draw.text((gx + 12, gy + 14), g, font=FKB, fill=ACCENT)
            # layout badges
            layouts = ["力导向", "树形", "同心圆", "网格", "螺旋", "分层"]
            lx = RSTART + 20
            for i, l in enumerate(layouts):
                lx, _ = draw_badge(draw, l, lx, CBOT - 40, PRIMARY)
                lx += 14
        elif slide.seg_id == "#2":
            # Agent categories
            cats = [
                ("教学辅导", 5, GREEN), ("内容生成", 3, PRIMARY),
                ("评价批阅", 5, ORANGE), ("辅助工具", 11, PURPLE),
            ]
            for i, (name, count, color) in enumerate(cats):
                cx = RSTART + 30 + (i % 2) * 290
                cy = CTOP + 100 + (i // 2) * 160
                draw.rounded_rectangle((cx, cy, cx + 260, cy + 130), radius=14, fill=PANEL, outline=color, width=2)
                draw.text((cx + 85, cy + 20), str(count), font=_font(48, True), fill=color)
                draw.text((cx + 75, cy + 75), name, font=_font(22, True), fill=TEXT)
        elif slide.seg_id == "#3":
            # Digital twin body metaphor
            parts = [
                ("骨架", "学习路径"), ("肌肉", "测验掌握"),
                ("心脏", "实验能力"), ("脉搏", "活动热力"),
                ("血液", "知识流动"), ("眼睛", "美德画像"),
            ]
            for i, (name, desc) in enumerate(parts):
                px = RSTART + 20 + (i % 3) * 200
                py = CTOP + 80 + (i // 3) * 170
                draw.rounded_rectangle((px, py, px + 180, py + 145), radius=12, fill=PANEL, outline=ACCENT, width=1)
                draw.text((px + 20, py + 20), name, font=_font(28, True), fill=ACCENT)
                draw.text((px + 20, py + 70), desc, font=FK, fill=TEXT_MUTE)
            # Radar placeholder
            draw.rounded_rectangle(
                (RSTART + 20, CTOP + 360, W - MARGIN - 20, CBOT - 10),
                radius=10, fill=(28, 46, 78), outline=GREEN, width=1,
            )
            draw.text((RSTART + 60, CTOP + 380), "五维雷达图 + OBE 达成度 (0.68 → 0.76)", font=FKB, fill=GREEN)
        elif slide.seg_id == "#4":
            # Platform cards
            platforms = [
                ("Windows", "55MB EXE", PRIMARY),
                ("Android", "35MB APK", GREEN),
                ("Web", "10MB 静态", ORANGE),
                ("HarmonyOS", "64MB HAP", PURPLE),
            ]
            for i, (plat, size, color) in enumerate(platforms):
                px = RSTART + 15 + (i % 2) * 310
                py = CTOP + 80 + (i // 2) * 190
                draw.rounded_rectangle((px, py, px + 285, py + 165), radius=14, fill=PANEL, outline=color, width=2)
                draw.text((px + 60, py + 25), plat, font=_font(38, True), fill=color)
                draw.text((px + 60, py + 85), size, font=_font(22, True), fill=TEXT)
                draw.text((px + 60, py + 115), "Flutter 一套代码", font=FK, fill=TEXT_MUTE)
            # Mic icon for voice
            draw.ellipse((RSTART + 350, CTOP + 310, RSTART + 460, CTOP + 420), fill=ACCENT)
            draw.text((RSTART + 380, CTOP + 343), "🎤", font=_font(48, True), fill=(0, 0, 0))
            draw.text((RSTART + 480, CTOP + 355), "语音导航 < 500ms", font=FKB, fill=ACCENT)
        elif slide.seg_id == "#5":
            # Big KPI cards
            kpis = [
                ("137", "学生使用", PRIMARY),
                ("4", "学期迭代", GREEN),
                ("92%", "满意度", ORANGE),
                ("12%", "达成度提升", PURPLE),
            ]
            for i, (val, label, color) in enumerate(kpis):
                kx = RSTART + 15 + (i % 2) * 310
                ky = CTOP + 60 + (i // 2) * 175
                draw_data_card(draw, kx, ky, 285, 150, val, label, color)
            # License + URLs
            draw.text((RSTART + 30, CBOT - 70), "MIT License", font=_font(24, True), fill=GREEN)
            draw.text((RSTART + 30, CBOT - 40), "gitee.com/osgisOne/mad-fd", font=FK, fill=TEXT_MUTE)
            draw.text((RSTART + 360, CBOT - 40), "github.com/dll/mad-fd", font=FK, fill=TEXT_MUTE)

    # ── subtitle area placeholder ─────────────────────────────────────────
    draw.rectangle((0, CONTENT_H, W, H), fill=SUB_BG)
    draw.line((0, CONTENT_H, W, CONTENT_H), fill=(40, 44, 68), width=1)
    return img


def add_subtitle_bar(base: Image.Image, sentence: str) -> Image.Image:
    result = base.copy()
    draw = ImageDraw.Draw(result)
    draw.rectangle((0, CONTENT_H, W, H), fill=SUB_BG)
    draw.line((0, CONTENT_H, W, CONTENT_H), fill=(48, 52, 78), width=1)
    if not sentence:
        return result
    font = FSU
    max_w = W - 180
    lines = _wrap(draw, sentence, font, max_w)[:3]
    lh = _tsz(draw, "M", font)[1] + 7
    total = len(lines) * lh
    y = CONTENT_H + (SUB_H - total) // 2
    for line in lines:
        lw, _ = _tsz(draw, line, font)
        draw.text(((W - lw) // 2, y), line, font=font, fill=WHITE)
        y += lh
    return result


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE DEFINITIONS
# ═══════════════════════════════════════════════════════════════════════════
SEGMENTS = [
    # ═══════════════════════════════════════════════════════════════════════
    # #0 开场 (3 slides, ~90s)
    # ═══════════════════════════════════════════════════════════════════════
    {
        "name": "#0 开场",
        "filename": "demo_00_opening",
        "slides": [
            SlideSpec(
                seg_id="#0",
                title="MAD-KGDT 整体定位",
                subtitle="移动应用开发知识图谱与数字孪生教学系统",
                image_path=_ss("01_login.png"),
                bullets=[
                    "4 大支柱：知识图谱 / AI Agent / 数字孪生 / Flutter 四端",
                    "解决 4 大痛点：知识碎片化、实践薄弱、个性化不足、数据孤岛",
                    "技术栈：Flutter + Dart + SQLite + 24 AI Agent + RAG",
                    "目标：教—学—练—评—管 全链路数字化",
                ],
                narration="MAD-KGDT，移动应用开发知识图谱与数字孪生教学系统，"
                         "面向《移动应用开发》课程的全平台教学平台。"
                         "系统以知识图谱、AI Agent、数字孪生和Flutter四端为四大支柱，"
                         "致力于解决知识碎片化、实践薄弱、个性化不足和数据孤岛四大教学痛点。"
                         "采用Flutter加Dart技术栈，集成24个AI Agent和RAG检索增强，"
                         "覆盖教、学、练、评、管全链路数字化。",
                voice_segments=[
                    "MAD-KGDT，移动应用开发知识图谱与数字孪生教学系统，",
                    "面向《移动应用开发》课程的全平台教学平台。",
                    "系统以知识图谱、AI Agent、数字孪生和Flutter四端为四大支柱，",
                    "致力于解决知识碎片化、实践薄弱、个性化不足和数据孤岛四大教学痛点。",
                    "采用Flutter加Dart技术栈，集成24个AI Agent和RAG检索增强。",
                ],
            ),
            SlideSpec(
                seg_id="#0",
                title="13.3 万行代码 · 263 文件",
                subtitle="四端架构 · 66 DB 表 · 133 页面",
                image_path=_ss("02_home.png"),
                bullets=[
                    "148,503 Dart 行 · 263 Dart 文件 · 66 数据库表",
                    "4 平台：Android / Windows / Web / HarmonyOS",
                    "24 AI 专业 Agent + Orchestrator 协调",
                    "6 张知识图谱 · 200+ 节点 · 11 种布局算法",
                ],
                narration="项目共十三万三千行Dart代码，263个Dart文件，66张数据库表。"
                         "覆盖Android、Windows、Web和HarmonyOS四个平台。"
                         "代码采用三层架构：页面层只负责展示，DAO层只访问数据库，"
                         "Service层组合业务逻辑。集成24个AI专业智能体和协调器，"
                         "内置6张知识图谱，200多个节点，支持11种布局算法。",
                voice_segments=[
                    "项目共13万3千行Dart代码，263个Dart文件，66张数据库表。",
                    "覆盖Android、Windows、Web和HarmonyOS四个平台。",
                    "代码采用三层架构：页面层展示，DAO层访问数据库，Service层组合业务逻辑。",
                    "集成24个AI专业智能体和协调器，内置6张知识图谱，200多个节点。",
                ],
            ),
            SlideSpec(
                seg_id="#0",
                title="四端构建产物一览",
                subtitle="Windows 55M / Android 35M / Web 10M / HarmonyOS 64M",
                image_path=_ss("10_settings.png"),
                bullets=[
                    "Windows：55MB EXE，libmpv 视频解码，noir 设计主题",
                    "Android：35MB APK，minSdk 21，arm64-v8a + armeabi-v7a",
                    "Web：10MB 静态站，base-href /mad-fd/ 适配 GitHub Pages",
                    "HarmonyOS：64MB HAP，OpenHarmony 调试签名，arm64-v8a 真机",
                ],
                narration="四个平台的构建产物均已优化到最小体积。"
                         "Windows桌面端55兆，集成libmpv视频解码库，采用noir编辑感设计主题。"
                         "Android包35兆，兼容arm64-v8a和armeabi-v7a两种架构。"
                         "Web端仅10兆静态资源，部署在GitHub Pages上。"
                         "HarmonyOS鸿蒙端64兆，已配好OpenHarmony调试签名，直接装到鸿蒙真机就能跑。"
                         "Flutter一套代码编译四端，每端只打包该平台需要的原生代码。",
                voice_segments=[
                    "四个平台的构建产物均已优化到最小体积。",
                    "Windows桌面端55兆，集成libmpv视频解码库，采用noir编辑感设计主题。",
                    "Android包35兆，兼容arm64-v8a和armeabi-v7a两种架构。",
                    "Web端仅10兆静态资源，部署在GitHub Pages上。",
                    "HarmonyOS鸿蒙端64兆，已配好OpenHarmony调试签名。",
                    "Flutter一套代码编译四端，每端只打包该平台需要的原生代码。",
                ],
            ),
        ],
    },
    # ═══════════════════════════════════════════════════════════════════════
    # #1 知识图谱 (3 slides, ~90s)
    # ═══════════════════════════════════════════════════════════════════════
    {
        "name": "#1 知识图谱",
        "filename": "demo_01_knowledge_graph",
        "slides": [
            SlideSpec(
                seg_id="#1",
                title="6 张图谱 · 200+ 节点",
                subtitle="11 种布局算法 · 双指缩放 · 学习路径计算",
                image_path=_ss("03_graph_list.png"),
                bullets=[
                    "6 张主题图谱：课程 / 技术栈 / 实验 / 项目 / 教学 / 学习",
                    "CustomPainter + InteractiveViewer 高性能渲染",
                    "11 种布局算法可切换（力导向 / 树形 / 同心圆 / 网格…）",
                    "模糊搜索 + 自动定位高亮 + 最短学习路径",
                ],
                narration="知识图谱模块包含6张主题图谱，涵盖课程体系、技术栈、实验项目、"
                         "项目实践、教学大纲和学习路径六大维度。"
                         "采用CustomPainter绘制节点边，InteractiveViewer支持双指缩放和鼠标拖拽，"
                         "即使200多个节点同时渲染也能保持60帧流畅交互。"
                         "支持11种布局算法一键切换，包括力导向布局、树形布局、同心圆布局、"
                         "网格布局和螺旋布局等。搜索框支持模糊匹配，自动定位并高亮目标节点，"
                         "还能计算两个节点之间的最短学习路径。",
                voice_segments=[
                    "知识图谱模块包含6张主题图谱，涵盖课程体系、技术栈、实验项目、",
                    "项目实践、教学大纲和学习路径六大维度。",
                    "采用CustomPainter绘制节点边，InteractiveViewer支持双指缩放和鼠标拖拽，",
                    "即使200多个节点同时渲染也能保持60帧流畅交互。",
                    "支持11种布局算法一键切换，包括力导向、树形、同心圆、网格和螺旋布局等。",
                    "搜索框支持模糊匹配，自动定位并高亮目标节点，",
                    "还能计算两个节点之间的最短学习路径。",
                ],
            ),
            SlideSpec(
                seg_id="#1",
                title="节点详情 · 资源联动",
                subtitle="视频 / 文档 / 测验 一站直达",
                image_path=_ss("04_graph_detail.png"),
                bullets=[
                    "点击节点 → 详情面板 4 入口（视频/文档/测验/相关节点）",
                    "节点跨资源联通：学习资源与知识概念双向映射",
                    "收藏节点 → 自动聚合到学习计划",
                    "学习进度与图谱节点同步跟踪",
                ],
                narration="点击任意节点会弹出详情面板，提供视频、文档、测验和相关节点四个学习入口。"
                         "比如点击Flutter状态管理节点，右侧直接展示对应的教学视频、PDF课件、"
                         "章节测验题和相关知识点。节点与学习资源之间通过resource_chapter_mapping表双向映射，"
                         "收藏节点后自动聚合到学习计划。学习进度与图谱节点同步跟踪，"
                         "学完一个节点后图谱上自动打勾标记，形成完整的学习闭环。",
                voice_segments=[
                    "点击任意节点会弹出详情面板，",
                    "提供视频、文档、测验和相关节点四个学习入口。",
                    "比如点击Flutter状态管理节点，右侧直接展示对应的教学视频、PDF课件和测验题。",
                    "节点与学习资源之间通过resource_chapter_mapping表双向映射。",
                    "收藏节点后自动聚合到学习计划，学习进度与图谱节点同步跟踪。",
                ],
            ),
            SlideSpec(
                seg_id="#1",
                title="图谱布局算法原理",
                subtitle="力导向 · 树形 · 同心圆 · 网格 · 螺旋",
                bullets=[
                    "力导向：模拟物理弹簧系统，自动散开重叠节点",
                    "树形：按层级父子关系垂直展开，适合课程体系",
                    "同心圆：按节点权重半径排列，核心概念在圆心",
                    "网格 + 螺旋：均匀分布和环绕排列，适合等权节点",
                ],
                narration="图谱模块实现了多种成熟的布局算法。"
                         "力导向布局模拟物理弹簧系统，节点之间有引力也有斥力，"
                         "自动把重叠的节点弹开，适合浏览复杂关系网络。"
                         "树形布局按照父子层级关系垂直展开，从根节点到叶子节点层次分明，"
                         "最适合展示课程章节体系。同心圆布局按节点权重分配到不同半径的圆环上，"
                         "核心概念在最内圈。网格和螺旋布局则适用于等权节点的均匀排列。"
                         "所有布局算法都可以在设置面板中一键切换，切换时还有平滑过渡动画。",
                voice_segments=[
                    "图谱模块实现了多种成熟的布局算法。",
                    "力导向布局模拟物理弹簧系统，自动把重叠的节点弹开。",
                    "树形布局按照父子层级关系垂直展开，最适合展示课程章节体系。",
                    "同心圆布局按节点权重分配到不同半径的圆环上。",
                    "网格和螺旋布局适用于等权节点的均匀排列。",
                    "所有布局算法都可以在设置面板中一键切换，切换时还有平滑过渡动画。",
                ],
            ),
        ],
    },
    # ═══════════════════════════════════════════════════════════════════════
    # #2 AI 24 Agent (3 slides, ~95s)
    # ═══════════════════════════════════════════════════════════════════════
    {
        "name": "#2 AI 24 Agent",
        "filename": "demo_02_ai_agents",
        "slides": [
            SlideSpec(
                seg_id="#2",
                title="24 个专业 AI Agent",
                subtitle="4 大类别 · 秒级响应 · 10000+ 调用记录",
                image_path=_ss("05_agents.png"),
                bullets=[
                    "教学辅导 5 个：tutor / quiz / lab / path / learning",
                    "内容生成 3 个：course_gen / courseware / doc_converter",
                    "评价批阅 5 个：lab_grading / works_grading / assessment_grading",
                    "辅助工具 11 个：voice / graph / safety / ethics / repo …",
                ],
                narration="系统集成了24个专业AI智能体，统一由AgentRegistry单例管理。"
                         "每个Agent由AgentConfig定义人格、工具列表和用例，"
                         "通过BaseAgent基类统一处理会话管理和AI推理。"
                         "教学辅导类5个，包括课堂助教tutor、测验生成quiz、实验指导lab、"
                         "路径规划path和学习助手learning。"
                         "内容生成类3个，course_gen一键生课、courseware课件生成、doc_converter文档转换。"
                         "评价批阅类5个，覆盖实验、作品和考核的AI自动批阅。"
                         "辅助工具类11个，包括语音导航voice、图谱分析graph、安全审查safety、"
                         "学术伦理ethics、仓库分析repo等。",
                voice_segments=[
                    "系统集成了24个专业AI智能体，统一由AgentRegistry单例管理。",
                    "教学辅导类5个，包括课堂助教、测验生成、实验指导、路径规划和学生学习助手。",
                    "内容生成类3个，一键生课、课件生成和文档转换。",
                    "评价批阅类5个，覆盖实验、作品和考核的AI自动批阅。",
                    "辅助工具类11个，包括语音导航、图谱分析、安全审查、学术伦理和仓库分析等。",
                ],
            ),
            SlideSpec(
                seg_id="#2",
                title="多智能体协同 · 全程留痕",
                subtitle="AgentChatOverlay · 吉祥物入口 · 调用统计",
                image_path=_ss("05_agents.png"),
                bullets=[
                    "吉祥物悬浮按钮 → AgentChatOverlay 全局浮层",
                    "费曼 + 苏格拉底 + 类比 多教学法自动切换",
                    "course_gen：输入主题 → 自动出题 + 生成图谱 + 实验框架",
                    "AI 调用统计面板：24 Agent 排行 + 平均耗时",
                ],
                narration="通过吉祥物悬浮按钮可打开全局Agent对话浮层，支持7种导航动作。"
                         "辅导回答不是单一风格，而是根据问题类型自动切换费曼技巧、"
                         "苏格拉底问答法和类比教学法。一键生课Agent输入学科主题，"
                         "自动生成12节点知识图谱、30道测验题和实验框架，几分钟完成传统需要"
                         "几周的教学设计工作。AI调用统计面板记录所有Agent的调用排行和平均耗时，"
                         "全程留痕，目前已有超过一万条调用记录。",
                voice_segments=[
                    "通过吉祥物悬浮按钮可打开全局Agent对话浮层，支持7种导航动作。",
                    "辅导回答根据问题类型自动切换费曼技巧、苏格拉底问答法和类比教学法。",
                    "一键生课Agent输入学科主题，自动生成知识图谱、测验题和实验框架。",
                    "AI调用统计面板记录所有Agent的调用排行和平均耗时，全程留痕。",
                ],
            ),
            SlideSpec(
                seg_id="#2",
                title="RAG 检索增强 · Orchestrator 协调",
                subtitle="知识库 + Agent Chain = 高质量回答",
                bullets=[
                    "RAG 检索：基于课程内容知识库，自动注入相关文档片段",
                    "Orchestrator 协调器：编排多 Agent 工作链",
                    "例如实验批阅链：safety 审查 → grading 批阅 → ethics 伦理建议",
                    "AgentCallLog 记录每步耗时，仪表板按 chainId 聚合分析",
                ],
                narration="RAG检索增强服务是智能体的知识后盾。"
                         "每次提问时，RagService自动从课程内容知识库中检索最相关的文档片段，"
                         "注入到AI的提示词上下文中，确保回答基于课程教材而不是通用知识。"
                         "Orchestrator协调器更进一层，它可以编排多个Agent的工作链。"
                         "比如实验批阅链：先由safety_agent进行内容安全审查，"
                         "再由主批阅Agent打分，最后由ethics_agent给出学术伦理建议。"
                         "整个过程通过AgentCallLog记录每步耗时，仪表板按chainId聚合分析。",
                voice_segments=[
                    "RAG检索增强服务自动从课程内容知识库中检索相关文档片段。",
                    "注入到AI的提示词上下文中，确保回答基于课程教材。",
                    "Orchestrator协调器编排多个Agent的工作链。",
                    "比如实验批阅链：安全审查、AI批阅、伦理建议三步协同。",
                    "AgentCallLog记录每步耗时，仪表板按chainId聚合分析。",
                ],
            ),
        ],
    },
    # ═══════════════════════════════════════════════════════════════════════
    # #3 数字孪生 (4 slides, ~130s)
    # ═══════════════════════════════════════════════════════════════════════
    {
        "name": "#3 数字孪生",
        "filename": "demo_03_digital_twin",
        "slides": [
            SlideSpec(
                seg_id="#3",
                title="学生数字孪生 · 身体隐喻",
                subtitle="十一层信息 → 可感知的'另一个我'",
                image_path=_ss("06_twin.png"),
                bullets=[
                    "五维雷达图 + 30 天学习热力图",
                    "骨架→肌肉→心脏→脉搏→血液→眼睛 逐层展示",
                    "美德精灵画像：抽象数据 → 具象化自我认知",
                    "三级风险预警：不活跃 / 低分 / 下滑",
                ],
                narration="学生数字孪生模块是系统的核心创新之一。"
                         "它通过五维能力雷达图和30天学习热力图，把碎片化的学习行为数据"
                         "转化为直观的可视化仪表盘。设计上采用身体隐喻，从骨架、肌肉、心脏、"
                         "脉搏、血液到眼睛共十一层信息逐层展开。骨架是学习路径完成度，"
                         "肌肉是测验掌握程度，心脏是实验能力评分，脉搏是每日活动热力，"
                         "血液是知识流动网络，眼睛是美德精灵画像——把抽象数据变成可感知的另一个我。"
                         "同时三级风险预警机制自动标注不活跃、低分和成绩下滑的学生。",
                voice_segments=[
                    "学生数字孪生模块是系统的核心创新之一。",
                    "通过五维能力雷达图和30天学习热力图，把学习行为数据可视化。",
                    "采用身体隐喻，从骨架、肌肉、心脏、脉搏、血液到眼睛共十一层信息逐层展开。",
                    "把抽象数据变成可感知的另一个我。",
                    "三级风险预警自动标注不活跃、低分和成绩下滑的学生。",
                ],
            ),
            SlideSpec(
                seg_id="#3",
                title="五维雷达 · 热力图 · 数据采集",
                subtitle="学习行为 → 孪生数据 → 可视化呈现",
                bullets=[
                    "五维：知识广度 / 测验掌握 / 实验能力 / 活跃度 / 协作贡献",
                    "热力图：30 天逐日行为密度，颜色深浅反映投入程度",
                    "数据采集源：登录日志、页面浏览、测验记录、实验提交、协作消息",
                    "twin_service.dart 每日凌晨聚合计算，延迟不超过 1 小时",
                ],
                narration="五维雷达图从知识广度、测验掌握、实验能力、活跃度和协作贡献五个角度"
                         "评估学生的综合表现。30天热力图以日历形式展示每天的学习行为密度，"
                         "颜色越深表示当天学习投入越高。数据采集来源包括登录日志、"
                         "页面浏览行为、测验答题记录、实验提交记录和协作消息等六个维度。"
                         "数据由twin_service.dart在每日凌晨自动聚合计算，延迟不超过1小时，"
                         "保证学生每次打开数字孪生页面看到的都是最新的自己。",
                voice_segments=[
                    "五维雷达图从知识广度、测验掌握、实验能力、活跃度和协作贡献五方面综合评估。",
                    "30天热力图以日历形式展示每天的学习行为密度。",
                    "数据采集来源包括登录日志、页面浏览、测验答题、实验提交和协作消息。",
                    "twin_service.dart每日凌晨自动聚合，延迟不超过1小时。",
                ],
            ),
            SlideSpec(
                seg_id="#3",
                title="教师班级看板 · AI 批阅",
                subtitle="班级全景 · 五维评分 · 反馈周期 1 周 → 分钟级",
                image_path=_ss("07_classroom.png"),
                bullets=[
                    "班级看板：薄弱知识点 Top5 + 成绩分布",
                    "实验 AI 五维批阅：功能/代码/报告/分析/创新",
                    "反馈周期从 1 周 → 分钟级，效率提升 5 倍",
                    "全班自动批改，教师只审，不再逐份打分",
                ],
                narration="教师端班级看板提供班级全景视图，自动识别薄弱知识点前五名和成绩分布。"
                         "实验报告采用AI五维批阅，从功能完整性、代码质量、报告规范性、"
                         "分析深度和创新性五个维度自动评分并给出评语。"
                         "传统教学中，教师批阅40份实验报告需要整整一周时间，"
                         "现在全班提交后AI在30分钟内完成初评，教师只需审核确认，"
                         "反馈周期从一周缩短到分钟级，效率提升五倍。"
                         "lab_grading_agent已累计批阅超过80份实验报告。",
                voice_segments=[
                    "教师端班级看板提供班级全景视图，自动识别薄弱知识点前五名。",
                    "AI五维批阅从功能完整性、代码质量、报告规范性、分析深度和创新性自动评分。",
                    "传统批阅40份实验报告需要一周，现在AI 30分钟完成初评。",
                    "反馈周期从一周缩短到分钟级，效率提升五倍。",
                ],
            ),
            SlideSpec(
                seg_id="#3",
                title="OBE 达成度 · 工程认证",
                subtitle="三维加权 · 四步计算 · 持续改进策略",
                image_path=_ss("08_obe.png"),
                bullets=[
                    "OBE 达成度 0.68 → 0.76 ↑12%",
                    "三维分数：平时 20% + 实验 30% + 考核 50%",
                    "四步计算：权重分配 → 均值计算 → 达成判定 → 雷达出图",
                    "AI 持续改进策略：重大干预 / 中度改进 / 巩固 / 优化",
                ],
                narration="OBE达成度从最初的0.68提升到了0.76，提升了12个百分点。"
                         "采用三维加权计算匹配工程教育认证标准：平时成绩占20%，"
                         "实验占30%，期末考核占50%。四步计算流程包括权重分配、"
                         "均值计算、达成判定和雷达出图。四个课程目标分别计算达成度，"
                         "低于0.6标红预警。AI自动生成持续改进策略，分为四个级别："
                         "重大干预用于低于0.6的目标，中度改进用于0.6到0.7之间，"
                         "巩固用于0.7到0.8之间，优化用于0.8以上。"
                         "改进建议直接关联具体教学环节，比如增加实验课时占比、"
                         "强化过程性考核等可操作建议。",
                voice_segments=[
                    "OBE达成度从最初的0.68提升到了0.76，提升了12个百分点。",
                    "三维加权：平时20%，实验30%，考核50%，匹配工程教育认证标准。",
                    "四步计算流程：权重分配、均值计算、达成判定和雷达出图。",
                    "AI持续改进策略分四级：重大干预、中度改进、巩固和优化。",
                    "改进建议直接关联具体教学环节，可操作可落地。",
                ],
            ),
        ],
    },
    # ═══════════════════════════════════════════════════════════════════════
    # #4 四端 + 同步 + 语音 (3 slides, ~85s)
    # ═══════════════════════════════════════════════════════════════════════
    {
        "name": "#4 四端 + 同步 + 语音",
        "filename": "demo_04_multiplatform",
        "slides": [
            SlideSpec(
                seg_id="#4",
                title="Flutter 一套代码 · 四端运行",
                subtitle="Android 35M / Windows 55M / Web 10M / HarmonyOS 64M",
                image_path=_ss("02_home.png"),
                bullets=[
                    "统一 UI 风格：noir 编辑感设计系统",
                    "Android / Windows / Web / HarmonyOS 四端真机可跑",
                    "同一套 Dart 代码 → 4 个平台原生编译",
                    "鸿蒙端 OpenHarmony 调试签名 HAP 已就绪",
                ],
                narration="Flutter一套代码编译四个平台，包大小因平台原生依赖而异。"
                         "Android包35兆，Windows桌面55兆，Web静态10兆，"
                         "HarmonyOS鸿蒙64兆。四端保持统一的noir编辑感设计风格——"
                         "深色背景搭配琥珀色强调色，模拟经典代码编辑器的视觉体验。"
                         "跨平台适配策略采用条件导入native_stub.dart的模式，"
                         "文件系统、路径服务等平台相关功能各有独立实现。"
                         "鸿蒙端已配置OpenHarmony调试签名，可直接安装到真机，"
                         "目前仅支持arm64-v8a架构。",
                voice_segments=[
                    "Flutter一套代码编译四个平台。",
                    "Android 35兆，Windows 55兆，Web 10兆，HarmonyOS 64兆。",
                    "四端保持统一的noir编辑感设计风格，深色背景搭配琥珀色强调色。",
                    "跨平台适配采用条件导入native_stub.dart的模式。",
                    "鸿蒙端已配置OpenHarmony调试签名，可直接安装到真机。",
                ],
            ),
            SlideSpec(
                seg_id="#4",
                title="语音导航 · 全链路 < 500ms",
                subtitle="讯飞 STT → VoiceAgent → NavigationService",
                image_path=_ss("09_lab.png"),
                bullets=[
                    "吉祥物语音按钮 → 讯飞 WebSocket STT 实时转写",
                    "4 层路由：快速路径 → Tab 映射 → 子页面 → AI 兜底",
                    "VoiceAgent 意图识别 → JSON 结构化导航指令",
                    "支持 7 种导航动作：navigate_tab / navigate_sub_page / go_back …",
                ],
                narration="语音导航是系统的亮点交互方式。"
                         "用户点击吉祥物语音按钮后，语音通过讯飞WebSocket API实时转写为文本，"
                         "再经过VoiceAgent的AI意图识别，解析出结构化导航指令。"
                         "系统采用四层路由机制确保准确率：第一层快速路径处理返回和退出等固定指令，"
                         "第二层NavigationService的Tab映射匹配首页、图谱等一级页面，"
                         "第三层匹配30多个子页面，第四层由VoiceAgent做AI兜底。"
                         "全链路从语音结束到页面跳转不超过500毫秒。",
                voice_segments=[
                    "语音导航通过吉祥物语音按钮触发。",
                    "语音通过讯飞WebSocket API实时转写为文本。",
                    "VoiceAgent进行AI意图识别，解析出结构化导航指令。",
                    "四层路由机制确保准确率，全链路不超过500毫秒。",
                ],
            ),
            SlideSpec(
                seg_id="#4",
                title="Gitee 双向同步 · 零服务器",
                subtitle="学生 → JSON 文件 → Gitee 仓库 → 教师",
                image_path=_ss("09_lab.png"),
                bullets=[
                    "同步流程：学生设备 → JSON 文件 → Gitee 仓库 → 教师设备",
                    "task_id 重映射：跨设备自增 ID 不同，通过 title 自然键匹配",
                    "批改数据保护：已有 score/feedback 的不被覆盖",
                    "即时同步：提交实验后立即触发 uploadStudentData()",
                ],
                narration="数据同步通过Gitee仓库实现零服务器双向同步。"
                         "学生提交实验或完成测验后，数据自动序列化为JSON文件，"
                         "通过Gitee API上传到仓库，教师端定时拉取更新。"
                         "由于每台设备的lab_tasks自增ID不同，同步时通过title字段做自然键匹配，"
                         "构建ID映射表。已批改的实验提交带有score和feedback字段，"
                         "导入时受到保护不被覆盖。而且提交实验报告后立即触发上传，"
                         "不等定时器，确保教师尽快看到学生提交。",
                voice_segments=[
                    "数据同步通过Gitee仓库实现零服务器双向同步。",
                    "学生提交实验后数据自动序列化为JSON文件上传到仓库。",
                    "跨设备ID通过title自然键匹配构建映射表。",
                    "已批改数据受到保护不被覆盖。",
                    "提交后立即触发上传，不等定时器。",
                ],
            ),
        ],
    },
    # ═══════════════════════════════════════════════════════════════════════
    # #5 收尾 (2 slides, ~50s)
    # ═══════════════════════════════════════════════════════════════════════
    {
        "name": "#5 收尾",
        "filename": "demo_05_closing",
        "slides": [
            SlideSpec(
                seg_id="#5",
                title="137 名学生 · 4 个学期 · 92% 满意度",
                subtitle="2024 秋 → 2026 春 四学期持续迭代",
                bullets=[
                    "课程目标达成度 0.68 → 0.76 ↑12%",
                    "AI 批阅效率 5 倍提升，反馈周期 1 周 → 分钟级",
                    "数字孪生学生认可度 87%",
                    "MIT 开源，欢迎兄弟院校共建",
                ],
                narration="项目自2024年秋季到2026年春季，历经四个学期持续迭代。"
                         "累计137名学生使用，课程满意度92%。"
                         "课程目标达成度从0.68提升到0.76，提升12个百分点。"
                         "AI批阅使反馈效率提升5倍。"
                         "数字孪生学生认可度达到87%。"
                         "这些数据来自滁州学院软件工程专业四学期的真实教学实践。",
                voice_segments=[
                    "项目自2024年秋季到2026年春季，历经四个学期持续迭代。",
                    "累计137名学生使用，课程满意度92%。",
                    "课程目标达成度从0.68提升到0.76，提升12个百分点。",
                    "AI批阅使反馈效率提升5倍，数字孪生学生认可度达到87%。",
                ],
            ),
            SlideSpec(
                seg_id="#5",
                title="未来规划 · 开源共建",
                subtitle="MIT License · Gitee + GitHub 双平台",
                bullets=[
                    "Gitee: gitee.com/osgisOne/mad-fd   (主仓库)",
                    "GitHub: github.com/dll/mad-fd   (镜像)",
                    "Web 在线: dll.github.io/mad-fd",
                    "6 轮 AI 自审，评分 4.4/5，完整 audit 报告公开",
                ],
                narration="项目采用MIT开源协议，在Gitee和GitHub双平台维护。"
                         "Gitee主仓库同步师生数据，GitHub镜像做Web部署和社区协作。"
                         "Web在线版可以直接在浏览器访问体验，无需安装。"
                         "项目经历了六轮AI自审，评分从3.6提升到4.4，"
                         "每轮审查报告和修改记录都完整保存在docs/目录下。"
                         "后续计划包括增加更多课程适配、完善鸿蒙原生体验、"
                         "以及建设开源社区，欢迎兄弟院校共建。",
                voice_segments=[
                    "项目采用MIT开源协议，Gitee和GitHub双平台维护。",
                    "Web在线版可直接在浏览器访问体验，无需安装。",
                    "六轮AI自审评分从3.6提升到4.4，审查报告公开。",
                    "后续计划增加更多课程适配和鸿蒙原生体验，欢迎兄弟院校共建。",
                ],
            ),
        ],
    },
]


# ═══════════════════════════════════════════════════════════════════════════
# SRT
# ═══════════════════════════════════════════════════════════════════════════
def _ts(s: float) -> str:
    h = int(s // 3600)
    m = int((s % 3600) // 60)
    sc = int(s % 60)
    ms = int(round((s - int(s)) * 1000))
    return f"{h:02d}:{m:02d}:{sc:02d},{ms:03d}"


def write_srt(entries: list[tuple[str, float, float]], srt_path: Path) -> None:
    lines: list[str] = []
    for i, (text, start, end) in enumerate(entries, 1):
        lines += [str(i), f"{_ts(start)} --> {_ts(end)}", text, ""]
    srt_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"SRT  → {srt_path}  ({len(entries)} entries)")


# ═══════════════════════════════════════════════════════════════════════════
# SCRIPT MARKDOWN
# ═══════════════════════════════════════════════════════════════════════════
def write_script(slides_map: dict[str, list[SlideSpec]], scp_path: Path) -> None:
    lines = [
        "# MAD-KGDT 案例评比 Demo 视频脚本",
        "",
        "## 技术特性",
        "",
        "- **edge_tts** (zh-CN-XiaoxiaoNeural) 高质量中文语音",
        "- **moviepy** 精确帧级同步：每句独立帧，一次性编码，彻底消除漂移",
        "- 字幕逐句烧录到画面；SRT 时间戳来自真实 MP3 时长",
        "",
        "## 分段结构",
        "",
    ]
    seg_counter = 1
    for seg_name, seg_slides in slides_map.items():
        for s in seg_slides:
            lines.append(f"{seg_counter}. **{s.title}** — {s.subtitle}")
            seg_counter += 1
    lines += ["", "## 讲解稿", ""]
    seg_counter = 1
    for seg_name, seg_slides in slides_map.items():
        for s in seg_slides:
            lines += [f"### {seg_counter}. {s.title}", f"> {s.subtitle}", "", f"**旁白：** {s.narration}", ""]
            seg_counter += 1
    scp_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"Script → {scp_path}")


# ═══════════════════════════════════════════════════════════════════════════
# MAIN VIDEO PIPELINE (per segment)
# ═══════════════════════════════════════════════════════════════════════════
def build_segment(
    seg_name: str, seg_filename: str, slides: list[SlideSpec],
) -> tuple[Path | None, bool]:
    """Build one video segment. Returns (video_path, ok)."""
    print(f"\n{'=' * 60}")
    print(f"  Segment: {seg_name}")
    print(f"{'=' * 60}")

    if not _MOVIEPY_OK:
        print("ERROR: moviepy is required. pip install moviepy==1.0.3")
        return None, False

    video_path = OUT_DIR / f"{seg_filename}.mp4"
    srt_path = DEMO_DIR / f"{seg_filename}.srt"

    all_clips: list = []
    srt_entries: list[tuple[str, float, float]] = []
    t = 0.0
    total = len(slides)

    for si, slide in enumerate(slides, 1):
        print(f"\n  [{si:02d}/{total}] {slide.title}")

        # render base slide
        base_img = render_demo_slide(slide, si, total)
        base_path = SLIDES_DIR / f"{seg_filename}_{si:02d}.png"
        base_path.parent.mkdir(parents=True, exist_ok=True)
        base_img.save(base_path, "PNG")

        # split narration into sentences
        if slide.voice_segments:
            sentences = [s.strip() for s in slide.voice_segments if s.strip()]
        else:
            text = slide.narration.strip()
            chunks = []
            cur = ""
            for ch in text:
                cur += ch
                if ch in "。！？；":
                    t2 = cur.strip()
                    if t2:
                        chunks.append(t2)
                    cur = ""
            if cur.strip():
                chunks.append(cur.strip())
            merged = []
            i = 0
            while i < len(chunks):
                if i + 1 < len(chunks) and len(chunks[i]) < 8:
                    merged.append(chunks[i] + chunks[i + 1])
                    i += 2
                else:
                    merged.append(chunks[i])
                    i += 1
            sentences = merged if merged else [text]

        nsent = len(sentences)
        for ki, sentence in enumerate(sentences):
            is_last = ki == nsent - 1
            tail = SLIDE_TAIL if is_last else CLIP_TAIL

            audio_stem = AUDIO_DIR / f"{seg_filename}_{si:02d}_{ki + 1:02d}"
            print(f"    TTS [{ki + 1}/{nsent}] → {sentence[:55]}")

            audio_path = generate_tts(sentence, audio_stem)

            if audio_path is None or not audio_path.exists():
                audio_path = audio_stem.with_suffix(".wav")
                estimated = max(1.8, len(sentence) * 0.14)
                make_silent_wav(estimated, audio_path)

            real_dur = get_duration(audio_path)
            if real_dur < 0.05:
                real_dur = max(1.8, len(sentence) * 0.14)

            # per-sentence frame
            frame_img = add_subtitle_bar(base_img, sentence)
            frame_path = SENT_DIR / f"{seg_filename}_{si:02d}_{ki + 1:02d}.png"
            frame_path.parent.mkdir(parents=True, exist_ok=True)
            frame_img.save(frame_path, "PNG")

            # moviepy clip
            clip_dur = real_dur + tail
            try:
                audio_clip = AudioFileClip(str(audio_path))
                vid_clip = (
                    ImageClip(str(frame_path))
                    .set_duration(clip_dur)
                    .set_audio(audio_clip)
                )
            except Exception as e:
                print(f"    clip error: {e} (silent)")
                vid_clip = ImageClip(str(frame_path)).set_duration(clip_dur)

            all_clips.append(vid_clip)
            srt_entries.append((sentence, t, t + real_dur))
            t += clip_dur

    # write SRT
    write_srt(srt_entries, srt_path)

    # concatenate and encode
    print(f"\n  Concatenating {len(all_clips)} clips (total ≈ {t:.1f}s)...")
    try:
        final = concatenate_videoclips(all_clips, method="compose")
        print(f"  Encoding → {video_path}")
        final.write_videofile(
            str(video_path),
            fps=FPS,
            codec="libx264",
            audio_codec="aac",
            audio_bitrate="192k",
            temp_audiofile=str(TEMP_DIR / f"{seg_filename}_tmp.m4a"),
            remove_temp=True,
            verbose=False,
            write_logfile=False,
            threads=4,
        )
        final.close()
        print(f"  ✓ {video_path}")
        return video_path, True
    except Exception as e:
        print(f"  ERROR: {e}")
        import traceback
        traceback.print_exc()
        return None, False


# ═══════════════════════════════════════════════════════════════════════════
# PPTX BUILDER (with real screenshots)
# ═══════════════════════════════════════════════════════════════════════════
try:
    from pptx import Presentation as PptxPres
    from pptx.util import Inches, Pt as PptxPt
    from pptx.dml.color import RGBColor

    _PPTX_OK = True
except Exception:
    PptxPres = None
    Inches = Pt = None
    _PPTX_OK = False


def build_pptx(slides_map: dict[str, list[SlideSpec]], pptx_path: Path) -> None:
    """Build PPTX with screenshot slides + script in notes."""
    if not _PPTX_OK or PptxPres is None:
        print("python-pptx not available, skipping PPTX.")
        return
    prs = PptxPres()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for seg_name, seg_slides in slides_map.items():
        for slide in seg_slides:
            sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            # Title
            tb = sl.shapes.add_textbox(Inches(0.38), Inches(0.12), Inches(10.5), Inches(0.52))
            p = tb.text_frame.paragraphs[0]
            p.text = f"[{slide.seg_id}] {slide.title}"
            p.font.size = PptxPt(28)
            p.font.bold = True
            # Subtitle
            tb2 = sl.shapes.add_textbox(Inches(0.40), Inches(0.62), Inches(12.2), Inches(0.34))
            p2 = tb2.text_frame.paragraphs[0]
            p2.text = slide.subtitle
            p2.font.size = PptxPt(14)
            p2.font.color.rgb = RGBColor(0x66, 0x7E, 0xEA)
            # Screenshot image
            if slide.image_path and slide.image_path.exists():
                sl.shapes.add_picture(
                    str(slide.image_path),
                    Inches(0.28), Inches(1.02),
                    width=Inches(12.78),
                    height=Inches(6.18),
                )
            else:
                tb3 = sl.shapes.add_textbox(Inches(3), Inches(3), Inches(7), Inches(1))
                p3 = tb3.text_frame.paragraphs[0]
                p3.text = "（请在此处插入应用截图）"
                p3.font.size = PptxPt(20)
                p3.font.italic = True
            # Narration in notes
            notes_slide = sl.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = f"[{seg_name}] {slide.title}\n{slide.subtitle}\n\n旁白：\n{slide.narration}"
    prs.save(str(pptx_path))
    print(f"PPTX → {pptx_path}  ({len(prs.slides)} slides)")


# ═══════════════════════════════════════════════════════════════════════════
# ENTRY POINT
# ═══════════════════════════════════════════════════════════════════════════
def main() -> None:
    for d in [OUT_DIR, DEMO_DIR, SLIDES_DIR, SENT_DIR, AUDIO_DIR, TEMP_DIR]:
        d.mkdir(parents=True, exist_ok=True)

    print("=" * 60)
    print("  MAD-KGDT 案例评比 Demo 视频生成器")
    print("  edge_tts + moviepy  严格帧级同步")
    print("=" * 60)

    results: list[tuple[str, Path | None, bool]] = []
    for seg in SEGMENTS:
        vpath, ok = build_segment(seg["name"], seg["filename"], seg["slides"])
        results.append((seg["name"], vpath, ok))

    # Write combined script
    all_slides: dict[str, list[SlideSpec]] = {}
    for seg in SEGMENTS:
        all_slides[seg["name"]] = seg["slides"]
    write_script(all_slides, DEMO_DIR / "demo_script_output.md")

    # Build PPTX
    pptx_path = OUT_DIR / "demo_slides.pptx"
    build_pptx(all_slides, pptx_path)

    # Summary
    print("\n" + "=" * 60)
    print("  Generation Summary")
    print("=" * 60)
    for name, vpath, ok in results:
        status = "✓" if ok else "✗"
        size = ""
        if ok and vpath:
            mb = vpath.stat().st_size / (1024 * 1024)
            size = f" ({mb:.1f} MB)"
        print(f"  {status} {name}: {vpath}{size}")

    all_ok = all(r[2] for r in results)
    if all_ok:
        # Write a concat script for merging
        concat_path = DEMO_DIR / "concat.txt"
        lines = []
        for seg in SEGMENTS:
            mp4 = OUT_DIR / f"{seg['filename']}.mp4"
            if mp4.exists():
                lines.append(f"file '{mp4}'")
        concat_path.write_text("\n".join(lines), encoding="utf-8")
        print(f"\n  Concat list → {concat_path}")
        print("  Merge with:  ffmpeg -f concat -safe 0 -i concat.txt -c copy demo_full.mp4")
    else:
        print("\n  Some segments failed. Check errors above.")

    print("=" * 60)


if __name__ == "__main__":
    main()
