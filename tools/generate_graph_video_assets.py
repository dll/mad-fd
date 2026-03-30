from __future__ import annotations

import math
import os
import subprocess
import textwrap
from dataclasses import dataclass
from pathlib import Path
from typing import List, Tuple

from PIL import Image, ImageDraw, ImageFont

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
DIAGRAMS_DIR = DOCS_DIR / "diagrams"
VIDEO_DIR = DOCS_DIR / "video"
SLIDES_DIR = VIDEO_DIR / "slides"
GENERATED_DIR = VIDEO_DIR / "generated"
VIDEO_OUTPUT_DIR = ROOT / "video_output"

PUML_ARCH_PATH = DIAGRAMS_DIR / "knowledge_graph_feature_architecture.puml"
PUML_FLOW_PATH = DIAGRAMS_DIR / "knowledge_graph_feature_flow.puml"
ARCH_PNG_PATH = GENERATED_DIR / "knowledge_graph_feature_architecture.png"
FLOW_PNG_PATH = GENERATED_DIR / "knowledge_graph_feature_flow.png"
SCRIPT_MD_PATH = VIDEO_DIR / "video_script.md"
PPTX_PATH = VIDEO_OUTPUT_DIR / "知识图谱核心功能_图谱功能演示.pptx"
VIDEO_PATH = VIDEO_OUTPUT_DIR / "知识图谱核心功能_01_图谱功能.mp4"
SLIDES_LIST_PATH = GENERATED_DIR / "slides_manifest.txt"

WIDTH = 1920
HEIGHT = 1080
BG_TOP = (241, 246, 255)
BG_BOTTOM = (226, 236, 255)
PRIMARY = (76, 110, 245)
PRIMARY_DARK = (48, 77, 201)
ACCENT = (18, 184, 134)
TEXT = (35, 45, 65)
MUTED = (92, 104, 128)
WHITE = (255, 255, 255)
BORDER = (205, 216, 238)
ORANGE = (255, 146, 43)
RED = (234, 67, 53)
PURPLE = (125, 92, 255)


@dataclass
class Slide:
    title: str
    subtitle: str
    bullets: List[str]
    footer: str
    diagram: str
    duration: int = 6


def ensure_dirs() -> None:
    for path in [DIAGRAMS_DIR, VIDEO_DIR, SLIDES_DIR, GENERATED_DIR, VIDEO_OUTPUT_DIR]:
        path.mkdir(parents=True, exist_ok=True)


def load_font(size: int, bold: bool = False):
    candidates = []
    if os.name == "nt":
        if bold:
            candidates += [
                r"C:\Windows\Fonts\msyhbd.ttc",
                r"C:\Windows\Fonts\simhei.ttf",
                r"C:\Windows\Fonts\arialbd.ttf",
            ]
        else:
            candidates += [
                r"C:\Windows\Fonts\msyh.ttc",
                r"C:\Windows\Fonts\simsun.ttc",
                r"C:\Windows\Fonts\arial.ttf",
            ]
    else:
        candidates += [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        ]

    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size=size)
        except Exception:
            continue
    return ImageFont.load_default()


FONT_H1 = load_font(54, bold=True)
FONT_H2 = load_font(30, bold=False)
FONT_BODY = load_font(28, bold=False)
FONT_BODY_BOLD = load_font(28, bold=True)
FONT_SMALL = load_font(22, bold=False)
FONT_TAG = load_font(20, bold=True)
FONT_BOX = load_font(24, bold=True)
FONT_BOX_SMALL = load_font(20, bold=False)


def draw_vertical_gradient(
    img: Image.Image, top: Tuple[int, int, int], bottom: Tuple[int, int, int]
) -> None:
    draw = ImageDraw.Draw(img)
    for y in range(img.height):
        ratio = y / max(1, img.height - 1)
        color = tuple(int(top[i] * (1 - ratio) + bottom[i] * ratio) for i in range(3))
        draw.line([(0, y), (img.width, y)], fill=color)


def rounded_box(draw: ImageDraw.ImageDraw, xy, fill, outline=None, radius=28, width=2):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def draw_text_block(
    draw: ImageDraw.ImageDraw,
    text: str,
    box: Tuple[int, int, int, int],
    font,
    fill,
    line_spacing: int = 10,
) -> int:
    x1, y1, x2, y2 = box
    max_width = x2 - x1
    lines: List[str] = []
    for paragraph in text.splitlines() or [""]:
        if not paragraph.strip():
            lines.append("")
            continue
        current = ""
        for ch in paragraph:
            test = current + ch
            bbox = draw.textbbox((0, 0), test, font=font)
            if bbox[2] - bbox[0] <= max_width or not current:
                current = test
            else:
                lines.append(current)
                current = ch
        if current:
            lines.append(current)
    y = y1
    for line in lines:
        draw.text((x1, y), line, font=font, fill=fill)
        bbox = draw.textbbox((x1, y), line or "中", font=font)
        y += (bbox[3] - bbox[1]) + line_spacing
        if y > y2:
            break
    return y


def center_text(draw: ImageDraw.ImageDraw, xy, text: str, font, fill):
    bbox = draw.textbbox((0, 0), text, font=font)
    x = xy[0] - (bbox[2] - bbox[0]) / 2
    y = xy[1] - (bbox[3] - bbox[1]) / 2
    draw.text((x, y), text, font=font, fill=fill)


def write_puml_files() -> None:
    arch = """@startuml
skinparam backgroundColor #F8FAFF
skinparam componentStyle rectangle
skinparam shadowing false
skinparam defaultFontName Microsoft YaHei
skinparam rectangle {
  BorderColor #4C6EF5
  RoundCorner 18
  BackgroundColor #FFFFFF
}
title 知识图谱核心功能架构图

actor 学生
actor 教师
actor 管理员

rectangle "Flutter UI" {
  rectangle "LoginPage\\n登录入口"
  rectangle "HomePage\\n导航中心"
  rectangle "GraphListPage\\n图谱列表"
  rectangle "GraphDetailPage\\n图谱详情与节点交互"
  rectangle "DocumentListPage\\nPPT/PDF资料"
  rectangle "VideoListPage\\n视频资源"
  rectangle "QuizPage\\n章节测验"
  rectangle "ProgressPage\\n学习进度"
  rectangle "FavoritesPage\\n收藏内容"
}

rectangle "Service / DAO" {
  rectangle "AuthService"
  rectangle "GraphDao"
  rectangle "QuizDao"
  rectangle "FavoriteDao"
  rectangle "LearningRecordDao"
  rectangle "UserDao"
  rectangle "DatabaseHelper"
}

database "SQLite / assets/learning_data.db" as DB

学生 --> LoginPage
教师 --> LoginPage
管理员 --> LoginPage

LoginPage --> AuthService
AuthService --> UserDao
UserDao --> DatabaseHelper
DatabaseHelper --> DB

HomePage --> GraphListPage
HomePage --> DocumentListPage
HomePage --> VideoListPage
HomePage --> QuizPage
HomePage --> ProgressPage
HomePage --> FavoritesPage

GraphListPage --> GraphDao
GraphDetailPage --> GraphDao
GraphDetailPage --> FavoriteDao
GraphDetailPage --> LearningRecordDao
QuizPage --> QuizDao
ProgressPage --> QuizDao
ProgressPage --> LearningRecordDao
FavoritesPage --> FavoriteDao

GraphDao --> DatabaseHelper
QuizDao --> DatabaseHelper
FavoriteDao --> DatabaseHelper
LearningRecordDao --> DatabaseHelper

@enduml
"""
    flow = """@startuml
skinparam backgroundColor #F8FAFF
skinparam defaultFontName Microsoft YaHei
skinparam shadowing false
title 图谱功能演示流程

start
:进入 HomePage;
:切换到底部导航“图谱”;
:GraphListPage 加载图谱列表;
if (是否存在图谱数据?) then (是)
  :点击某个图谱;
  :GraphDetailPage 查询 nodes / edges;
  :InteractiveViewer 缩放与拖拽;
  :点击节点;
  :显示节点详情;
  if (执行学习动作?) then (开始学习)
    :记录 learning_records;
  else (收藏)
    :写入 favorites;
  endif
  :返回图谱继续浏览;
else (否)
  :显示“暂无图谱数据”;
endif
stop
@enduml
"""
    PUML_ARCH_PATH.write_text(arch, encoding="utf-8")
    PUML_FLOW_PATH.write_text(flow, encoding="utf-8")


def render_architecture_png() -> None:
    img = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw_vertical_gradient(img, BG_TOP, BG_BOTTOM)
    draw = ImageDraw.Draw(img)

    draw.text((90, 54), "知识图谱核心功能架构图", font=FONT_H1, fill=TEXT)
    draw.text(
        (92, 126),
        "Flutter 页面层 · DAO / Service 层 · 本地 SQLite 数据层",
        font=FONT_H2,
        fill=MUTED,
    )

    # left actors
    actor_x = 130
    actor_ys = [300, 470, 640]
    actor_names = ["学生", "教师", "管理员"]
    for y, name in zip(actor_ys, actor_names):
        draw.ellipse(
            (actor_x, y, actor_x + 70, y + 70),
            outline=PRIMARY_DARK,
            width=5,
            fill=(255, 255, 255),
        )
        draw.line(
            (actor_x + 35, y + 70, actor_x + 35, y + 145), fill=PRIMARY_DARK, width=5
        )
        draw.line(
            (actor_x - 8, y + 96, actor_x + 78, y + 96), fill=PRIMARY_DARK, width=5
        )
        draw.line(
            (actor_x + 35, y + 145, actor_x - 5, y + 205), fill=PRIMARY_DARK, width=5
        )
        draw.line(
            (actor_x + 35, y + 145, actor_x + 75, y + 205), fill=PRIMARY_DARK, width=5
        )
        center_text(draw, (actor_x + 35, y + 235), name, FONT_BODY, TEXT)

    # UI layer
    rounded_box(
        draw,
        (300, 210, 1170, 860),
        fill=(255, 255, 255),
        outline=BORDER,
        radius=32,
        width=3,
    )
    draw.text((330, 228), "Flutter UI 页面层", font=FONT_H2, fill=PRIMARY_DARK)

    ui_boxes = [
        ("LoginPage", "登录入口", (340, 300, 580, 400), PRIMARY),
        ("HomePage", "导航中心", (620, 300, 860, 400), PRIMARY),
        ("GraphListPage", "图谱列表", (900, 300, 1140, 400), PRIMARY),
        ("GraphDetailPage", "图谱详情与节点交互", (340, 450, 620, 570), ACCENT),
        ("DocumentListPage", "PPT / PDF资料", (660, 450, 900, 570), ORANGE),
        ("VideoListPage", "视频资源", (940, 450, 1140, 570), RED),
        ("QuizPage", "章节测验", (340, 620, 580, 740), PURPLE),
        ("ProgressPage", "学习进度", (620, 620, 860, 740), ACCENT),
        ("FavoritesPage", "收藏内容", (900, 620, 1140, 740), ORANGE),
    ]

    for title, sub, rect, color in ui_boxes:
        rounded_box(draw, rect, fill=WHITE, outline=color, radius=22, width=4)
        center_text(
            draw, ((rect[0] + rect[2]) / 2, rect[1] + 32), title, FONT_BOX, color
        )
        draw_text_block(
            draw,
            sub,
            (rect[0] + 16, rect[1] + 58, rect[2] - 16, rect[3] - 16),
            FONT_BOX_SMALL,
            TEXT,
            6,
        )

    # service layer
    rounded_box(
        draw,
        (1210, 210, 1760, 760),
        fill=(255, 255, 255),
        outline=BORDER,
        radius=32,
        width=3,
    )
    draw.text((1240, 228), "Service / DAO 层", font=FONT_H2, fill=PRIMARY_DARK)

    dao_boxes = [
        ("AuthService", (1250, 300, 1510, 380), PRIMARY),
        ("GraphDao", (1530, 300, 1720, 380), ACCENT),
        ("QuizDao", (1250, 420, 1450, 500), PURPLE),
        ("FavoriteDao", (1470, 420, 1720, 500), ORANGE),
        ("LearningRecordDao", (1250, 540, 1520, 620), RED),
        ("UserDao", (1540, 540, 1720, 620), PRIMARY),
        ("DatabaseHelper", (1330, 660, 1640, 740), PRIMARY_DARK),
    ]
    for label, rect, color in dao_boxes:
        rounded_box(draw, rect, fill=WHITE, outline=color, radius=18, width=4)
        center_text(
            draw,
            ((rect[0] + rect[2]) / 2, (rect[1] + rect[3]) / 2),
            label,
            FONT_BOX,
            color,
        )

    # database
    rounded_box(
        draw,
        (1320, 840, 1650, 970),
        fill=(255, 255, 255),
        outline=PRIMARY_DARK,
        radius=32,
        width=4,
    )
    draw.ellipse(
        (1395, 825, 1575, 875), fill=(233, 239, 255), outline=PRIMARY_DARK, width=4
    )
    draw.rectangle(
        (1395, 850, 1575, 950), fill=(233, 239, 255), outline=PRIMARY_DARK, width=4
    )
    draw.ellipse(
        (1395, 925, 1575, 975), fill=(233, 239, 255), outline=PRIMARY_DARK, width=4
    )
    center_text(draw, (1485, 905), "SQLite / learning_data.db", FONT_BOX, PRIMARY_DARK)

    # connections
    def arrow(p1, p2, color=PRIMARY_DARK, width=4):
        draw.line((p1, p2), fill=color, width=width)
        angle = math.atan2(p2[1] - p1[1], p2[0] - p1[0])
        ah = 14
        left = (
            p2[0] - ah * math.cos(angle - math.pi / 6),
            p2[1] - ah * math.sin(angle - math.pi / 6),
        )
        right = (
            p2[0] - ah * math.cos(angle + math.pi / 6),
            p2[1] - ah * math.sin(angle + math.pi / 6),
        )
        draw.polygon([p2, left, right], fill=color)

    # actor to login
    arrow((235, 335), (340, 335))
    arrow((235, 505), (340, 345))
    arrow((235, 675), (340, 355))

    # home to pages
    arrow((740, 400), (740, 445), color=PRIMARY)
    arrow((860, 350), (900, 350), color=PRIMARY)
    arrow((740, 400), (1040, 445), color=PRIMARY)
    arrow((740, 400), (460, 620), color=PRIMARY)
    arrow((740, 400), (740, 620), color=PRIMARY)
    arrow((740, 400), (1020, 620), color=PRIMARY)

    # login to auth / dao
    arrow((580, 340), (1250, 340), color=PRIMARY_DARK)
    arrow((460, 570), (1620, 340), color=ACCENT)
    arrow((460, 570), (1595, 460), color=ORANGE)
    arrow((460, 570), (1385, 580), color=RED)
    arrow((460, 680), (1345, 460), color=PURPLE)
    arrow((740, 680), (1385, 580), color=RED)
    arrow((1020, 680), (1595, 460), color=ORANGE)

    # dao to db helper
    arrow((1380, 380), (1485, 660), color=PRIMARY_DARK)
    arrow((1625, 380), (1485, 660), color=PRIMARY_DARK)
    arrow((1350, 500), (1485, 660), color=PRIMARY_DARK)
    arrow((1595, 500), (1485, 660), color=PRIMARY_DARK)
    arrow((1385, 620), (1485, 660), color=PRIMARY_DARK)
    arrow((1630, 620), (1485, 660), color=PRIMARY_DARK)
    arrow((1485, 740), (1485, 840), color=PRIMARY_DARK)

    draw.text(
        (90, 1008),
        "说明：图谱模块是核心入口，节点交互会联动收藏、学习记录与进度统计。",
        font=FONT_SMALL,
        fill=MUTED,
    )

    ARCH_PNG_PATH.parent.mkdir(parents=True, exist_ok=True)
    img.save(ARCH_PNG_PATH)


def render_flow_png() -> None:
    img = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw_vertical_gradient(img, (248, 251, 255), (236, 244, 255))
    draw = ImageDraw.Draw(img)

    draw.text((90, 54), "图谱功能演示流程图", font=FONT_H1, fill=TEXT)
    draw.text(
        (92, 126), "从底部导航进入图谱，到节点学习与收藏联动", font=FONT_H2, fill=MUTED
    )

    boxes = [
        ("1", "HomePage\n进入系统首页", (160, 260, 430, 420), PRIMARY),
        ("2", "GraphListPage\n加载图谱列表", (500, 260, 770, 420), ACCENT),
        ("3", "GraphDetailPage\n查询 nodes / edges", (840, 260, 1170, 420), ORANGE),
        (
            "4",
            "InteractiveViewer\n缩放、拖拽、浏览结构",
            (1240, 260, 1610, 420),
            PURPLE,
        ),
        ("5", "点击节点\n展示详情卡片", (330, 590, 650, 750), RED),
        ("6", "开始学习\n写入 learning_records", (810, 590, 1160, 750), ACCENT),
        ("7", "收藏知识点\n写入 favorites", (1320, 590, 1660, 750), ORANGE),
    ]

    for tag, title, rect, color in boxes:
        rounded_box(draw, rect, fill=WHITE, outline=color, radius=24, width=5)
        rounded_box(
            draw,
            (rect[0] + 18, rect[1] + 18, rect[0] + 78, rect[1] + 78),
            fill=color,
            outline=color,
            radius=18,
            width=1,
        )
        center_text(draw, (rect[0] + 48, rect[1] + 48), tag, FONT_TAG, WHITE)
        draw_text_block(
            draw,
            title,
            (rect[0] + 96, rect[1] + 36, rect[2] - 24, rect[3] - 24),
            FONT_BOX,
            TEXT,
            8,
        )

    def arrow(p1, p2, color=PRIMARY_DARK, width=5):
        draw.line((p1, p2), fill=color, width=width)
        angle = math.atan2(p2[1] - p1[1], p2[0] - p1[0])
        ah = 16
        left = (
            p2[0] - ah * math.cos(angle - math.pi / 6),
            p2[1] - ah * math.sin(angle - math.pi / 6),
        )
        right = (
            p2[0] - ah * math.cos(angle + math.pi / 6),
            p2[1] - ah * math.sin(angle + math.pi / 6),
        )
        draw.polygon([p2, left, right], fill=color)

    arrow((430, 340), (500, 340), PRIMARY)
    arrow((770, 340), (840, 340), ACCENT)
    arrow((1170, 340), (1240, 340), ORANGE)
    arrow((1425, 420), (1425, 500), PURPLE)
    arrow((1425, 500), (490, 590), PRIMARY_DARK)
    arrow((650, 670), (810, 670), RED)
    arrow((650, 690), (1320, 690), ORANGE)

    draw.text((170, 840), "结果：", font=FONT_BODY_BOLD, fill=TEXT)
    draw.text(
        (260, 840),
        "用户可以从图谱总览快速进入知识点，查看节点内容，并把学习行为沉淀为记录与收藏。",
        font=FONT_BODY,
        fill=MUTED,
    )
    draw.text((170, 895), "价值：", font=FONT_BODY_BOLD, fill=TEXT)
    draw.text(
        (260, 895),
        "图谱不是静态展示，而是学习入口、行为采集入口和后续测验/进度分析的上游。",
        font=FONT_BODY,
        fill=MUTED,
    )

    FLOW_PNG_PATH.parent.mkdir(parents=True, exist_ok=True)
    img.save(FLOW_PNG_PATH)


def build_slides() -> List[Slide]:
    return [
        Slide(
            title="知识图谱核心功能演示",
            subtitle="第 1 支视频：图谱功能",
            bullets=[
                "本视频聚焦项目最核心的知识图谱能力，而不是登录注册。",
                "演示主线：图谱入口、图谱列表、节点交互、学习入口、收藏联动。",
                "项目已有 23 张图谱、537 个节点、616 条边，具备完整演示基础。",
            ],
            footer="目标：让观看者快速理解“图谱为什么是系统核心”。",
            diagram="arch",
            duration=6,
        ),
        Slide(
            title="项目定位与核心能力",
            subtitle="这是一个移动应用开发知识图谱学习系统",
            bullets=[
                "图谱：把课程知识结构化，建立章节、概念、能力之间的关系。",
                "资料：支持 PDF / PPT 课件访问，形成配套学习材料。",
                "视频：课程视频资源统一挂接，便于按章节学习。",
                "测验：章节题目与学习结果汇总，构成闭环。",
            ],
            footer="图谱负责组织知识，其他功能围绕图谱形成学习闭环。",
            diagram="arch",
            duration=6,
        ),
        Slide(
            title="图谱页面入口",
            subtitle="HomePage 作为统一导航中心",
            bullets=[
                "底部导航栏直接提供“图谱”入口，图谱是一级主功能。",
                "HomePage 同时串联测验、视频、资料、进度与计划。",
                "这意味着图谱并不是附属模块，而是核心学习视图之一。",
            ],
            footer="入口清晰，用户从首页即可快速进入图谱学习。",
            diagram="flow",
            duration=6,
        ),
        Slide(
            title="图谱浏览与节点交互",
            subtitle="GraphListPage + GraphDetailPage",
            bullets=[
                "图谱列表页从本地数据库读取全部图谱，支持刷新与进入详情。",
                "详情页读取 nodes / edges，并通过 InteractiveViewer 支持缩放与拖拽。",
                "点击节点后，下方弹出详情卡片，展示标题、类型、内容与操作按钮。",
            ],
            footer="用户既能看全局结构，也能深入单个知识点。",
            diagram="flow",
            duration=7,
        ),
        Slide(
            title="从图谱到学习行为沉淀",
            subtitle="学习记录与收藏联动",
            bullets=[
                "在节点详情中可执行“开始学习”，用于形成 learning_records。",
                "也可执行“收藏”，把重要知识点加入 favorites。",
                "这些数据后续会进入进度统计、学习建议和复习路径中。",
            ],
            footer="图谱是行为采集入口，而不是只读展示页。",
            diagram="arch",
            duration=6,
        ),
        Slide(
            title="系列视频规划",
            subtitle="本次先生成图谱视频，后续扩展 4 支",
            bullets=[
                "01 图谱：入口、列表、节点交互、学习与收藏。",
                "02 路径：学习记录、进度统计、学习建议。",
                "03 资料：PDF / PPT 文档管理与播放讲解。",
                "04 视频：视频资源组织与播放流程。",
                "05 测验：章节测验、错题本、成绩汇总。",
            ],
            footer="本文件已同时生成 UML、脚本、PPT 与第 1 支视频素材。",
            diagram="arch",
            duration=7,
        ),
    ]


def draw_header(draw: ImageDraw.ImageDraw, slide: Slide, index: int, total: int):
    draw.text((90, 60), slide.title, font=FONT_H1, fill=TEXT)
    draw.text((92, 138), slide.subtitle, font=FONT_H2, fill=MUTED)
    rounded_box(
        draw, (1620, 58, 1820, 118), fill=PRIMARY, outline=PRIMARY, radius=22, width=1
    )
    center_text(draw, (1720, 88), f"{index}/{total}", FONT_TAG, WHITE)


def draw_bullets(
    draw: ImageDraw.ImageDraw, bullets: List[str], box: Tuple[int, int, int, int]
) -> None:
    x1, y1, x2, y2 = box
    y = y1
    for bullet in bullets:
        rounded_box(
            draw,
            (x1, y + 8, x1 + 22, y + 30),
            fill=PRIMARY,
            outline=PRIMARY,
            radius=8,
            width=1,
        )
        y = (
            draw_text_block(draw, bullet, (x1 + 42, y, x2, y2), FONT_BODY, TEXT, 10)
            + 18
        )


def paste_diagram(
    base: Image.Image, diagram_path: Path, box: Tuple[int, int, int, int]
) -> None:
    if not diagram_path.exists():
        return
    diagram = Image.open(diagram_path).convert("RGB")
    target_w = box[2] - box[0]
    target_h = box[3] - box[1]
    ratio = min(target_w / diagram.width, target_h / diagram.height)
    new_size = (max(1, int(diagram.width * ratio)), max(1, int(diagram.height * ratio)))
    diagram = diagram.resize(new_size)
    x = box[0] + (target_w - new_size[0]) // 2
    y = box[1] + (target_h - new_size[1]) // 2
    base.paste(diagram, (x, y))


def render_slide(slide: Slide, index: int, total: int) -> Path:
    img = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw_vertical_gradient(img, BG_TOP, BG_BOTTOM)
    draw = ImageDraw.Draw(img)

    draw_header(draw, slide, index, total)

    rounded_box(
        draw, (80, 220, 900, 920), fill=WHITE, outline=BORDER, radius=32, width=3
    )
    rounded_box(
        draw, (940, 220, 1840, 920), fill=WHITE, outline=BORDER, radius=32, width=3
    )

    draw.text((120, 255), "讲解要点", font=FONT_H2, fill=PRIMARY_DARK)
    draw_bullets(draw, slide.bullets, (120, 320, 840, 820))

    draw.text((980, 255), "模型图", font=FONT_H2, fill=PRIMARY_DARK)
    diagram_path = ARCH_PNG_PATH if slide.diagram == "arch" else FLOW_PNG_PATH
    paste_diagram(img, diagram_path, (980, 315, 1800, 840))

    rounded_box(
        draw,
        (100, 945, 1820, 1015),
        fill=(238, 243, 255),
        outline=(218, 228, 250),
        radius=22,
        width=2,
    )
    draw.text((130, 964), slide.footer, font=FONT_SMALL, fill=MUTED)

    path = SLIDES_DIR / f"graph_feature_slide_{index:02d}.png"
    img.save(path)
    return path


def generate_slides(slides: List[Slide]) -> List[Path]:
    paths = []
    for idx, slide in enumerate(slides, start=1):
        paths.append(render_slide(slide, idx, len(slides)))
    return paths


def generate_markdown_script(slides: List[Slide]) -> None:
    lines = [
        "# 知识图谱核心功能视频脚本",
        "",
        "## 系列信息",
        "",
        "- 系列主题：移动应用开发知识图谱学习系统",
        "- 当前视频：01 图谱功能",
        "- 表现形式：PPT讲解 + 模型图 + 自动生成 MP4",
        "- 讲解重点：不讲登录注册，聚焦知识图谱核心能力",
        "",
        "## 五支视频规划",
        "",
        "1. 图谱功能视频",
        "2. 路径与学习记录视频",
        "3. 资料 / PPT 播放视频",
        "4. 视频资源播放视频",
        "5. 测验与错题本视频",
        "",
        "---",
        "",
        "## 01 图谱功能视频脚本",
        "",
    ]

    for idx, slide in enumerate(slides, start=1):
        lines += [
            f"### 第 {idx} 页：{slide.title}",
            "",
            f"**副标题**：{slide.subtitle}",
            "",
            "**讲解要点**：",
            "",
        ]
        for bullet in slide.bullets:
            lines.append(f"- {bullet}")
        lines += [
            "",
            f"**收束语**：{slide.footer}",
            "",
            f"**建议时长**：{slide.duration} 秒",
            "",
        ]

    lines += [
        "---",
        "",
        "## 可生成产物",
        "",
        "- UML 源文件：`docs/diagrams/knowledge_graph_feature_architecture.puml`",
        "- UML 源文件：`docs/diagrams/knowledge_graph_feature_flow.puml`",
        "- 模型图：`docs/video/generated/knowledge_graph_feature_architecture.png`",
        "- 模型图：`docs/video/generated/knowledge_graph_feature_flow.png`",
        "- PPT：`video_output/知识图谱核心功能_图谱功能演示.pptx`",
        "- 视频：`video_output/知识图谱核心功能_01_图谱功能.mp4`",
    ]

    SCRIPT_MD_PATH.write_text("\n".join(lines), encoding="utf-8")


def generate_pptx(slides: List[Slide], image_paths: List[Path]) -> None:
    if Presentation is None:
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data, image_path in zip(slides, image_paths):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(7.0), Inches(0.8)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = slide_data.title
        p.font.size = Pt(28)
        p.font.bold = True

        subtitle_box = slide.shapes.add_textbox(
            Inches(0.52), Inches(0.95), Inches(6.5), Inches(0.5)
        )
        p = subtitle_box.text_frame.paragraphs[0]
        p.text = slide_data.subtitle
        p.font.size = Pt(16)

        bullet_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.6), Inches(4.8), Inches(4.8)
        )
        tf = bullet_box.text_frame
        tf.clear()
        for i, bullet in enumerate(slide_data.bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = bullet
            para.level = 0
            para.font.size = Pt(18)

        slide.shapes.add_picture(
            str(image_path),
            Inches(6.0),
            Inches(1.55),
            width=Inches(6.7),
            height=Inches(4.75),
        )

        footer_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(6.65), Inches(11.8), Inches(0.35)
        )
        p = footer_box.text_frame.paragraphs[0]
        p.text = slide_data.footer
        p.font.size = Pt(12)

    if len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    prs.save(PPTX_PATH)


def build_ffmpeg_concat_file(image_paths: List[Path], slides: List[Slide]) -> None:
    lines = []
    for path, slide in zip(image_paths, slides):
        unix_path = path.as_posix().replace("'", "'\\''")
        lines.append(f"file '{unix_path}'")
        lines.append(f"duration {slide.duration}")
    if image_paths:
        unix_path = image_paths[-1].as_posix().replace("'", "'\\''")
        lines.append(f"file '{unix_path}'")
    SLIDES_LIST_PATH.write_text("\n".join(lines), encoding="utf-8")


def which_ffmpeg() -> str | None:
    candidates = [
        "ffmpeg",
        r"C:\Users\ldl\AppData\Local\Microsoft\WinGet\Links\ffmpeg.exe",
    ]
    for candidate in candidates:
        try:
            result = subprocess.run(
                [candidate, "-version"],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                check=False,
            )
            if result.returncode == 0:
                return candidate
        except Exception:
            continue
    return None


def generate_video(image_paths: List[Path], slides: List[Slide]) -> bool:
    ffmpeg = which_ffmpeg()
    if not ffmpeg:
        return False

    build_ffmpeg_concat_file(image_paths, slides)

    cmd = [
        ffmpeg,
        "-y",
        "-f",
        "concat",
        "-safe",
        "0",
        "-i",
        str(SLIDES_LIST_PATH),
        "-vf",
        "fps=30,format=yuv420p",
        "-pix_fmt",
        "yuv420p",
        str(VIDEO_PATH),
    ]
    result = subprocess.run(
        cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, check=False
    )
    return result.returncode == 0


def print_summary(image_paths: List[Path], video_ok: bool) -> None:
    print("已生成以下文件：")
    print(f"- {PUML_ARCH_PATH}")
    print(f"- {PUML_FLOW_PATH}")
    print(f"- {ARCH_PNG_PATH}")
    print(f"- {FLOW_PNG_PATH}")
    print(f"- {SCRIPT_MD_PATH}")
    for img in image_paths:
        print(f"- {img}")
    if PPTX_PATH.exists():
        print(f"- {PPTX_PATH}")
    if video_ok and VIDEO_PATH.exists():
        print(f"- {VIDEO_PATH}")
    else:
        print("- 视频未成功生成，请检查 ffmpeg 是否可用。")


def main() -> None:
    ensure_dirs()
    write_puml_files()
    render_architecture_png()
    render_flow_png()

    slides = build_slides()
    image_paths = generate_slides(slides)
    generate_markdown_script(slides)
    generate_pptx(slides, image_paths)
    video_ok = generate_video(image_paths, slides)

    print_summary(image_paths, video_ok)


if __name__ == "__main__":
    main()
