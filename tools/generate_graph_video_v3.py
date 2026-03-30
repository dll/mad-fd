from __future__ import annotations

import math
import os
import shutil
import struct
import subprocess
import wave
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional, Tuple

from PIL import Image, ImageDraw, ImageFont

try:
    import pyttsx3
except Exception:
    pyttsx3 = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
DIAGRAMS_V3_DIR = DOCS_DIR / "diagrams" / "v3"
VIDEO_V3_DIR = DOCS_DIR / "video" / "v3"
TESTING_DIR = DOCS_DIR / "testing"
VIDEO_OUTPUT_DIR = ROOT / "video_output"

SLIDES_DIR = VIDEO_V3_DIR / "slides"
AUDIO_DIR = VIDEO_V3_DIR / "audio"
CLIPS_DIR = VIDEO_V3_DIR / "clips"
TEMP_DIR = VIDEO_V3_DIR / "temp"

FRAMEWORK_IMG = DIAGRAMS_V3_DIR / "flutter_dart_framework_architecture.png"
CLASS_IMG = DIAGRAMS_V3_DIR / "flutter_dart_core_class_diagram.png"
SEQUENCE_IMG = DIAGRAMS_V3_DIR / "graph_feature_sequence_diagram.png"
PROCESS_IMG = DIAGRAMS_V3_DIR / "knowledge_graph_development_process.png"

SCRIPT_PATH = VIDEO_V3_DIR / "video_script_v3_compact.md"
PPTX_PATH = VIDEO_OUTPUT_DIR / "知识图谱核心功能_图谱功能教程_v3.pptx"
VIDEO_PATH = VIDEO_OUTPUT_DIR / "知识图谱核心功能_图谱功能教程_v3.mp4"

MODEL_TEST_PATH = ROOT / "test" / "models" / "model_test.dart"
HOME_WIDGET_TEST_PATH = ROOT / "test" / "widgets" / "home_page_widget_test.dart"
LOGIN_WIDGET_TEST_PATH = ROOT / "test" / "widget_test.dart"
TEST_CASES_PATH = TESTING_DIR / "test_cases.md"
TEST_REPORT_PATH = TESTING_DIR / "test_report.md"

WIDTH = 1920
HEIGHT = 1080

BG_TOP = (244, 248, 255)
BG_BOTTOM = (230, 238, 252)
WHITE = (255, 255, 255)
TEXT = (34, 44, 66)
TEXT_LIGHT = (92, 104, 128)
PRIMARY = (76, 110, 245)
PRIMARY_DARK = (48, 77, 201)
ACCENT = (18, 184, 134)
ORANGE = (255, 146, 43)
RED = (235, 87, 87)
PURPLE = (125, 92, 255)
BORDER = (210, 220, 240)
PANEL_BG = (250, 252, 255)
TAG_BG = (233, 239, 255)
CODE_BG = (245, 247, 251)

TITLE_FONT_SIZE = 56
SUBTITLE_FONT_SIZE = 28
SECTION_FONT_SIZE = 28
BODY_FONT_SIZE = 26
SMALL_FONT_SIZE = 20
CODE_FONT_SIZE = 21


@dataclass
class SlideSpec:
    title: str
    subtitle: str
    bullets: List[str]
    narration: str
    image_path: Optional[Path] = None
    image_caption: Optional[str] = None
    code_title: Optional[str] = None
    code_text: Optional[str] = None


def ensure_dirs() -> None:
    for path in [
        VIDEO_V3_DIR,
        SLIDES_DIR,
        AUDIO_DIR,
        CLIPS_DIR,
        TEMP_DIR,
        VIDEO_OUTPUT_DIR,
    ]:
        path.mkdir(parents=True, exist_ok=True)


def load_font(size: int, bold: bool = False):
    candidates = []
    if os.name == "nt":
        if bold:
            candidates += [
                r"C:\Windows\Fonts\msyhbd.ttc",
                r"C:\Windows\Fonts\simhei.ttf",
                r"C:\Windows\Fonts\arialbd.ttf",
                r"C:\Windows\Fonts\calibrib.ttf",
            ]
        else:
            candidates += [
                r"C:\Windows\Fonts\msyh.ttc",
                r"C:\Windows\Fonts\simsun.ttc",
                r"C:\Windows\Fonts\arial.ttf",
                r"C:\Windows\Fonts\calibri.ttf",
            ]
    else:
        candidates += [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        ]

    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size=size)
        except Exception:
            continue
    return ImageFont.load_default()


FONT_TITLE = load_font(TITLE_FONT_SIZE, True)
FONT_SUBTITLE = load_font(SUBTITLE_FONT_SIZE, False)
FONT_SECTION = load_font(SECTION_FONT_SIZE, True)
FONT_BODY = load_font(BODY_FONT_SIZE, False)
FONT_BODY_BOLD = load_font(BODY_FONT_SIZE, True)
FONT_SMALL = load_font(SMALL_FONT_SIZE, False)
FONT_SMALL_BOLD = load_font(SMALL_FONT_SIZE, True)
FONT_CODE = load_font(CODE_FONT_SIZE, False)


def draw_vertical_gradient(
    img: Image.Image, top: Tuple[int, int, int], bottom: Tuple[int, int, int]
) -> None:
    draw = ImageDraw.Draw(img)
    for y in range(img.height):
        ratio = y / max(1, img.height - 1)
        color = tuple(int(top[i] * (1 - ratio) + bottom[i] * ratio) for i in range(3))
        draw.line([(0, y), (img.width, y)], fill=color)


def rounded_rect(draw: ImageDraw.ImageDraw, xy, fill, outline=None, radius=24, width=2):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def text_size(draw: ImageDraw.ImageDraw, text: str, font) -> Tuple[int, int]:
    bbox = draw.textbbox((0, 0), text, font=font)
    return bbox[2] - bbox[0], bbox[3] - bbox[1]


def center_text(
    draw: ImageDraw.ImageDraw, x: int, y: int, text: str, font, fill
) -> None:
    w, h = text_size(draw, text, font)
    draw.text((x - w / 2, y - h / 2), text, font=font, fill=fill)


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font, max_width: int) -> List[str]:
    lines: List[str] = []
    for paragraph in text.splitlines():
        p = paragraph.strip()
        if not p:
            lines.append("")
            continue
        current = ""
        for ch in p:
            candidate = current + ch
            w, _ = text_size(draw, candidate, font)
            if w <= max_width or not current:
                current = candidate
            else:
                lines.append(current)
                current = ch
        if current:
            lines.append(current)
    return lines


def draw_wrapped_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    box: Tuple[int, int, int, int],
    font,
    fill,
    line_spacing: int = 8,
) -> int:
    x1, y1, x2, y2 = box
    y = y1
    max_width = x2 - x1
    for line in wrap_text(draw, text, font, max_width):
        if y > y2:
            break
        draw.text((x1, y), line, font=font, fill=fill)
        _, h = text_size(draw, line or "中", font)
        y += h + line_spacing
    return y


def safe_name(text: str) -> str:
    bad = '\\/:*?"<>| '
    result = "".join("_" if ch in bad else ch for ch in text)
    while "__" in result:
        result = result.replace("__", "_")
    return result.strip("_") or "slide"


def read_text(path: Path, fallback: str = "") -> str:
    if not path.exists():
        return fallback
    return path.read_text(encoding="utf-8")


def extract_block_by_marker(path: Path, marker: str, max_lines: int = 26) -> str:
    if not path.exists():
        return f"// missing file: {path.name}"
    lines = path.read_text(encoding="utf-8").splitlines()
    start = 0
    for i, line in enumerate(lines):
        if marker in line:
            start = i
            break
    end = min(len(lines), start + max_lines)
    return "\n".join(lines[start:end]).strip()


def select_lines(
    text: str, include_keywords: Iterable[str], max_lines: int = 18
) -> str:
    lines = [line.rstrip() for line in text.splitlines() if line.strip()]
    picked: List[str] = []
    lowered = [k.lower() for k in include_keywords]
    for line in lines:
        content = line.lower()
        if any(k in content for k in lowered):
            picked.append(line)
        if len(picked) >= max_lines:
            break
    if not picked:
        picked = lines[:max_lines]
    return "\n".join(picked)


def load_test_summary() -> str:
    report = read_text(TEST_REPORT_PATH)
    return select_lines(
        report,
        include_keywords=[
            "通过",
            "测试",
            "apk",
            "graphmodel",
            "nodemodel",
            "edgemodel",
            "questionmodel",
            "quizresultmodel",
            "usermodel",
            "homepage",
        ],
        max_lines=16,
    )


def load_test_case_summary() -> str:
    cases = read_text(TEST_CASES_PATH)
    return select_lines(
        cases,
        include_keywords=[
            "tc-",
            "图谱",
            "登录",
            "测验",
            "资料",
            "视频",
            "progress",
            "model",
        ],
        max_lines=16,
    )


def load_code_summaries() -> Tuple[str, str]:
    model_code = extract_block_by_marker(MODEL_TEST_PATH, "group('GraphModel'", 28)
    widget_code = extract_block_by_marker(
        HOME_WIDGET_TEST_PATH, "group('HomePage widget tests'", 28
    )
    return model_code, widget_code


def build_script_markdown(slides: List[SlideSpec]) -> None:
    lines = [
        "# 知识图谱核心功能教程脚本 v3（成片展示版）",
        "",
        "## 说明",
        "",
        "- 本文件用于 v3 视频生成。",
        "- 视频画面不展示“旁白 / 时长 / 画面说明”等脚本字段。",
        "- 成片只展示标题、技术要点、UML 图、测试与开发总结。",
        "",
        "## 成片结构",
        "",
    ]
    for i, slide in enumerate(slides, start=1):
        lines.append(f"{i}. {slide.title} - {slide.subtitle}")
    lines += ["", "## 讲解稿", ""]
    for i, slide in enumerate(slides, start=1):
        lines += [
            f"### {i}. {slide.title}",
            "",
            f"- 副标题：{slide.subtitle}",
            f"- 关键词：{', '.join(slide.bullets)}",
            f"- 讲解：{slide.narration}",
            "",
        ]
    SCRIPT_PATH.write_text("\n".join(lines), encoding="utf-8")


def build_slides() -> List[SlideSpec]:
    model_test_code, widget_test_code = load_code_summaries()
    test_case_summary = load_test_case_summary()
    test_report_summary = load_test_summary()

    return [
        SlideSpec(
            title="Flutter + Dart 知识图谱学习系统",
            subtitle="核心功能教程 v3",
            bullets=[
                "Flutter UI",
                "Dart Logic",
                "SQLite",
                "UML 设计",
                "知识图谱",
            ],
            narration=(
                "本视频聚焦移动应用开发知识图谱学习系统的核心设计。"
                "重点说明 Flutter 与 Dart 的技术分工，解释系统框架图、类图、顺序图，"
                "并进一步说明知识图谱功能的开发过程，以及测试代码、测试用例和测试报告等工程化成果。"
            ),
            image_path=FRAMEWORK_IMG if FRAMEWORK_IMG.exists() else None,
            image_caption="项目总览：Flutter + Dart + SQLite + UML",
        ),
        SlideSpec(
            title="技术栈与分层定位",
            subtitle="Flutter 负责界面，Dart 负责逻辑，SQLite 负责持久化",
            bullets=[
                "Flutter 页面构建",
                "Dart 数据模型",
                "DAO 封装",
                "Material Design",
                "本地数据库驱动",
            ],
            narration=(
                "这个项目的技术栈非常清晰。Flutter 负责页面构建、导航与状态渲染，"
                "Dart 负责模型类、业务逻辑和数据库访问封装，底层由 SQLite 保存图谱、节点、边、题目、学习记录和收藏数据。"
                "这种分层结构使界面层、逻辑层和数据层边界明确，便于维护和扩展。"
            ),
            image_path=FRAMEWORK_IMG if FRAMEWORK_IMG.exists() else None,
            image_caption="技术栈：Flutter + Dart + SQLite",
            code_title="技术关键词",
            code_text=(
                "Flutter\n"
                "Dart\n"
                "MaterialApp / Widget Tree\n"
                "DAO Pattern\n"
                "SQLite / assets/learning_data.db\n"
                "GraphModel / NodeModel / EdgeModel"
            ),
        ),
        SlideSpec(
            title="系统框架图",
            subtitle="从 Flutter UI 到 DAO 再到本地数据库",
            bullets=[
                "UI Layer",
                "Service / DAO Layer",
                "Data Model Layer",
                "SQLite",
                "学习闭环",
            ],
            narration=(
                "这张框架图展示了系统整体架构。顶部和左侧是 Flutter 页面层，包括登录、首页、图谱、视频、资料、测验、进度与收藏等页面。"
                "中间是 Service 与 DAO 层，负责把页面请求转换为数据库读写。"
                "底部是 SQLite 本地数据库，用于保存图谱结构和学习行为。"
                "因此，知识图谱不仅是一个页面，而是整个学习系统的组织核心。"
            ),
            image_path=FRAMEWORK_IMG if FRAMEWORK_IMG.exists() else None,
            image_caption="框架图解释：UI、业务逻辑、数据存储三层分工",
        ),
        SlideSpec(
            title="框架图重点解释",
            subtitle="为什么图谱模块是系统核心",
            bullets=[
                "HomePage 导航中心",
                "GraphListPage 入口",
                "GraphDetailPage 核心交互",
                "资料 / 视频 / 测验联动",
                "进度反馈",
            ],
            narration=(
                "从框架图可以看到，HomePage 是统一导航中心，而 GraphListPage 和 GraphDetailPage 组成了知识图谱主流程。"
                "资料页、视频页、测验页和进度页并不是独立模块，而是围绕图谱形成学习闭环。"
                "也就是说，图谱负责组织课程知识结构，其他模块负责提供资源、验证效果并反馈学习结果。"
            ),
            image_path=FRAMEWORK_IMG if FRAMEWORK_IMG.exists() else None,
            image_caption="图谱处于学习闭环中心",
            code_title="框架图说明关键词",
            code_text=(
                "HomePage -> GraphListPage -> GraphDetailPage\n"
                "GraphDetailPage -> FavoriteDao\n"
                "GraphDetailPage -> LearningRecordDao\n"
                "ProgressPage -> QuizDao / LearningRecordDao\n"
                "DatabaseHelper -> SQLite"
            ),
        ),
        SlideSpec(
            title="核心类图",
            subtitle="Flutter 页面类、DAO 类和模型类的职责划分",
            bullets=[
                "UI Layer",
                "GraphDetailPage",
                "GraphDao",
                "DatabaseHelper",
                "Model Objects",
            ],
            narration=(
                "类图展示了项目中最关键的类以及它们之间的关系。"
                "页面层负责显示与交互，DAO 层负责查询与写入，模型层负责承接结构化数据。"
                "其中 GraphDetailPage 是图谱功能的核心页面，GraphDao 负责图谱、节点和边的读取，"
                "DatabaseHelper 统一管理 SQLite 初始化和表结构。"
            ),
            image_path=CLASS_IMG if CLASS_IMG.exists() else None,
            image_caption="类图：页面层、DAO 层、模型层职责清晰",
        ),
        SlideSpec(
            title="类图中的关键类解释",
            subtitle="围绕图谱功能的核心类协作",
            bullets=[
                "GraphListPage",
                "GraphDetailPage",
                "GraphPainter",
                "FavoriteDao",
                "LearningRecordDao",
            ],
            narration=(
                "在类图中，GraphListPage 负责图谱入口，GraphDetailPage 负责节点与边加载、布局计算和节点点击处理。"
                "GraphPainter 负责把图谱可视化绘制出来。"
                "当用户点击节点后，FavoriteDao 可以负责收藏写入，LearningRecordDao 可以负责学习记录写入。"
                "这种设计使图谱从纯展示层升级为行为驱动入口。"
            ),
            image_path=CLASS_IMG if CLASS_IMG.exists() else None,
            image_caption="图谱相关关键类：加载、绘制、交互、写入",
            code_title="核心测试代码片段",
            code_text=model_test_code,
        ),
        SlideSpec(
            title="图谱功能顺序图",
            subtitle="从进入图谱到节点交互的完整执行过程",
            bullets=[
                "HomePage",
                "GraphListPage",
                "GraphDao",
                "GraphDetailPage",
                "SQLite",
            ],
            narration=(
                "顺序图比类图更强调运行过程。"
                "用户首先从 HomePage 切换到底部导航中的图谱模块，GraphListPage 调用 GraphDao 查询所有图谱。"
                "当用户点击某个图谱后，GraphDetailPage 继续查询当前图谱对应的 nodes 和 edges，"
                "随后再把这些数据组织为可浏览、可缩放、可点击的图谱结构。"
            ),
            image_path=SEQUENCE_IMG if SEQUENCE_IMG.exists() else None,
            image_caption="顺序图：页面、DAO 与数据库的动态调用过程",
        ),
        SlideSpec(
            title="顺序图重点解释",
            subtitle="节点点击、开始学习与收藏是图谱价值的关键",
            bullets=[
                "节点点击",
                "详情卡片",
                "开始学习",
                "收藏",
                "学习行为沉淀",
            ],
            narration=(
                "顺序图最关键的部分是节点交互。"
                "用户点击节点后，GraphDetailPage 会定位选中节点并展示详情卡片。"
                "如果选择开始学习，就会调用 LearningRecordDao，把行为写入 learning_records。"
                "如果选择收藏，就会调用 FavoriteDao，把节点写入 favorites。"
                "这说明知识图谱不是静态图，而是学习行为采集入口。"
            ),
            image_path=SEQUENCE_IMG if SEQUENCE_IMG.exists() else None,
            image_caption="节点交互：从浏览转化为学习记录与收藏",
            code_title="HomePage 组件测试片段",
            code_text=widget_test_code,
        ),
        SlideSpec(
            title="知识图谱开发过程",
            subtitle="从需求分析到图谱交互，再到学习闭环",
            bullets=[
                "需求分析",
                "数据建模",
                "DAO 开发",
                "页面交互",
                "测试交付",
            ],
            narration=(
                "开发过程图说明了知识图谱功能是如何一步步实现出来的。"
                "首先要明确图谱是系统核心，然后定义 Graph、Node、Edge、Favorite、LearningRecord 等业务对象，"
                "再设计 SQLite 表结构，接着完成 DatabaseHelper 和 DAO 层，最后再实现 Flutter 页面和节点交互逻辑。"
                "完成图谱功能之后，继续扩展资料、视频、测验和进度模块，从而形成完整学习闭环。"
            ),
            image_path=PROCESS_IMG if PROCESS_IMG.exists() else None,
            image_caption="开发过程图：需求、建模、实现、测试、交付",
        ),
        SlideSpec(
            title="测试代码与测试用例",
            subtitle="模型测试、组件测试和构建验证共同保证质量",
            bullets=[
                "model_test.dart",
                "widget_test.dart",
                "home_page_widget_test.dart",
                "test_cases.md",
                "工程化验证",
            ],
            narration=(
                "在工程实现之外，项目还补充了测试代码与测试文档。"
                "模型测试覆盖 GraphModel、NodeModel、EdgeModel、QuestionModel、QuizResultModel 和 UserModel 的映射与计算逻辑。"
                "组件测试覆盖登录页和首页导航。"
                "同时，测试用例文档对图谱、测验、资料、视频和进度等模块给出了系统化验证点。"
            ),
            image_path=PROCESS_IMG if PROCESS_IMG.exists() else None,
            image_caption="测试资产：代码 + 用例 + 构建验证",
            code_title="测试用例摘要",
            code_text=test_case_summary,
        ),
        SlideSpec(
            title="测试报告与交付结果",
            subtitle="测试通过、APK 构建成功、教学展示资料齐备",
            bullets=[
                "23 项测试通过",
                "Debug APK",
                "Release APK",
                "测试报告",
                "视频与 PPT",
            ],
            narration=(
                "当前项目已经补充测试报告，并完成了测试执行。"
                "模型层测试与首页组件测试已经通过，同时 Android 的 Debug 和 Release APK 也已经能够成功构建。"
                "在此基础上，项目还输出了 UML 图、脚本、PPT 和视频，因此它不仅具备功能实现能力，也具备了教学展示与答辩汇报能力。"
            ),
            image_path=CLASS_IMG if CLASS_IMG.exists() else None,
            image_caption="测试与交付：代码质量与展示材料同时具备",
            code_title="测试报告摘要",
            code_text=test_report_summary,
        ),
        SlideSpec(
            title="总结",
            subtitle="Flutter + Dart + UML + 测试资产构成完整答辩材料",
            bullets=[
                "技术栈清晰",
                "架构分层明确",
                "图谱模块核心",
                "开发过程完整",
                "测试与交付齐全",
            ],
            narration=(
                "总结来看，这个项目的价值不只是把知识图谱画出来，而是通过 Flutter 和 Dart 建立了清晰的技术分层，"
                "通过 UML 图把设计说明清楚，通过测试代码、测试用例和测试报告把工程质量说明清楚，"
                "最终形成了适合课程汇报、项目答辩和教学展示的完整材料。"
            ),
            image_path=FRAMEWORK_IMG if FRAMEWORK_IMG.exists() else None,
            image_caption="总结：设计、实现、测试、展示四位一体",
        ),
    ]


def paste_image(
    base: Image.Image, source: Path, box: Tuple[int, int, int, int]
) -> None:
    if not source or not source.exists():
        return
    img = Image.open(source).convert("RGB")
    target_w = box[2] - box[0]
    target_h = box[3] - box[1]
    ratio = min(target_w / img.width, target_h / img.height)
    new_size = (max(1, int(img.width * ratio)), max(1, int(img.height * ratio)))
    img = img.resize(new_size)
    x = box[0] + (target_w - new_size[0]) // 2
    y = box[1] + (target_h - new_size[1]) // 2
    base.paste(img, (x, y))


def draw_bullets(
    draw: ImageDraw.ImageDraw, bullets: List[str], box: Tuple[int, int, int, int]
):
    x1, y1, x2, y2 = box
    y = y1
    for bullet in bullets:
        rounded_rect(draw, (x1, y + 8, x1 + 18, y + 26), PRIMARY, PRIMARY, 6, 1)
        y = (
            draw_wrapped_text(
                draw, bullet, (x1 + 32, y, x2, y + 60), FONT_BODY, TEXT, 8
            )
            + 10
        )
        if y > y2:
            break


def draw_code_block(
    draw: ImageDraw.ImageDraw,
    title: str,
    code: str,
    box: Tuple[int, int, int, int],
) -> None:
    rounded_rect(draw, box, CODE_BG, (224, 230, 242), 18, 2)
    x1, y1, x2, y2 = box
    draw.text((x1 + 18, y1 + 14), title, font=FONT_SMALL_BOLD, fill=PRIMARY_DARK)
    y = y1 + 54
    for line in code.splitlines():
        if y > y2 - 24:
            break
        draw.text((x1 + 18, y), line.rstrip(), font=FONT_CODE, fill=TEXT)
        _, h = text_size(draw, line or "中", FONT_CODE)
        y += h + 5


def render_slide(slide: SlideSpec, index: int, total: int) -> Path:
    img = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw_vertical_gradient(img, BG_TOP, BG_BOTTOM)
    draw = ImageDraw.Draw(img)

    rounded_rect(draw, (48, 34, 1870, 150), WHITE, BORDER, 28, 2)
    draw.text((84, 50), slide.title, font=FONT_TITLE, fill=TEXT)
    draw.text((88, 108), slide.subtitle, font=FONT_SUBTITLE, fill=TEXT_LIGHT)

    rounded_rect(draw, (1660, 54, 1835, 108), PRIMARY, PRIMARY, 18, 1)
    center_text(draw, 1748, 82, f"{index}/{total}", FONT_SMALL_BOLD, WHITE)

    left_box = (56, 180, 690, 1005)
    right_box = (720, 180, 1862, 735)
    bottom_box = (720, 760, 1862, 1005)

    rounded_rect(draw, left_box, WHITE, BORDER, 30, 3)
    rounded_rect(draw, right_box, WHITE, BORDER, 30, 3)
    rounded_rect(draw, bottom_box, WHITE, BORDER, 30, 3)

    draw.text((90, 210), "核心要点", font=FONT_SECTION, fill=PRIMARY_DARK)
    draw_bullets(draw, slide.bullets, (90, 280, 630, 580))

    rounded_rect(draw, (88, 615, 286, 660), TAG_BG, TAG_BG, 16, 1)
    draw.text((110, 627), "讲解关键词", font=FONT_SMALL_BOLD, fill=PRIMARY_DARK)

    keywords_text = " · ".join(slide.bullets)
    draw_wrapped_text(draw, keywords_text, (90, 685, 630, 770), FONT_BODY, TEXT)

    rounded_rect(draw, (88, 810, 286, 855), TAG_BG, TAG_BG, 16, 1)
    draw.text((110, 822), "图示说明", font=FONT_SMALL_BOLD, fill=PRIMARY_DARK)

    if slide.image_caption:
        draw_wrapped_text(
            draw, slide.image_caption, (90, 880, 630, 970), FONT_BODY, TEXT
        )

    draw.text((752, 210), "UML / 结构图", font=FONT_SECTION, fill=PRIMARY_DARK)
    paste_image(img, slide.image_path, (760, 265, 1818, 685))

    if slide.code_title and slide.code_text:
        draw_code_block(draw, slide.code_title, slide.code_text, (752, 810, 1818, 970))
    else:
        draw_code_block(
            draw,
            "本页说明",
            "本页重点展示结构图与核心关键词，不展示脚本原文字段。",
            (752, 810, 1818, 970),
        )

    out_path = SLIDES_DIR / f"{index:02d}_{safe_name(slide.title)}.png"
    img.save(out_path)
    return out_path


def render_all_slides(slides: List[SlideSpec]) -> List[Path]:
    return [
        render_slide(slide, i, len(slides)) for i, slide in enumerate(slides, start=1)
    ]


def choose_voice(engine, preferred_keywords: Iterable[str]) -> Optional[str]:
    voices = engine.getProperty("voices") or []
    preferred = [k.lower() for k in preferred_keywords]
    for voice in voices:
        text = f"{getattr(voice, 'id', '')} {getattr(voice, 'name', '')}".lower()
        if any(key in text for key in preferred):
            return voice.id
    return voices[0].id if voices else None


def clean_tts_text(text: str) -> str:
    replacements = {
        "Flutter": "Flutter",
        "Dart": "Dart",
        "SQLite": "S Q Lite",
        "GraphDao": "Graph Dao",
        "FavoriteDao": "Favorite Dao",
        "LearningRecordDao": "Learning Record Dao",
        "DatabaseHelper": "Database Helper",
        "GraphListPage": "Graph List Page",
        "GraphDetailPage": "Graph Detail Page",
        "HomePage": "Home Page",
        "QuizPage": "Quiz Page",
        "ProgressPage": "Progress Page",
        "favorites": "favorites",
        "learning_records": "learning records",
        "UML": "U M L",
        "APK": "A P K",
    }
    result = text
    for k, v in replacements.items():
        result = result.replace(k, v)
    return result


def save_tts_audio(text: str, out_path: Path, rate: int = 162) -> bool:
    if pyttsx3 is None:
        return False
    try:
        engine = pyttsx3.init()
        voice_id = choose_voice(
            engine, ["zh", "chinese", "huihui", "hanhan", "xiaoxiao"]
        )
        if voice_id:
            engine.setProperty("voice", voice_id)
        engine.setProperty("rate", rate)
        engine.save_to_file(clean_tts_text(text), str(out_path))
        engine.runAndWait()
        engine.stop()
        return out_path.exists() and out_path.stat().st_size > 0
    except Exception:
        return False


def create_silent_wav(duration_seconds: float, out_path: Path) -> None:
    sample_rate = 22050
    channels = 1
    sample_width = 2
    frames = max(1, int(duration_seconds * sample_rate))
    silence = struct.pack("<h", 0)
    with wave.open(str(out_path), "wb") as f:
        f.setnchannels(channels)
        f.setsampwidth(sample_width)
        f.setframerate(sample_rate)
        for _ in range(frames):
            f.writeframesraw(silence)


def wav_duration(path: Path) -> float:
    with wave.open(str(path), "rb") as f:
        return f.getnframes() / float(f.getframerate())


def which(names: List[str]) -> Optional[str]:
    for name in names:
        found = shutil.which(name)
        if found:
            return found
    return None


def ffmpeg_path() -> Optional[str]:
    return which(
        ["ffmpeg", r"C:\Users\ldl\AppData\Local\Microsoft\WinGet\Links\ffmpeg.exe"]
    )


def run_ffmpeg(args: List[str]) -> bool:
    ffmpeg = ffmpeg_path()
    if not ffmpeg:
        return False
    process = subprocess.run(
        [ffmpeg] + args,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
    )
    return process.returncode == 0


def build_clip(image_path: Path, audio_path: Path, output_path: Path) -> bool:
    return run_ffmpeg(
        [
            "-y",
            "-loop",
            "1",
            "-i",
            str(image_path),
            "-i",
            str(audio_path),
            "-c:v",
            "libx264",
            "-tune",
            "stillimage",
            "-c:a",
            "aac",
            "-b:a",
            "192k",
            "-pix_fmt",
            "yuv420p",
            "-shortest",
            "-vf",
            "scale=1920:1080,fps=30",
            str(output_path),
        ]
    )


def concat_clips(clips: List[Path], output_path: Path) -> bool:
    ffmpeg = ffmpeg_path()
    if not ffmpeg:
        return False
    manifest = TEMP_DIR / "clips_manifest.txt"
    manifest.write_text(
        "\n".join(f"file '{clip.as_posix()}'" for clip in clips), encoding="utf-8"
    )
    process = subprocess.run(
        [
            ffmpeg,
            "-y",
            "-f",
            "concat",
            "-safe",
            "0",
            "-i",
            str(manifest),
            "-c",
            "copy",
            str(output_path),
        ],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
    )
    return process.returncode == 0 and output_path.exists()


def narration_for_slide(slide: SlideSpec) -> str:
    parts = [
        slide.title,
        slide.subtitle,
        "本页重点包括",
        "，".join(slide.bullets),
        slide.narration,
    ]
    if slide.image_caption:
        parts.append(slide.image_caption)
    return "。".join(parts) + "。"


def build_video(slides: List[SlideSpec], images: List[Path]) -> bool:
    clips: List[Path] = []

    for idx, (slide, image_path) in enumerate(zip(slides, images), start=1):
        audio_path = AUDIO_DIR / f"{idx:02d}_{safe_name(slide.title)}.wav"
        clip_path = CLIPS_DIR / f"{idx:02d}_{safe_name(slide.title)}.mp4"

        ok = save_tts_audio(narration_for_slide(slide), audio_path)
        if not ok:
            create_silent_wav(8.0, audio_path)

        duration = max(7.0, wav_duration(audio_path))
        if duration < 7.0:
            create_silent_wav(7.0, audio_path)

        if not build_clip(image_path, audio_path, clip_path):
            return False
        clips.append(clip_path)

    return concat_clips(clips, VIDEO_PATH)


def build_pptx(slides: List[Path], specs: List[SlideSpec]) -> None:
    if Presentation is None:
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for img_path, spec in zip(slides, specs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(
            Inches(0.45), Inches(0.18), Inches(9.8), Inches(0.5)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = spec.title
        p.font.size = Pt(28)
        p.font.bold = True

        subtitle_box = slide.shapes.add_textbox(
            Inches(0.48), Inches(0.70), Inches(10.6), Inches(0.35)
        )
        p2 = subtitle_box.text_frame.paragraphs[0]
        p2.text = spec.subtitle
        p2.font.size = Pt(15)

        slide.shapes.add_picture(
            str(img_path),
            Inches(0.36),
            Inches(1.12),
            width=Inches(12.55),
            height=Inches(6.0),
        )

    prs.save(PPTX_PATH)


def print_summary(images: List[Path], video_ok: bool) -> None:
    print("已生成 v3 视频素材：")
    print(f"- {SCRIPT_PATH}")
    print(f"- {PPTX_PATH}")
    for img in images:
        print(f"- {img}")
    if video_ok:
        print(f"- {VIDEO_PATH}")
    else:
        print("- 视频未成功生成，请检查 ffmpeg 或音频环境。")


def main() -> None:
    ensure_dirs()

    missing = [
        path
        for path in [FRAMEWORK_IMG, CLASS_IMG, SEQUENCE_IMG, PROCESS_IMG]
        if not path.exists()
    ]
    if missing:
        print("缺少以下 UML 图片文件：")
        for item in missing:
            print(f"- {item}")
        raise SystemExit(1)

    slides = build_slides()
    build_script_markdown(slides)
    images = render_all_slides(slides)
    build_pptx(images, slides)
    video_ok = build_video(slides, images)
    print_summary(images, video_ok)


if __name__ == "__main__":
    main()
