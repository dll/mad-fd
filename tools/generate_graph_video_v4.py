from __future__ import annotations

import os
import shutil
import struct
import subprocess
import wave
from dataclasses import dataclass
from pathlib import Path
from typing import Protocol, cast

from PIL import Image, ImageDraw, ImageFont

PillowFont = ImageFont.ImageFont | ImageFont.FreeTypeFont | ImageFont.TransposedFont

try:
    import pyttsx3
except Exception:
    pyttsx3 = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None
    Inches = None
    Pt = None


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
DIAGRAMS_V3_DIR = DOCS_DIR / "diagrams" / "v3"
VIDEO_V4_DIR = DOCS_DIR / "video" / "v4"
TESTING_DIR = DOCS_DIR / "testing"
VIDEO_OUTPUT_DIR = ROOT / "video_output"

SLIDES_DIR = VIDEO_V4_DIR / "slides"
AUDIO_DIR = VIDEO_V4_DIR / "audio"
CLIPS_DIR = VIDEO_V4_DIR / "clips"
TEMP_DIR = VIDEO_V4_DIR / "temp"
CROPS_DIR = VIDEO_V4_DIR / "crops"

FRAMEWORK_IMG = DIAGRAMS_V3_DIR / "flutter_dart_framework_architecture.png"
CLASS_IMG = DIAGRAMS_V3_DIR / "flutter_dart_core_class_diagram.png"
SEQUENCE_IMG = DIAGRAMS_V3_DIR / "graph_feature_sequence_diagram.png"
PROCESS_IMG = DIAGRAMS_V3_DIR / "knowledge_graph_development_process.png"

SCRIPT_PATH = VIDEO_V4_DIR / "video_script_v4_compact.md"
PPTX_PATH = VIDEO_OUTPUT_DIR / "知识图谱核心功能_图谱功能教程_v4.pptx"
VIDEO_PATH = VIDEO_OUTPUT_DIR / "知识图谱核心功能_图谱功能教程_v4.mp4"
SRT_PATH = VIDEO_V4_DIR / "video_subtitles_v4.srt"

MODEL_TEST_PATH = ROOT / "test" / "models" / "model_test.dart"
HOME_WIDGET_TEST_PATH = ROOT / "test" / "widgets" / "home_page_widget_test.dart"
LOGIN_WIDGET_TEST_PATH = ROOT / "test" / "widget_test.dart"
TEST_CASES_PATH = TESTING_DIR / "test_cases.md"
TEST_REPORT_PATH = TESTING_DIR / "test_report.md"

LOGIN_GOLDEN_PATH = ROOT / "test" / "screenshots" / "goldens" / "login_page.png"
HOME_GOLDEN_PATH = ROOT / "test" / "screenshots" / "goldens" / "home_page.png"

WIDTH = 1920
HEIGHT = 1080

BG_TOP = (244, 248, 255)
BG_BOTTOM = (230, 238, 252)
WHITE = (255, 255, 255)
TEXT = (34, 44, 66)
TEXT_LIGHT = (92, 104, 128)
PRIMARY = (76, 110, 245)
PRIMARY_DARK = (48, 77, 201)
ACCENT = (18, 184, 134)
ORANGE = (255, 146, 43)
RED = (235, 87, 87)
PURPLE = (125, 92, 255)
BORDER = (210, 220, 240)
PANEL_BG = (250, 252, 255)
TAG_BG = (233, 239, 255)
CODE_BG = (245, 247, 251)
SUBTITLE_BG = (20, 24, 36)

TITLE_FONT_SIZE = 52
SUBTITLE_FONT_SIZE = 28
SECTION_FONT_SIZE = 28
BODY_FONT_SIZE = 25
SMALL_FONT_SIZE = 20
CODE_FONT_SIZE = 20
SUBTITLE_BAR_FONT_SIZE = 26


class VoiceLike(Protocol):
    id: str
    name: str


class TTSEngineLike(Protocol):
    def getProperty(self, name: str) -> object: ...
    def setProperty(self, name: str, value: str | float) -> None: ...
    def save_to_file(
        self, text: str, filename: str, name: str | None = None
    ) -> None: ...
    def runAndWait(self) -> None: ...
    def stop(self) -> None: ...


@dataclass
class SlideSpec:
    title: str
    subtitle: str
    bullets: list[str]
    narration: str
    image_path: Path | None = None
    image_caption: str | None = None
    code_title: str | None = None
    code_text: str | None = None
    subtitle_lines: list[str] | None = None
    teaching_tip: str | None = None
    chapter_title: str | None = None


def ensure_dirs() -> None:
    for path in [
        VIDEO_V4_DIR,
        SLIDES_DIR,
        AUDIO_DIR,
        CLIPS_DIR,
        TEMP_DIR,
        CROPS_DIR,
        VIDEO_OUTPUT_DIR,
    ]:
        path.mkdir(parents=True, exist_ok=True)


def load_font(size: int, bold: bool = False):
    candidates: list[str] = []
    if os.name == "nt":
        if bold:
            candidates.extend(
                [
                    r"C:\Windows\Fonts\msyhbd.ttc",
                    r"C:\Windows\Fonts\simhei.ttf",
                    r"C:\Windows\Fonts\arialbd.ttf",
                    r"C:\Windows\Fonts\calibrib.ttf",
                ]
            )
        else:
            candidates.extend(
                [
                    r"C:\Windows\Fonts\msyh.ttc",
                    r"C:\Windows\Fonts\simsun.ttc",
                    r"C:\Windows\Fonts\arial.ttf",
                    r"C:\Windows\Fonts\calibri.ttf",
                ]
            )
    else:
        candidates.extend(
            [
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            ]
        )

    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size=size)
        except Exception:
            continue
    return ImageFont.load_default()


FONT_TITLE = load_font(TITLE_FONT_SIZE, True)
FONT_SUBTITLE = load_font(SUBTITLE_FONT_SIZE, False)
FONT_SECTION = load_font(SECTION_FONT_SIZE, True)
FONT_BODY = load_font(BODY_FONT_SIZE, False)
FONT_BODY_BOLD = load_font(BODY_FONT_SIZE, True)
FONT_SMALL = load_font(SMALL_FONT_SIZE, False)
FONT_SMALL_BOLD = load_font(SMALL_FONT_SIZE, True)
FONT_CODE = load_font(CODE_FONT_SIZE, False)
FONT_SUBTITLE_BAR = load_font(SUBTITLE_BAR_FONT_SIZE, False)


def draw_vertical_gradient(
    img: Image.Image, top: tuple[int, int, int], bottom: tuple[int, int, int]
) -> None:
    draw = ImageDraw.Draw(img)
    for y in range(img.height):
        ratio = y / max(1, img.height - 1)
        color = tuple(int(top[i] * (1 - ratio) + bottom[i] * ratio) for i in range(3))
        draw.line([(0, y), (img.width, y)], fill=color)


def rounded_rect(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int, int, int],
    fill: tuple[int, int, int],
    outline: tuple[int, int, int] | None = None,
    radius: int = 24,
    width: int = 2,
) -> None:
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def text_size(
    draw: ImageDraw.ImageDraw, text: str, font: PillowFont
) -> tuple[int, int]:
    bbox = draw.textbbox((0, 0), text, font=font)
    return int(bbox[2] - bbox[0]), int(bbox[3] - bbox[1])


def center_text(
    draw: ImageDraw.ImageDraw,
    x: int,
    y: int,
    text: str,
    font: PillowFont,
    fill: tuple[int, int, int],
) -> None:
    w, h = text_size(draw, text, font)
    draw.text((x - w / 2, y - h / 2), text, font=font, fill=fill)


def wrap_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    font: PillowFont,
    max_width: int,
) -> list[str]:
    lines: list[str] = []
    for paragraph in text.splitlines():
        p = paragraph.strip()
        if not p:
            lines.append("")
            continue
        current = ""
        for ch in p:
            candidate = current + ch
            w, _ = text_size(draw, candidate, font)
            if w <= max_width or not current:
                current = candidate
            else:
                lines.append(current)
                current = ch
        if current:
            lines.append(current)
    return lines


def draw_wrapped_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    box: tuple[int, int, int, int],
    font: PillowFont,
    fill: tuple[int, int, int],
    line_spacing: int = 6,
) -> int:
    x1, y1, x2, y2 = box
    y = y1
    max_width = x2 - x1
    for line in wrap_text(draw, text, font, max_width):
        if y > y2:
            break
        draw.text((x1, y), line, font=font, fill=fill)
        _, h = text_size(draw, line or "中", font)
        y += h + line_spacing
    return y


def safe_name(text: str) -> str:
    bad = '\\/:*?"<>| '
    result = "".join("_" if ch in bad else ch for ch in text)
    while "__" in result:
        result = result.replace("__", "_")
    return result.strip("_") or "asset"


def read_text(path: Path, fallback: str = "") -> str:
    if not path.exists():
        return fallback
    return path.read_text(encoding="utf-8")


def extract_block_by_marker(path: Path, marker: str, max_lines: int = 28) -> str:
    if not path.exists():
        return f"// missing file: {path.name}"
    lines = path.read_text(encoding="utf-8").splitlines()
    start = 0
    for i, line in enumerate(lines):
        if marker in line:
            start = i
            break
    end = min(len(lines), start + max_lines)
    return "\n".join(lines[start:end]).strip()


def select_lines(text: str, include_keywords: list[str], max_lines: int = 18) -> str:
    lines = [line.rstrip() for line in text.splitlines() if line.strip()]
    picked: list[str] = []
    lowered = [k.lower() for k in include_keywords]
    for line in lines:
        lower_line = line.lower()
        if any(k in lower_line for k in lowered):
            picked.append(line)
        if len(picked) >= max_lines:
            break
    if not picked:
        picked = lines[:max_lines]
    return "\n".join(picked)


def load_test_summary() -> str:
    report = read_text(TEST_REPORT_PATH)
    return select_lines(
        report,
        [
            "通过",
            "测试",
            "apk",
            "homepage",
            "graphmodel",
            "nodemodel",
            "edgemodel",
            "questionmodel",
            "quizresultmodel",
            "usermodel",
        ],
        16,
    )


def load_test_case_summary() -> str:
    cases = read_text(TEST_CASES_PATH)
    return select_lines(
        cases,
        [
            "tc-",
            "图谱",
            "登录",
            "测验",
            "资料",
            "视频",
            "progress",
            "model",
        ],
        16,
    )


def load_code_summaries() -> tuple[str, str, str]:
    model_code = extract_block_by_marker(MODEL_TEST_PATH, "group('GraphModel'", 30)
    home_code = extract_block_by_marker(
        HOME_WIDGET_TEST_PATH, "group('HomePage widget tests'", 30
    )
    login_code = extract_block_by_marker(
        LOGIN_WIDGET_TEST_PATH, "testWidgets('Login page shows core UI elements'", 24
    )
    return model_code, home_code, login_code


def crop_and_save(
    source_path: Path,
    crop_box: tuple[int, int, int, int],
    output_name: str,
    border_color: tuple[int, int, int] = PRIMARY,
) -> Path:
    img = Image.open(source_path).convert("RGB")
    cropped = img.crop(crop_box)

    canvas = Image.new("RGB", (1600, 900), WHITE)
    draw_vertical_gradient(canvas, BG_TOP, BG_BOTTOM)
    draw = ImageDraw.Draw(canvas)

    rounded_rect(draw, (18, 18, 1582, 882), PANEL_BG, BORDER, 22, 3)

    target_w = 1530
    target_h = 840
    ratio = min(target_w / cropped.width, target_h / cropped.height)
    new_size = (max(1, int(cropped.width * ratio)), max(1, int(cropped.height * ratio)))
    resized = cropped.resize(new_size)
    x = 35 + (target_w - new_size[0]) // 2
    y = 30 + (target_h - new_size[1]) // 2
    canvas.paste(resized, (x, y))

    draw.rounded_rectangle(
        (18, 18, 1582, 882), radius=22, outline=border_color, width=4
    )
    out = CROPS_DIR / output_name
    canvas.save(out)
    return out


def generate_uml_crops() -> dict[str, Path]:
    crops: dict[str, Path] = {}

    if FRAMEWORK_IMG.exists():
        img = Image.open(FRAMEWORK_IMG)
        w, h = img.size
        crops["framework_full"] = FRAMEWORK_IMG
        crops["framework_ui"] = crop_and_save(
            FRAMEWORK_IMG,
            (int(w * 0.02), int(h * 0.02), int(w * 0.40), int(h * 0.56)),
            "framework_ui_crop.png",
            PRIMARY,
        )
        crops["framework_dao"] = crop_and_save(
            FRAMEWORK_IMG,
            (int(w * 0.26), int(h * 0.30), int(w * 0.64), int(h * 0.78)),
            "framework_dao_crop.png",
            ACCENT,
        )
        crops["framework_right"] = crop_and_save(
            FRAMEWORK_IMG,
            (int(w * 0.58), int(h * 0.20), int(w * 0.95), int(h * 0.66)),
            "framework_right_crop.png",
            ORANGE,
        )

    if CLASS_IMG.exists():
        img = Image.open(CLASS_IMG)
        w, h = img.size
        crops["class_full"] = CLASS_IMG
        crops["class_ui"] = crop_and_save(
            CLASS_IMG,
            (int(w * 0.02), int(h * 0.03), int(w * 0.42), int(h * 0.58)),
            "class_ui_crop.png",
            PRIMARY,
        )
        crops["class_dao"] = crop_and_save(
            CLASS_IMG,
            (int(w * 0.44), int(h * 0.22), int(w * 0.76), int(h * 0.64)),
            "class_dao_crop.png",
            ACCENT,
        )
        crops["class_model"] = crop_and_save(
            CLASS_IMG,
            (int(w * 0.56), int(h * 0.56), int(w * 0.94), int(h * 0.92)),
            "class_model_crop.png",
            PURPLE,
        )

    if SEQUENCE_IMG.exists():
        img = Image.open(SEQUENCE_IMG)
        w, h = img.size
        crops["sequence_full"] = SEQUENCE_IMG
        crops["sequence_enter"] = crop_and_save(
            SEQUENCE_IMG,
            (int(w * 0.04), 0, int(w * 0.96), int(h * 0.36)),
            "sequence_enter_crop.png",
            PRIMARY,
        )
        crops["sequence_detail"] = crop_and_save(
            SEQUENCE_IMG,
            (int(w * 0.04), int(h * 0.24), int(w * 0.96), int(h * 0.62)),
            "sequence_detail_crop.png",
            ACCENT,
        )
        crops["sequence_action"] = crop_and_save(
            SEQUENCE_IMG,
            (int(w * 0.04), int(h * 0.56), int(w * 0.96), int(h * 0.92)),
            "sequence_action_crop.png",
            ORANGE,
        )

    if PROCESS_IMG.exists():
        img = Image.open(PROCESS_IMG)
        w, h = img.size
        crops["process_full"] = PROCESS_IMG
        crops["process_early"] = crop_and_save(
            PROCESS_IMG,
            (0, 0, int(w * 1.00), int(h * 0.45)),
            "process_early_crop.png",
            PRIMARY,
        )
        crops["process_middle"] = crop_and_save(
            PROCESS_IMG,
            (0, int(h * 0.28), int(w * 1.00), int(h * 0.76)),
            "process_middle_crop.png",
            ACCENT,
        )
        crops["process_late"] = crop_and_save(
            PROCESS_IMG,
            (0, int(h * 0.58), int(w * 1.00), int(h * 1.00)),
            "process_late_crop.png",
            ORANGE,
        )

    return crops


def generate_mock_runtime_screenshots() -> dict[str, Path]:
    out: dict[str, Path] = {}

    def base_page(
        title: str,
        subtitle: str,
        chips: list[str],
        cards: list[tuple[str, str, tuple[int, int, int]]],
        output_name: str,
    ) -> Path:
        img = Image.new("RGB", (1080, 2400), WHITE)
        draw_vertical_gradient(img, (245, 247, 255), (233, 238, 252))
        draw = ImageDraw.Draw(img)

        rounded_rect(draw, (0, 0, 1080, 165), PRIMARY, PRIMARY, 0, 0)
        draw.text((40, 42), title, font=load_font(44, True), fill=WHITE)
        draw.text((40, 102), subtitle, font=load_font(24, False), fill=(235, 240, 255))

        y = 210
        x = 40
        for chip in chips:
            w, _ = text_size(draw, chip, load_font(26, True))
            rounded_rect(draw, (x, y, x + w + 40, y + 56), TAG_BG, TAG_BG, 18, 1)
            draw.text(
                (x + 20, y + 14), chip, font=load_font(26, True), fill=PRIMARY_DARK
            )
            x += w + 60

        y = 330
        for card_title, card_sub, color in cards:
            rounded_rect(draw, (40, y, 1040, y + 180), WHITE, BORDER, 26, 3)
            rounded_rect(draw, (70, y + 40, 150, y + 120), color, color, 20, 1)
            center_text(draw, 110, y + 80, "●", load_font(26, True), WHITE)
            draw.text((190, y + 34), card_title, font=load_font(34, True), fill=TEXT)
            draw.text(
                (190, y + 92), card_sub, font=load_font(24, False), fill=TEXT_LIGHT
            )
            y += 220

        out_path = CROPS_DIR / output_name
        img.save(out_path)
        return out_path

    if LOGIN_GOLDEN_PATH.exists():
        out["login"] = LOGIN_GOLDEN_PATH
    else:
        out["login"] = base_page(
            "移动应用开发知识图谱",
            "登录页运行效果（页面渲染截图）",
            ["Flutter", "Dart", "登录入口"],
            [
                ("学号 / 工号输入", "用于学生、教师、管理员统一登录入口", PRIMARY),
                ("密码输入", "学生支持后六位密码，管理员支持专用密码", ACCENT),
                ("快速登录区", "提供测试学生、教师、管理员快捷体验", ORANGE),
            ],
            "runtime_login_mock.png",
        )

    if HOME_GOLDEN_PATH.exists():
        out["home"] = HOME_GOLDEN_PATH
    else:
        out["home"] = base_page(
            "HomePage 导航中心",
            "首页运行效果（页面渲染截图）",
            ["首页", "图谱", "资料", "视频", "测验", "进度"],
            [
                ("知识图谱", "进入课程图谱、技术栈图谱和学习图谱等主题", PRIMARY),
                ("章节测验", "按章节组织题目并记录答题结果", ORANGE),
                ("视频教程", "统一管理课程视频资源", RED),
                ("课程资料", "支持 PDF / PPT 课件访问", ACCENT),
            ],
            "runtime_home_mock.png",
        )

    out["graph_list"] = base_page(
        "GraphListPage 图谱列表",
        "图谱页运行效果（页面示意图）",
        ["图谱入口", "列表加载", "图谱选择"],
        [
            ("01 课程图谱详细图谱", "展示课程整体知识结构与章节关系", PRIMARY),
            ("02 技术栈图谱详细图谱", "展示 Flutter、Dart、SQLite 等技术关系", ACCENT),
            ("06 学习图谱详细图谱", "展示学习路径、复习节点与实践内容", ORANGE),
        ],
        "runtime_graph_list_mock.png",
    )

    out["graph_detail"] = base_page(
        "GraphDetailPage 图谱详情",
        "图谱详情运行效果（页面示意图）",
        ["节点交互", "缩放拖拽", "学习记录", "收藏"],
        [
            (
                "图谱可视化区域",
                "通过 CustomPainter 与 InteractiveViewer 浏览结构",
                PRIMARY,
            ),
            ("节点详情卡片", "点击节点后显示标题、类型、内容和操作按钮", ACCENT),
            ("开始学习 / 收藏", "把浏览行为转化为学习记录和收藏数据", RED),
        ],
        "runtime_graph_detail_mock.png",
    )

    return out


def build_script_markdown(slides: list[SlideSpec]) -> None:
    lines = [
        "# 知识图谱核心功能教学视频脚本 v4",
        "",
        "## 说明",
        "",
        "- 本版面向教学视频展示，强调讲解节奏、学习要点和字幕可读性。",
        "- 成片中不展示脚本字段名称。",
        "- 成片中展示标题、教学图示、学习提示、测试结果与字幕。",
        "",
        "## 成片结构",
        "",
    ]
    for i, slide in enumerate(slides, start=1):
        lines.append(f"{i}. {slide.title} - {slide.subtitle}")
    lines += ["", "## 讲解稿", ""]
    for i, slide in enumerate(slides, start=1):
        lines += [
            f"### {i}. {slide.title}",
            f"- 副标题：{slide.subtitle}",
            f"- 学习要点：{', '.join(slide.bullets)}",
            f"- 讲解：{slide.narration}",
            "",
        ]
    _ = SCRIPT_PATH.write_text("\n".join(lines), encoding="utf-8")


def build_slides(
    crops: dict[str, Path], runtime_shots: dict[str, Path]
) -> list[SlideSpec]:
    model_code, home_code, login_code = load_code_summaries()
    test_case_summary = load_test_case_summary()
    test_report_summary = load_test_summary()

    return [
        SlideSpec(
            title="课程导入",
            subtitle="先了解学习目标与观看重点",
            bullets=[
                "了解系统整体结构",
                "理解核心页面职责",
                "掌握图谱数据流与交互流程",
                "观察测试与运行结果",
            ],
            narration=(
                "欢迎进入知识图谱核心功能教学视频。"
                "本节内容会按照系统结构、页面职责、交互流程、开发过程和测试结果逐步展开。"
                "建议你先带着学习目标观看，再结合后续图示理解整个功能链路。"
            ),
            image_path=crops.get("framework_full"),
            image_caption="课程导入：先建立整体认识，再逐步进入细节讲解",
            subtitle_lines=[
                "欢迎进入知识图谱核心功能教学视频。",
                "本节将按结构、页面、流程、测试四个维度展开讲解。",
            ],
            teaching_tip="先记住学习目标，再带着问题观看后续内容。",
            chapter_title="第 0 讲 课程导入",
        ),
        SlideSpec(
            title="项目总览",
            subtitle="认识知识图谱学习系统的整体结构",
            bullets=["Flutter", "Dart", "SQLite", "UML", "知识图谱", "测试资产"],
            narration=(
                "本视频聚焦 Flutter 与 Dart 技术栈下的知识图谱学习系统设计。"
                "重点说明系统架构、核心类、顺序流程、开发过程、测试资产以及运行效果。"
            ),
            image_path=crops.get("framework_full"),
            image_caption="项目总览：技术栈、UML 与知识图谱功能一体化展示",
            subtitle_lines=[
                "本节先从项目总览开始。",
                "重点认识技术栈、结构图和知识图谱在系统中的位置。",
            ],
            teaching_tip="先看整体，再看细节，避免一开始陷入局部实现。",
            chapter_title="第 1 讲 项目总览",
        ),
        SlideSpec(
            title="页面职责分析",
            subtitle="从页面分工理解系统学习闭环",
            bullets=[
                "LoginPage 统一登录入口",
                "HomePage 导航中心",
                "GraphListPage 图谱入口",
                "GraphDetailPage 核心交互页",
                "Quiz / Video / Document / Progress 联动",
            ],
            narration=(
                "项目中的页面并不是简单罗列名称。LoginPage 负责统一登录入口，HomePage 负责导航分发，"
                "GraphListPage 负责图谱主题选择，GraphDetailPage 负责图谱核心交互。"
                "另外，资料、视频、测验和进度页面围绕知识图谱共同构成学习闭环。"
            ),
            image_path=runtime_shots.get("home"),
            image_caption="页面作用说明：入口、导航、图谱、资料、视频、测验、进度",
            code_title="页面职责摘要",
            code_text=(
                "LoginPage   -> 用户登录入口\n"
                "HomePage    -> 统一导航中心\n"
                "GraphListPage -> 图谱主题列表\n"
                "GraphDetailPage -> 节点交互与学习动作\n"
                "DocumentListPage -> PDF / PPT 资料\n"
                "VideoListPage -> 视频资源\n"
                "QuizPage -> 章节测验\n"
                "ProgressPage -> 学习反馈"
            ),
            subtitle_lines=[
                "先看页面分工，不要只记页面名称。",
                "更重要的是理解每个页面在学习闭环中的作用。",
            ],
            teaching_tip="把入口页、导航页、功能页和反馈页区分开来理解。",
            chapter_title="第 2 讲 页面职责分析",
        ),
        SlideSpec(
            title="系统框架总览",
            subtitle="先建立页面层、逻辑层和数据层的整体认识",
            bullets=[
                "Flutter UI 层",
                "DAO / Service 层",
                "SQLite 数据层",
                "图谱驱动学习",
            ],
            narration=(
                "这张框架图展示了系统的整体技术结构。"
                "上层是 Flutter 页面层，中层是业务逻辑与 DAO 层，下层是 SQLite 数据层。"
                "图谱功能位于整个系统核心位置，向上连接页面，向下连接数据，向外联动测验、资料、视频与进度。"
            ),
            image_path=crops.get("framework_full"),
            image_caption="系统框架图：从页面层到数据库层的总体结构",
            subtitle_lines=[
                "这一页先建立整体框架认识。",
                "请重点观察页面层、逻辑层和数据层之间的关系。",
            ],
            teaching_tip="看框架图时，先分层，再看层与层之间如何协作。",
            chapter_title="第 3 讲 系统框架总览",
        ),
        SlideSpec(
            title="UI 层讲解",
            subtitle="把页面名称与页面职责一一对应起来",
            bullets=[
                "LoginPage：登录入口",
                "HomePage：导航中心",
                "GraphListPage：图谱列表",
                "GraphDetailPage：节点交互",
                "其他页面：资料、视频、测验、进度、收藏",
            ],
            narration=(
                "这里放大的是框架图中的 UI 层。"
                "LoginPage 是统一登录入口，HomePage 是导航中心，GraphListPage 是图谱列表入口，"
                "GraphDetailPage 是图谱核心交互页面。"
                "DocumentListPage、VideoListPage、QuizPage、ProgressPage 和 FavoritesPage 则分别承担资源、测验和学习反馈功能。"
            ),
            image_path=crops.get("framework_ui"),
            image_caption="UI 层放大图：页面名称之外，更重要的是页面职责",
            subtitle_lines=[
                "现在放大 UI 层。",
                "请把页面名称和页面职责一一对应起来理解。",
            ],
            teaching_tip="不要孤立记页面名，要结合它在学习流程中的作用记忆。",
            chapter_title="第 4 讲 UI 层讲解",
        ),
        SlideSpec(
            title="DAO 与数据层讲解",
            subtitle="理解页面背后的业务逻辑与数据库访问",
            bullets=[
                "AuthService 管理登录状态",
                "GraphDao 读取图谱、节点、边",
                "FavoriteDao 管理收藏",
                "LearningRecordDao 管理学习记录",
                "DatabaseHelper 统一管理 SQLite",
            ],
            narration=(
                "这里放大的是框架图中的 DAO 与数据层。"
                "AuthService 管理用户状态，GraphDao 负责图谱结构读取，FavoriteDao 负责收藏读写，"
                "LearningRecordDao 负责学习记录统计与保存，DatabaseHelper 负责统一管理 SQLite 初始化和建表逻辑。"
                "这部分说明 Dart 业务层不是附属代码，而是整个系统的运行核心。"
            ),
            image_path=crops.get("framework_dao"),
            image_caption="DAO / 数据层放大图：页面请求如何转化为数据库操作",
            subtitle_lines=[
                "接下来观察 DAO 与数据层。",
                "这里对应的是页面背后的核心业务逻辑和数据库访问过程。",
            ],
            teaching_tip="把页面请求、DAO 调用和数据库读写串成一条完整链路。",
            chapter_title="第 5 讲 DAO 与数据层讲解",
        ),
        SlideSpec(
            title="核心类图总览",
            subtitle="从类图理解页面层、DAO 层和模型层",
            bullets=[
                "UI Layer 负责交互",
                "DAO Layer 负责数据访问",
                "Model Layer 负责结构化对象",
                "GraphDetailPage 是核心类",
            ],
            narration=(
                "类图展示了项目中最核心的类以及它们之间的职责关系。"
                "页面层负责显示与交互，DAO 层负责数据访问，模型层负责结构化数据对象。"
                "其中，GraphDetailPage 是知识图谱功能的核心类。"
            ),
            image_path=crops.get("class_full"),
            image_caption="类图整体：页面层、DAO 层、模型层的职责关系",
            subtitle_lines=[
                "这一页从类图角度理解系统。",
                "重点不是记住类名，而是看清职责边界。",
            ],
            teaching_tip="类图学习的关键不是背类名，而是掌握职责边界。",
            chapter_title="第 6 讲 核心类图总览",
        ),
        SlideSpec(
            title="UI 类职责讲解",
            subtitle="区分应用启动、导航分发与图谱交互类",
            bullets=[
                "MyApp：应用启动",
                "HomePage：导航中心",
                "GraphListPage：图谱入口",
                "GraphDetailPage：图谱交互核心",
                "Progress / Favorites：学习反馈与复习入口",
            ],
            narration=(
                "放大 UI 类之后可以更清楚看到，页面类并不是简单堆叠。"
                "MyApp 负责应用启动，HomePage 负责主导航，GraphListPage 负责图谱入口分发，"
                "GraphDetailPage 负责图谱节点、边、布局和交互处理，"
                "ProgressPage 和 FavoritesPage 则负责结果反馈与复习入口。"
            ),
            image_path=crops.get("class_ui"),
            image_caption="类图 UI 层放大图：每个页面都有明确职责",
            subtitle_lines=[
                "这里聚焦 UI 类。",
                "请把应用启动、导航分发和图谱交互这三层作用区分开。",
            ],
            teaching_tip="先判断类属于启动、导航还是功能交互，再理解其作用。",
            chapter_title="第 7 讲 UI 类职责讲解",
        ),
        SlideSpec(
            title="DAO 与模型层讲解",
            subtitle="理解数据查询、封装与对象协同",
            bullets=[
                "GraphDao：图谱查询",
                "FavoriteDao：收藏管理",
                "LearningRecordDao：学习记录",
                "GraphModel / NodeModel / EdgeModel：图结构对象",
            ],
            narration=(
                "放大类图的 DAO 和模型层之后，可以看到 Dart 逻辑层是整个系统的数据骨架。"
                "GraphDao 负责图谱查询，FavoriteDao 负责收藏管理，LearningRecordDao 负责学习记录。"
                "GraphModel、NodeModel 和 EdgeModel 则把数据库结果转换成可以被页面层直接消费的对象结构。"
            ),
            image_path=crops.get("class_dao"),
            image_caption="类图 DAO / Model 放大图：逻辑层与对象层的协同",
            code_title="模型测试代码片段",
            code_text=model_code,
            subtitle_lines=[
                "这一页重点理解逻辑层与模型层。",
                "它们决定了图谱数据如何被查询、封装和使用。",
            ],
            teaching_tip="把 DAO 看成数据入口，把 Model 看成页面可直接使用的数据对象。",
            chapter_title="第 8 讲 DAO 与模型层讲解",
        ),
        SlideSpec(
            title="进入图谱模块流程",
            subtitle="理解用户进入图谱模块的调用链",
            bullets=[
                "HomePage 切换图谱导航",
                "GraphListPage 加载图谱列表",
                "GraphDao 查询数据库",
                "SQLite 返回图谱数据",
            ],
            narration=(
                "顺序图的第一段描述的是用户如何进入图谱模块。"
                "用户从 HomePage 切换到底部导航中的图谱页，GraphListPage 启动后调用 GraphDao，"
                "GraphDao 再通过 DatabaseHelper 访问 SQLite，最终把图谱列表返回给页面。"
            ),
            image_path=crops.get("sequence_enter"),
            image_caption="顺序图放大图：进入图谱模块的调用链",
            subtitle_lines=[
                "现在开始看顺序图。",
                "先理解用户是如何一步步进入图谱模块的。",
            ],
            teaching_tip="顺序图阅读时要从用户动作出发，沿着调用顺序向下看。",
            chapter_title="第 9 讲 进入图谱模块流程",
        ),
        SlideSpec(
            title="图谱详情加载流程",
            subtitle="理解节点、边与布局计算的加载过程",
            bullets=[
                "点击图谱卡片进入详情页",
                "读取 nodes",
                "读取 edges",
                "计算布局并显示图谱",
            ],
            narration=(
                "顺序图的第二段描述的是进入图谱详情页之后的逻辑。"
                "当用户点击某个图谱卡片后，GraphDetailPage 会分别查询 nodes 和 edges，"
                "随后进行节点布局计算，并把结果渲染成可浏览、可缩放、可拖拽的图谱视图。"
            ),
            image_path=crops.get("sequence_detail"),
            image_caption="顺序图放大图：详情页查询节点、边并构建图谱视图",
            subtitle_lines=[
                "这一步进入图谱详情页。",
                "请观察节点、边和布局计算是如何串起来的。",
            ],
            teaching_tip="把数据读取和布局计算分成两个阶段理解，会更清楚。",
            chapter_title="第 10 讲 图谱详情加载流程",
        ),
        SlideSpec(
            title="图谱交互动作流程",
            subtitle="理解开始学习与收藏如何沉淀为行为数据",
            bullets=[
                "点击节点显示详情卡片",
                "开始学习写入 learning_records",
                "收藏写入 favorites",
                "图谱从展示页变成行为入口",
            ],
            narration=(
                "顺序图最后一段最关键。"
                "用户点击节点后，GraphDetailPage 会展示详情卡片。"
                "如果点击开始学习，就把记录写入 learning_records；如果点击收藏，就写入 favorites。"
                "这说明知识图谱不仅是结构展示页，更是学习行为入口。"
            ),
            image_path=crops.get("sequence_action"),
            image_caption="顺序图放大图：节点交互如何沉淀为学习行为",
            code_title="首页组件测试代码片段",
            code_text=home_code,
            subtitle_lines=[
                "这一页是交互流程的核心。",
                "知识图谱从展示工具变成了真正的学习行为入口。",
            ],
            teaching_tip="重点关注交互动作如何写入学习记录和收藏数据。",
            chapter_title="第 11 讲 图谱交互动作流程",
        ),
        SlideSpec(
            title="开发过程：前期阶段",
            subtitle="需求分析、业务建模与数据库设计",
            bullets=[
                "明确图谱是系统核心",
                "定义 Graph / Node / Edge / Favorite / LearningRecord",
                "设计 SQLite 表结构",
            ],
            narration=(
                "开发过程的前期重点是需求分析和数据建模。"
                "首先要明确知识图谱是系统核心，再定义 Graph、Node、Edge、Favorite 和 LearningRecord 等业务对象，"
                "随后把它们映射为 SQLite 表结构，为后续开发奠定数据基础。"
            ),
            image_path=crops.get("process_early"),
            image_caption="开发过程放大图：需求分析、业务对象定义与数据库设计",
            subtitle_lines=[
                "接下来回到开发过程。",
                "前期工作的重点是需求分析、建模和数据库设计。",
            ],
            teaching_tip="前期要先把业务对象和数据表设计清楚，再进入编码阶段。",
            chapter_title="第 12 讲 开发过程：前期阶段",
        ),
        SlideSpec(
            title="开发过程：中期阶段",
            subtitle="DAO 开发、页面实现与图谱交互构建",
            bullets=[
                "实现 DatabaseHelper",
                "实现 GraphDao / FavoriteDao / LearningRecordDao",
                "实现 GraphListPage / GraphDetailPage",
                "构建图谱布局和节点交互",
            ],
            narration=(
                "开发过程的中期重点是实现 DAO 层和页面层。"
                "先实现 DatabaseHelper，再实现 GraphDao、FavoriteDao 和 LearningRecordDao，"
                "接着完成 GraphListPage 和 GraphDetailPage，最终实现图谱布局、缩放拖拽和节点交互能力。"
            ),
            image_path=crops.get("process_middle"),
            image_caption="开发过程放大图：DAO 层与图谱页面的主要实现阶段",
            subtitle_lines=[
                "中期开发进入实现阶段。",
                "这时要把数据库、页面和图谱交互真正串联起来。",
            ],
            teaching_tip="中期开发的目标是把数据层、页面层和交互层真正打通。",
            chapter_title="第 13 讲 开发过程：中期阶段",
        ),
        SlideSpec(
            title="开发过程：后期阶段",
            subtitle="测试验证、构建发布与教学材料交付",
            bullets=[
                "补充模型测试与组件测试",
                "生成测试用例与测试报告",
                "构建 Debug / Release APK",
                "生成 UML、PPT、视频材料",
            ],
            narration=(
                "开发过程的后期重点是验证与交付。"
                "项目补充了模型测试和组件测试，生成了测试用例与测试报告，"
                "同时成功构建 Debug 和 Release APK，最后再输出 UML 图、PPT 和视频材料，"
                "形成适合教学展示与项目答辩的完整交付物。"
            ),
            image_path=crops.get("process_late"),
            image_caption="开发过程放大图：测试验证、APK 构建与展示材料生成",
            subtitle_lines=[
                "后期工作的核心是验证与交付。",
                "教学展示需要的不只是代码，还包括测试和成片资料。",
            ],
            teaching_tip="后期不仅要保证功能可用，还要保证成果能够被清晰展示。",
            chapter_title="第 14 讲 开发过程：后期阶段",
        ),
        SlideSpec(
            title="测试资产讲解",
            subtitle="从测试代码与测试用例验证系统质量",
            bullets=[
                "model_test.dart",
                "widget_test.dart",
                "home_page_widget_test.dart",
                "test_cases.md",
                "23 项测试通过",
            ],
            narration=(
                "项目不仅有 UML 图和页面实现，也补充了测试代码与测试用例。"
                "模型测试覆盖 GraphModel、NodeModel、EdgeModel、QuestionModel、QuizResultModel 和 UserModel。"
                "组件测试覆盖登录页和首页导航。"
                "测试用例文档则从图谱、资料、视频、测验和进度等多个模块给出了系统化验证点。"
            ),
            image_path=runtime_shots.get("login"),
            image_caption="测试资产：测试代码、测试用例、页面渲染截图共同支撑展示",
            code_title="测试用例摘要",
            code_text=test_case_summary,
            subtitle_lines=[
                "这一页看测试资产。",
                "教学展示不仅要讲功能，也要证明功能经过了验证。",
            ],
            teaching_tip="测试资产的价值在于证明系统不仅能展示，还能被验证。",
            chapter_title="第 15 讲 测试资产讲解",
        ),
        SlideSpec(
            title="测试结果与交付成果",
            subtitle="理解测试、构建与展示材料的闭环关系",
            bullets=[
                "23 项测试通过",
                "Debug APK 构建成功",
                "Release APK 构建成功",
                "测试报告已整理",
                "视频 / PPT / UML 已生成",
            ],
            narration=(
                "测试报告总结了当前工程状态。"
                "目前模型层测试与组件测试已经执行通过，共通过 23 项测试。"
                "同时，Android 的 Debug 与 Release APK 也已经成功构建。"
                "在此基础上，项目已经具备测试报告、UML 图、PPT 和视频等完整展示资料。"
            ),
            image_path=runtime_shots.get("graph_list"),
            image_caption="测试与交付：代码质量验证与展示资料输出同时完成",
            code_title="测试报告摘要",
            code_text=test_report_summary,
            subtitle_lines=[
                "这里总结验证结果与交付状态。",
                "可以看到测试、构建和展示材料已经形成闭环。",
            ],
            teaching_tip="展示项目成熟度时，要把测试结果和交付成果一起说明。",
            chapter_title="第 16 讲 测试结果与交付成果",
        ),
        SlideSpec(
            title="实际运行界面讲解",
            subtitle="把结构图理解与真实页面效果对应起来",
            bullets=[
                "登录页：统一入口",
                "首页：导航中心",
                "图谱列表页：图谱主题选择",
                "图谱详情页：节点交互与学习动作",
            ],
            narration=(
                "除了 UML 图，本版本还补充了页面运行效果截图。"
                "登录页负责统一入口，首页负责导航中心，图谱列表页负责主题选择，"
                "图谱详情页负责节点交互、开始学习和收藏动作。"
                "这样在视频中不仅能看到结构图，也能看到页面层面对应的实际界面。"
            ),
            image_path=runtime_shots.get("graph_detail"),
            image_caption="运行界面说明：页面不是只有名称，而是有清晰的业务职责",
            code_title="登录页组件测试片段",
            code_text=login_code,
            subtitle_lines=[
                "最后补充实际运行界面。",
                "这样你可以把结构图理解和真实页面效果对应起来。",
            ],
            teaching_tip="学到最后一定要把结构图、流程图和真实页面对应起来。",
            chapter_title="第 17 讲 实际运行界面讲解",
        ),
        SlideSpec(
            title="本节回顾",
            subtitle="回顾本节知识点与教学展示优化内容",
            bullets=[
                "增加字幕",
                "UML 放大裁切",
                "页面作用说明",
                "测试代码 / 用例 / 报告",
                "运行截图补充",
            ],
            narration=(
                "v4 版本重点解决了前一版中 UML 图不清晰、页面说明不足、测试展示不完整以及运行截图缺失的问题。"
                "现在视频同时具备字幕、局部放大 UML 图、页面职责解释、测试资产和运行界面展示，"
                "更适合课程汇报、项目答辩和教学展示。"
            ),
            image_path=crops.get("framework_right"),
            image_caption="v4 版本总结：结构更清晰、解释更完整、展示更可读",
            subtitle_lines=[
                "这一页回顾本版优化点。",
                "重点是让结构更清晰、讲解更完整、展示更适合教学。",
            ],
            teaching_tip="回顾时优先抓住结构、流程、测试和界面四条主线。",
            chapter_title="第 18 讲 本节回顾",
        ),
        SlideSpec(
            title="课程总结与下节预告",
            subtitle="本节回顾与后续学习方向说明",
            bullets=[
                "本节回顾：系统框架、类图、顺序图",
                "本节回顾：页面职责、测试资产、运行界面",
                "下节预告：数据库初始化与图谱数据装载",
                "下节预告：节点布局与交互源码阅读",
            ],
            narration=(
                "本节教学视频到这里结束。"
                "你可以先回看系统框架图、核心类图和顺序图，巩固知识图谱功能的整体认识。"
                "下一步建议继续学习数据库初始化、图谱数据装载、节点布局计算和交互源码阅读，这样可以把展示理解进一步过渡到实现理解。"
            ),
            image_path=crops.get("framework_right"),
            image_caption="课程片尾：回顾本节重点，并为后续源码学习建立方向",
            subtitle_lines=[
                "本节内容到这里结束。",
                "下节建议继续学习数据库初始化、数据装载与交互实现。",
            ],
            teaching_tip="课程结束后，先回顾结构与流程，再进入更细的源码实现学习。",
            chapter_title="第 19 讲 课程总结与下节预告",
        ),
    ]


def paste_image(
    base: Image.Image, source: Path | None, box: tuple[int, int, int, int]
) -> None:
    if source is None or not source.exists():
        return
    img = Image.open(source).convert("RGB")
    target_w = box[2] - box[0]
    target_h = box[3] - box[1]
    ratio = min(target_w / img.width, target_h / img.height)
    new_size = (max(1, int(img.width * ratio)), max(1, int(img.height * ratio)))
    img = img.resize(new_size)
    x = box[0] + (target_w - new_size[0]) // 2
    y = box[1] + (target_h - new_size[1]) // 2
    base.paste(img, (x, y))


def draw_bullets(
    draw: ImageDraw.ImageDraw, bullets: list[str], box: tuple[int, int, int, int]
) -> None:
    x1, y1, x2, y2 = box
    y = y1
    for bullet in bullets:
        rounded_rect(draw, (x1, y + 8, x1 + 18, y + 26), PRIMARY, PRIMARY, 6, 1)
        y = (
            draw_wrapped_text(
                draw, bullet, (x1 + 32, y, x2, y + 64), FONT_BODY, TEXT, 8
            )
            + 10
        )
        if y > y2:
            break


def draw_code_block(
    draw: ImageDraw.ImageDraw, title: str, code: str, box: tuple[int, int, int, int]
) -> None:
    rounded_rect(draw, box, CODE_BG, (224, 230, 242), 18, 2)
    x1, y1, _x2, y2 = box
    draw.text((x1 + 18, y1 + 14), title, font=FONT_SMALL_BOLD, fill=PRIMARY_DARK)
    y = y1 + 52
    for line in code.splitlines():
        if y > y2 - 24:
            break
        draw.text((x1 + 18, y), line.rstrip(), font=FONT_CODE, fill=TEXT)
        _, h = text_size(draw, line or "中", FONT_CODE)
        y += h + 5


def build_subtitle_text(slide: SlideSpec) -> str:
    if slide.subtitle_lines:
        return "\n".join(slide.subtitle_lines[:2])

    text = slide.narration.strip()
    max_chars = 34
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1] + "…"


def split_subtitle_segments(slide: SlideSpec) -> list[str]:
    if slide.subtitle_lines:
        return [line.strip() for line in slide.subtitle_lines if line.strip()]

    text = slide.narration.strip()
    chunks: list[str] = []
    current = ""

    for ch in text:
        current += ch
        if ch in "，。！？；":
            trimmed = current.strip()
            if trimmed:
                chunks.append(trimmed)
            current = ""

    tail = current.strip()
    if tail:
        chunks.append(tail)

    if not chunks:
        return [slide.title]

    merged: list[str] = []
    current = ""
    for chunk in chunks:
        candidate = current + chunk
        if len(candidate) <= 24 or not current:
            current = candidate
        else:
            merged.append(current)
            current = chunk
    if current:
        merged.append(current)

    return merged


def render_slide(slide: SlideSpec, index: int, total: int) -> Path:
    img = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw_vertical_gradient(img, BG_TOP, BG_BOTTOM)
    draw = ImageDraw.Draw(img)

    rounded_rect(draw, (48, 34, 1870, 150), WHITE, BORDER, 28, 2)
    draw.text((84, 50), slide.title, font=FONT_TITLE, fill=TEXT)
    draw.text((88, 108), slide.subtitle, font=FONT_SUBTITLE, fill=TEXT_LIGHT)

    rounded_rect(draw, (1660, 54, 1835, 108), PRIMARY, PRIMARY, 18, 1)
    center_text(draw, 1748, 82, f"{index}/{total}", FONT_SMALL_BOLD, WHITE)

    left_box = (56, 180, 650, 910)
    right_box = (680, 180, 1862, 765)
    bottom_box = (680, 790, 1862, 930)

    rounded_rect(draw, left_box, WHITE, BORDER, 30, 3)
    rounded_rect(draw, right_box, WHITE, BORDER, 30, 3)
    rounded_rect(draw, bottom_box, WHITE, BORDER, 30, 3)

    draw.text((90, 210), "学习要点", font=FONT_SECTION, fill=PRIMARY_DARK)
    draw_bullets(draw, slide.bullets, (90, 282, 590, 610))

    rounded_rect(draw, (88, 645, 360, 690), TAG_BG, TAG_BG, 16, 1)
    _ = draw.text((110, 657), "教学提示", font=FONT_SMALL_BOLD, fill=PRIMARY_DARK)
    tip_text = (
        slide.teaching_tip
        or slide.image_caption
        or "请结合当前图示归纳本页关键知识点。"
    )
    _ = draw_wrapped_text(draw, tip_text, (90, 706, 595, 790), FONT_BODY, TEXT)

    rounded_rect(draw, (88, 804, 360, 848), TAG_BG, TAG_BG, 16, 1)
    _ = draw.text((110, 816), "图示说明", font=FONT_SMALL_BOLD, fill=PRIMARY_DARK)
    if slide.image_caption:
        _ = draw_wrapped_text(
            draw, slide.image_caption, (90, 864, 595, 910), FONT_SMALL, TEXT
        )

    _ = draw.text((712, 210), "教学图示 / 运行图", font=FONT_SECTION, fill=PRIMARY_DARK)
    _ = paste_image(img, slide.image_path, (680, 248, 1860, 780))

    if slide.code_title and slide.code_text:
        draw_code_block(draw, slide.code_title, slide.code_text, bottom_box)
    else:
        draw_code_block(
            draw,
            "课程说明",
            "本页重点展示课程结构、学习要点和教学解释，不直接显示脚本字段。",
            bottom_box,
        )

    out_path = SLIDES_DIR / f"{index:02d}_{safe_name(slide.title)}.png"
    img.save(out_path)
    return out_path


def render_all_slides(slides: list[SlideSpec]) -> list[Path]:
    return [
        render_slide(slide, i, len(slides)) for i, slide in enumerate(slides, start=1)
    ]


def choose_voice(engine: TTSEngineLike, preferred_keywords: list[str]) -> str | None:
    voices_obj = engine.getProperty("voices")
    voices = cast(list[VoiceLike], voices_obj) if isinstance(voices_obj, list) else []
    lowered = [key.lower() for key in preferred_keywords]
    for voice in voices:
        combined = f"{voice.id} {voice.name}".lower()
        if any(key in combined for key in lowered):
            return voice.id
    return voices[0].id if voices else None


def clean_tts_text(text: str) -> str:
    replacements = {
        "Flutter": "Flutter",
        "Dart": "Dart",
        "SQLite": "S Q Lite",
        "GraphDao": "Graph Dao",
        "FavoriteDao": "Favorite Dao",
        "LearningRecordDao": "Learning Record Dao",
        "DatabaseHelper": "Database Helper",
        "GraphListPage": "Graph List Page",
        "GraphDetailPage": "Graph Detail Page",
        "HomePage": "Home Page",
        "QuizPage": "Quiz Page",
        "ProgressPage": "Progress Page",
        "LoginPage": "Login Page",
        "UML": "U M L",
        "APK": "A P K",
        "favorites": "favorites",
        "learning_records": "learning records",
    }
    result = text
    for old, new in replacements.items():
        result = result.replace(old, new)
    return result


def save_tts_audio(text: str, out_path: Path, rate: int = 158) -> bool:
    if pyttsx3 is None:
        return False
    try:
        engine = cast(TTSEngineLike, pyttsx3.init())
        voice_id = choose_voice(
            engine, ["zh", "chinese", "huihui", "hanhan", "xiaoxiao"]
        )
        if voice_id:
            _ = engine.setProperty("voice", voice_id)
        _ = engine.setProperty("rate", rate)
        _ = engine.save_to_file(clean_tts_text(text), str(out_path))
        _ = engine.runAndWait()
        _ = engine.stop()
        return out_path.exists() and out_path.stat().st_size > 0
    except Exception:
        return False


def create_silent_wav(duration_seconds: float, out_path: Path) -> None:
    sample_rate = 22050
    frames = max(1, int(duration_seconds * sample_rate))
    silence = struct.pack("<h", 0)
    with wave.open(str(out_path), "wb") as wav_file:
        wav_file.setnchannels(1)
        wav_file.setsampwidth(2)
        wav_file.setframerate(sample_rate)
        for _ in range(frames):
            wav_file.writeframesraw(silence)


def wav_duration(path: Path) -> float:
    with wave.open(str(path), "rb") as wav_file:
        return wav_file.getnframes() / float(wav_file.getframerate())


def which(names: list[str]) -> str | None:
    for name in names:
        found = shutil.which(name)
        if found:
            return found
    return None


def ffmpeg_path() -> str | None:
    return which(
        ["ffmpeg", r"C:\Users\ldl\AppData\Local\Microsoft\WinGet\Links\ffmpeg.exe"]
    )


def run_ffmpeg(args: list[str]) -> bool:
    ffmpeg = ffmpeg_path()
    if not ffmpeg:
        return False
    process = subprocess.run(
        [ffmpeg] + args, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False
    )
    return process.returncode == 0


def narration_for_slide(slide: SlideSpec) -> str:
    return slide.narration.strip() + "。"


def build_clip(image_path: Path, audio_path: Path, output_path: Path) -> bool:
    return run_ffmpeg(
        [
            "-y",
            "-loop",
            "1",
            "-i",
            str(image_path),
            "-i",
            str(audio_path),
            "-c:v",
            "libx264",
            "-tune",
            "stillimage",
            "-c:a",
            "aac",
            "-b:a",
            "192k",
            "-pix_fmt",
            "yuv420p",
            "-shortest",
            "-vf",
            "scale=1920:1080,fps=30",
            str(output_path),
        ]
    )


def concat_clips(clips: list[Path], output_path: Path) -> bool:
    ffmpeg = ffmpeg_path()
    if not ffmpeg:
        return False
    manifest = TEMP_DIR / "clips_manifest_v4.txt"
    _ = manifest.write_text(
        "\n".join(f"file '{clip.as_posix()}'" for clip in clips), encoding="utf-8"
    )
    process = subprocess.run(
        [
            ffmpeg,
            "-y",
            "-f",
            "concat",
            "-safe",
            "0",
            "-i",
            str(manifest),
            "-c",
            "copy",
            str(output_path),
        ],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
    )
    return process.returncode == 0 and output_path.exists()


def seconds_to_srt(ts: float) -> str:
    hours = int(ts // 3600)
    minutes = int((ts % 3600) // 60)
    seconds = int(ts % 60)
    millis = int((ts - int(ts)) * 1000)
    return f"{hours:02d}:{minutes:02d}:{seconds:02d},{millis:03d}"


def generate_srt(slides: list[SlideSpec], durations: list[float]) -> None:
    lines: list[str] = []
    current = 0.0
    subtitle_index = 1

    for slide, duration in zip(slides, durations):
        text = narration_for_slide(slide).strip()
        chunks: list[str] = []
        current_chunk = ""

        for ch in text:
            current_chunk += ch
            if ch in "，。！？；":
                trimmed = current_chunk.strip()
                if trimmed:
                    chunks.append(trimmed)
                current_chunk = ""

        tail = current_chunk.strip()
        if tail:
            chunks.append(tail)

        cleaned_chunks = [
            chunk
            for chunk in chunks
            if any(ch.isalnum() or "\u4e00" <= ch <= "\u9fff" for ch in chunk)
        ]
        if not cleaned_chunks:
            cleaned_chunks = [slide.narration.strip() or slide.title]

        merged: list[str] = []
        current_line = ""
        for chunk in cleaned_chunks:
            candidate = current_line + chunk
            if len(candidate) <= 22 or not current_line:
                current_line = candidate
            else:
                merged.append(current_line)
                current_line = chunk
        if current_line:
            merged.append(current_line)

        weights = [max(1, len(seg.replace("\n", ""))) for seg in merged]
        total_weight = sum(weights)
        segment_start = current

        for idx, segment in enumerate(merged):
            if idx == len(merged) - 1:
                segment_end = current + duration
            else:
                part = duration * (weights[idx] / total_weight)
                segment_end = segment_start + part

            lines.append(str(subtitle_index))
            lines.append(
                f"{seconds_to_srt(segment_start)} --> {seconds_to_srt(segment_end)}"
            )
            lines.append(segment)
            lines.append("")

            subtitle_index += 1
            segment_start = segment_end

        current += duration

    _ = SRT_PATH.write_text("\n".join(lines), encoding="utf-8")


def build_video(slides: list[SlideSpec], images: list[Path]) -> bool:
    clips: list[Path] = []
    durations: list[float] = []

    for idx, (slide, image_path) in enumerate(zip(slides, images), start=1):
        audio_path = AUDIO_DIR / f"{idx:02d}_{safe_name(slide.title)}.wav"
        clip_path = CLIPS_DIR / f"{idx:02d}_{safe_name(slide.title)}.mp4"

        ok = save_tts_audio(narration_for_slide(slide), audio_path)
        if not ok:
            create_silent_wav(9.0, audio_path)

        duration = max(8.0, wav_duration(audio_path))
        if duration < 8.0:
            create_silent_wav(8.0, audio_path)
            duration = 8.0

        if not build_clip(image_path, audio_path, clip_path):
            return False

        clips.append(clip_path)
        durations.append(duration)

    generate_srt(slides, durations)

    plain_video_path = TEMP_DIR / "knowledge_graph_video_v4_plain.mp4"
    if not concat_clips(clips, plain_video_path):
        return False

    ffmpeg = ffmpeg_path()
    if not ffmpeg:
        return False

    subtitle_path = (
        SRT_PATH.resolve().as_posix().replace(":", "\\:").replace(",", "\\,")
    )
    return run_ffmpeg(
        [
            "-y",
            "-i",
            str(plain_video_path),
            "-vf",
            f"subtitles='{subtitle_path}':force_style='FontName=Microsoft YaHei,FontSize=20,PrimaryColour=&H00FFFFFF,OutlineColour=&H00000000,BorderStyle=3,Outline=1,Shadow=0,MarginV=24'",
            "-c:a",
            "copy",
            str(VIDEO_PATH),
        ]
    )


def build_pptx(images: list[Path], slides: list[SlideSpec]) -> None:
    if Presentation is None or Inches is None or Pt is None:
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for img_path, slide_spec in zip(images, slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(
            Inches(0.45), Inches(0.18), Inches(9.8), Inches(0.5)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = slide_spec.title
        p.font.size = Pt(28)
        p.font.bold = True

        subtitle_box = slide.shapes.add_textbox(
            Inches(0.48), Inches(0.70), Inches(10.5), Inches(0.35)
        )
        p2 = subtitle_box.text_frame.paragraphs[0]
        p2.text = slide_spec.subtitle
        p2.font.size = Pt(15)

        _ = slide.shapes.add_picture(
            str(img_path),
            Inches(0.36),
            Inches(1.12),
            width=Inches(12.55),
            height=Inches(6.0),
        )

    prs.save(str(PPTX_PATH))


def print_summary(
    images: list[Path],
    video_ok: bool,
    crops: dict[str, Path],
    runtime_shots: dict[str, Path],
) -> None:
    print("已生成 v4 视频素材：")
    print(f"- {SCRIPT_PATH}")
    print(f"- {SRT_PATH}")
    print(f"- {PPTX_PATH}")
    for key, value in sorted(crops.items()):
        print(f"- {key}: {value}")
    for key, value in sorted(runtime_shots.items()):
        print(f"- {key}: {value}")
    for img in images:
        print(f"- {img}")
    if video_ok:
        print(f"- {VIDEO_PATH}")
    else:
        print("- 视频未成功生成，请检查 ffmpeg 或音频环境。")


def main() -> None:
    ensure_dirs()

    required = [FRAMEWORK_IMG, CLASS_IMG, SEQUENCE_IMG, PROCESS_IMG]
    missing = [path for path in required if not path.exists()]
    if missing:
        print("缺少以下 UML 图片文件：")
        for path in missing:
            print(f"- {path}")
        raise SystemExit(1)

    crops = generate_uml_crops()
    runtime_shots = generate_mock_runtime_screenshots()
    slides = build_slides(crops, runtime_shots)
    build_script_markdown(slides)
    images = render_all_slides(slides)
    build_pptx(images, slides)
    video_ok = build_video(slides, images)
    print_summary(images, video_ok, crops, runtime_shots)


if __name__ == "__main__":
    main()
