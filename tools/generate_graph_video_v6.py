#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
知识图谱核心功能教学视频生成器 v6

同步机制（参考 MTCS v4.2 实现）：
  1. 旁白按标点自动分句，或由 voice_segments 手动指定
  2. 每句用 edge_tts (zh-CN-XiaoxiaoNeural) 生成独立 MP3，获取真实时长
  3. 为每句渲染专属 PNG（内容相同，仅底部字幕条文字不同）
  4. moviepy: clip = ImageClip(png).set_duration(real_dur + tail).set_audio(AudioFileClip(mp3))
  5. concatenate_videoclips → write_videofile 一次性编码，彻底消除逐片段 AAC 漂移
  6. SRT 时间戳 = 各句 MP3 真实时长累加，与烧录字幕 100% 一致
"""

import asyncio
import os
import sys
import time
import urllib.request
import wave
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

# ── edge_tts ──────────────────────────────────────────────────────────────
try:
    import edge_tts

    _EDGE_OK = True
except ImportError:
    _EDGE_OK = False
    print("Warning: edge-tts not found.  pip install edge-tts")

# ── moviepy ───────────────────────────────────────────────────────────────
try:
    from moviepy.editor import AudioFileClip, ImageClip, concatenate_videoclips

    _MOVIEPY_OK = True
except ImportError:
    _MOVIEPY_OK = False
    print("Error: moviepy not found.  pip install moviepy==1.0.3")

# ── pyttsx3 fallback ──────────────────────────────────────────────────────
try:
    import pyttsx3 as _pyttsx3_mod  # type: ignore[import]

    _PYTTSX3_OK = True
except Exception:
    _pyttsx3_mod = None
    _PYTTSX3_OK = False

# ── python-pptx ───────────────────────────────────────────────────────────
try:
    from pptx import Presentation  # type: ignore[import]
    from pptx.util import Inches, Pt  # type: ignore[import]

    _PPTX_OK = True
except Exception:
    Presentation = Inches = Pt = None
    _PPTX_OK = False


# ═══════════════════════════════════════════════════════════════════════════
# PATHS
# ═══════════════════════════════════════════════════════════════════════════
ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
DIAGRAMS_DIR = DOCS_DIR / "diagrams" / "v3"
TESTING_DIR = DOCS_DIR / "testing"
OUT_DIR = ROOT / "video_output"
V6_DIR = DOCS_DIR / "video" / "v6"

SLIDES_DIR = V6_DIR / "slides"  # base slide PNGs (no subtitle)
SENT_DIR = V6_DIR / "sent"  # per-sentence PNGs (subtitle burned in)
AUDIO_DIR = V6_DIR / "audio"
CROPS_DIR = V6_DIR / "crops"
ASSETS_DIR = V6_DIR / "assets"
TEMP_DIR = V6_DIR / "temp"

VIDEO_PATH = OUT_DIR / "知识图谱核心功能_图谱功能教程_v6.mp4"
PPTX_PATH = OUT_DIR / "知识图谱核心功能_图谱功能教程_v6.pptx"
SRT_PATH = V6_DIR / "subtitles_v6.srt"
SCRIPT_PATH = V6_DIR / "script_v6.md"

FLUTTER_IMG = ASSETS_DIR / "flutter_official.png"
DART_IMG = ASSETS_DIR / "dart_official.png"

FRAMEWORK_IMG = DIAGRAMS_DIR / "flutter_dart_framework_architecture.png"
CLASS_IMG = DIAGRAMS_DIR / "flutter_dart_core_class_diagram.png"
SEQUENCE_IMG = DIAGRAMS_DIR / "graph_feature_sequence_diagram.png"
PROCESS_IMG = DIAGRAMS_DIR / "knowledge_graph_development_process.png"

MODEL_TEST_PATH = ROOT / "test" / "models" / "model_test.dart"
TEST_CASES_PATH = TESTING_DIR / "test_cases.md"
TEST_REPORT_PATH = TESTING_DIR / "test_report.md"
LOGIN_GOLDEN = ROOT / "test" / "screenshots" / "goldens" / "login_page.png"
HOME_GOLDEN = ROOT / "test" / "screenshots" / "goldens" / "home_page.png"


# ═══════════════════════════════════════════════════════════════════════════
# VIDEO PARAMETERS
# ═══════════════════════════════════════════════════════════════════════════
W, H = 1920, 1080
FPS = 30
TTS_VOICE = "zh-CN-XiaoxiaoNeural"
TTS_RATE = "-5%"
CLIP_TAIL = 0.45  # seconds held after speech for mid-slide sentences
SLIDE_TAIL = 1.2  # extra tail for last sentence of each slide
SUB_H = 115  # subtitle bar height (px)
CONTENT_H = H - SUB_H  # 965 px – main content area
TITLE_H = 84  # title bar height


# ═══════════════════════════════════════════════════════════════════════════
# COLOURS
# ═══════════════════════════════════════════════════════════════════════════
BG_TOP = (20, 36, 68)
BG_BOT = (32, 54, 96)
TITLE_BG = (13, 24, 50)
ACCENT = (255, 198, 55)
TEXT = (236, 241, 252)
TEXT_MUTE = (165, 185, 220)
PRIMARY = (92, 172, 255)
GREEN = (72, 208, 126)
ORANGE = (255, 152, 55)
PURPLE = (168, 110, 255)
BORDER = (52, 84, 140)
CODE_BG = (24, 30, 44)
PANEL = (26, 44, 76)
SUB_BG = (6, 6, 12)
WHITE = (255, 255, 255)


# ═══════════════════════════════════════════════════════════════════════════
# FONTS
# ═══════════════════════════════════════════════════════════════════════════
def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    paths = (
        [r"C:\Windows\Fonts\msyhbd.ttc", r"C:\Windows\Fonts\simhei.ttf"]
        if bold
        else [r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\simsun.ttc"]
    ) + [r"C:\Windows\Fonts\arial.ttf"]
    for p in paths:
        try:
            return ImageFont.truetype(p, size)
        except Exception:
            continue
    return ImageFont.load_default()


FT = _font(48, True)  # slide title
FS = _font(26, False)  # slide subtitle
FB = _font(24, True)  # section label
FN = _font(22, False)  # body / bullets
FK = _font(19, False)  # code / small
FKB = _font(19, True)  # code bold / badge
FSU = _font(27, False)  # subtitle bar text


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE SPEC
# ═══════════════════════════════════════════════════════════════════════════
@dataclass
class SlideSpec:
    title: str
    subtitle: str
    bullets: list[str]
    narration: str
    image_path: Path | None = None
    image_caption: str | None = None
    code_title: str | None = None
    code_text: str | None = None
    voice_segments: list[str] | None = None


# ═══════════════════════════════════════════════════════════════════════════
# DIRECTORY SETUP + ASSET DOWNLOAD
# ═══════════════════════════════════════════════════════════════════════════
def ensure_dirs() -> None:
    for d in [
        OUT_DIR,
        V6_DIR,
        SLIDES_DIR,
        SENT_DIR,
        AUDIO_DIR,
        CROPS_DIR,
        ASSETS_DIR,
        TEMP_DIR,
    ]:
        d.mkdir(parents=True, exist_ok=True)


def _download(url: str, out: Path) -> Path | None:
    if out.exists() and out.stat().st_size > 500:
        return out
    try:
        urllib.request.urlretrieve(url, str(out))
        return out if out.exists() and out.stat().st_size > 500 else None
    except Exception as e:
        print(f"  Download failed ({out.name}): {e}")
        return None


def ensure_official_assets() -> dict[str, Path]:
    assets: dict[str, Path] = {}
    f = _download(
        "https://flutter.dev/assets/lockup_flutter_horizontal.d0515092173211776ceed19b39c2a041.png",
        FLUTTER_IMG,
    )
    if f:
        assets["flutter"] = f
    d = _download(
        "https://dart.dev/assets/img/logo/dart-logo-for-shares.png",
        DART_IMG,
    )
    if d:
        assets["dart"] = d
    return assets


# ═══════════════════════════════════════════════════════════════════════════
# DRAWING UTILITIES
# ═══════════════════════════════════════════════════════════════════════════
def _tsz(draw: ImageDraw.ImageDraw, text: str, font) -> tuple[int, int]:
    bb = draw.textbbox((0, 0), text, font=font)
    return int(bb[2] - bb[0]), int(bb[3] - bb[1])


def _wrap(draw: ImageDraw.ImageDraw, text: str, font, max_w: int) -> list[str]:
    lines: list[str] = []
    for para in text.splitlines():
        cur = ""
        for ch in para.strip():
            cand = cur + ch
            if _tsz(draw, cand, font)[0] > max_w and cur:
                lines.append(cur)
                cur = ch
            else:
                cur = cand
        if cur:
            lines.append(cur)
    return lines


def _draw_wrapped(
    draw: ImageDraw.ImageDraw,
    text: str,
    x: int,
    y: int,
    font,
    fill,
    max_w: int,
    gap: int = 5,
) -> int:
    for line in _wrap(draw, text, font, max_w):
        draw.text((x, y), line, font=font, fill=fill)
        y += _tsz(draw, line or "M", font)[1] + gap
    return y


def _gradient(img: Image.Image) -> None:
    draw = ImageDraw.Draw(img)
    for y in range(img.height):
        r = BG_TOP[0] + int((BG_BOT[0] - BG_TOP[0]) * y / img.height)
        g = BG_TOP[1] + int((BG_BOT[1] - BG_TOP[1]) * y / img.height)
        b = BG_TOP[2] + int((BG_BOT[2] - BG_TOP[2]) * y / img.height)
        draw.line([(0, y), (img.width, y)], fill=(r, g, b))


def _paste_fit(
    base: Image.Image, src: Path | None, box: tuple[int, int, int, int]
) -> None:
    if not src or not src.exists():
        return
    try:
        im = Image.open(src).convert("RGB")
    except Exception:
        return
    bw, bh = box[2] - box[0], box[3] - box[1]
    ratio = min(bw / im.width, bh / im.height)
    nw, nh = max(1, int(im.width * ratio)), max(1, int(im.height * ratio))
    im = im.resize((nw, nh), Image.LANCZOS)
    base.paste(im, (box[0] + (bw - nw) // 2, box[1] + (bh - nh) // 2))


# ═══════════════════════════════════════════════════════════════════════════
# TEXT FILE UTILITIES
# ═══════════════════════════════════════════════════════════════════════════
def _read(path: Path) -> str:
    return path.read_text(encoding="utf-8") if path.exists() else ""


def _extract(path: Path, marker: str, n: int = 20) -> str:
    text = _read(path)
    for i, line in enumerate(text.splitlines()):
        if marker in line:
            return "\n".join(text.splitlines()[i : i + n])
    return "\n".join(text.splitlines()[:n])


def _select_lines(text: str, keywords: list[str], n: int = 10) -> str:
    picked: list[str] = []
    for line in text.splitlines():
        s = line.strip()
        if s and any(k.lower() in s.lower() for k in keywords):
            picked.append(s)
            if len(picked) >= n:
                break
    return "\n".join(picked) if picked else "\n".join(text.splitlines()[:n])


def load_code_summaries() -> tuple[str, str]:
    model = _extract(MODEL_TEST_PATH, "group('GraphModel'", 20)
    cases = _select_lines(
        _read(TEST_CASES_PATH),
        ["tc-", "图谱", "登录", "测验", "资料", "视频", "model"],
        10,
    )
    return model, cases


def load_test_summary() -> str:
    return _select_lines(
        _read(TEST_REPORT_PATH),
        ["23", "pass", "通过", "apk", "debug", "release", "测试", "构建"],
        10,
    )


def _safe(text: str) -> str:
    keep = []
    for ch in text:
        if ch.isalnum() or ch in ("-", "_") or "\u4e00" <= ch <= "\u9fff":
            keep.append(ch)
        else:
            keep.append("_")
    return "".join(keep).strip("_") or "slide"


# ═══════════════════════════════════════════════════════════════════════════
# UML CROPS
# ═══════════════════════════════════════════════════════════════════════════
def _crop_save(
    src: Path,
    box: tuple[float, float, float, float],
    name: str,
    border: tuple[int, int, int] = PRIMARY,
) -> Path:
    img = Image.open(src).convert("RGB")
    iw, ih = img.size
    cropped = img.crop(
        (int(iw * box[0]), int(ih * box[1]), int(iw * box[2]), int(ih * box[3]))
    )
    canvas = Image.new("RGB", (1600, 900), BG_TOP)
    draw = ImageDraw.Draw(canvas)
    tw, th = 1560, 860
    ratio = min(tw / cropped.width, th / cropped.height)
    nw = max(1, int(cropped.width * ratio))
    nh = max(1, int(cropped.height * ratio))
    resized = cropped.resize((nw, nh), Image.LANCZOS)
    canvas.paste(resized, (20 + (tw - nw) // 2, 20 + (th - nh) // 2))
    draw.rounded_rectangle((10, 10, 1590, 890), radius=16, outline=border, width=4)
    out = CROPS_DIR / name
    canvas.save(out)
    return out


def generate_uml_crops() -> dict[str, Path]:
    crops: dict[str, Path] = {}
    if FRAMEWORK_IMG.exists():
        crops["framework_full"] = FRAMEWORK_IMG
        crops["framework_ui"] = _crop_save(
            FRAMEWORK_IMG, (0.01, 0.01, 0.40, 0.55), "fw_ui.png", PRIMARY
        )
        crops["framework_dao"] = _crop_save(
            FRAMEWORK_IMG, (0.26, 0.30, 0.64, 0.78), "fw_dao.png", GREEN
        )
        crops["framework_right"] = _crop_save(
            FRAMEWORK_IMG, (0.58, 0.20, 0.95, 0.66), "fw_right.png", ORANGE
        )
    if CLASS_IMG.exists():
        crops["class_full"] = CLASS_IMG
        crops["class_ui"] = _crop_save(
            CLASS_IMG, (0.01, 0.02, 0.42, 0.58), "cls_ui.png", PRIMARY
        )
        crops["class_dao"] = _crop_save(
            CLASS_IMG, (0.44, 0.22, 0.76, 0.64), "cls_dao.png", GREEN
        )
    if SEQUENCE_IMG.exists():
        crops["sequence_full"] = SEQUENCE_IMG
        crops["sequence_enter"] = _crop_save(
            SEQUENCE_IMG, (0.04, 0.00, 0.96, 0.36), "seq_enter.png", PRIMARY
        )
        crops["sequence_detail"] = _crop_save(
            SEQUENCE_IMG, (0.04, 0.24, 0.96, 0.62), "seq_detail.png", GREEN
        )
        crops["sequence_action"] = _crop_save(
            SEQUENCE_IMG, (0.04, 0.56, 0.96, 0.92), "seq_action.png", ORANGE
        )
    if PROCESS_IMG.exists():
        crops["process_full"] = PROCESS_IMG
        crops["process_early"] = _crop_save(
            PROCESS_IMG, (0.02, 0.06, 0.34, 0.46), "proc_early.png", PRIMARY
        )
        crops["process_mid"] = _crop_save(
            PROCESS_IMG, (0.30, 0.28, 0.68, 0.68), "proc_mid.png", GREEN
        )
        crops["process_late"] = _crop_save(
            PROCESS_IMG, (0.60, 0.50, 0.97, 0.92), "proc_late.png", PURPLE
        )
    return crops


# ═══════════════════════════════════════════════════════════════════════════
# MOCK RUNTIME SCREENSHOTS
# ═══════════════════════════════════════════════════════════════════════════
def _mock_page(
    title: str,
    subtitle: str,
    tags: list[str],
    cards: list[tuple[str, str, tuple[int, int, int]]],
    filename: str,
) -> Path:
    img = Image.new("RGB", (1440, 900), BG_TOP)
    _gradient(img)
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 1440, 88), fill=TITLE_BG)
    draw.line((0, 88, 1440, 88), fill=ACCENT, width=3)
    draw.text((40, 14), title, font=_font(38, True), fill=ACCENT)
    draw.text((42, 58), subtitle, font=_font(17), fill=TEXT_MUTE)
    x = 40
    for tag in tags:
        tw, _ = _tsz(draw, tag, _font(17, True))
        draw.rounded_rectangle(
            (x, 104, x + tw + 26, 134),
            radius=10,
            fill=(36, 58, 102),
            outline=BORDER,
            width=1,
        )
        draw.text((x + 13, 107), tag, font=_font(17, True), fill=PRIMARY)
        x += tw + 40
    y = 158
    for ctitle, cdesc, color in cards:
        draw.rounded_rectangle(
            (40, y, 1400, y + 128),
            radius=14,
            fill=(28, 46, 78),
            outline=BORDER,
            width=2,
        )
        draw.rounded_rectangle(
            (58, y + 18, 86, y + 46), radius=8, fill=color, outline=color, width=1
        )
        draw.text((98, y + 12), ctitle, font=_font(24, True), fill=TEXT)
        _draw_wrapped(draw, cdesc, 98, y + 48, _font(19), TEXT_MUTE, 1270, 4)
        y += 148
    out = CROPS_DIR / filename
    img.save(out)
    return out


def generate_runtime_shots() -> dict[str, Path]:
    shots: dict[str, Path] = {}
    if LOGIN_GOLDEN.exists():
        shots["login"] = LOGIN_GOLDEN
    else:
        shots["login"] = _mock_page(
            "LoginPage 登录页",
            "统一入口 · 表单校验 · 跳转首页",
            ["统一入口", "表单输入", "登录校验"],
            [
                ("用户名输入", "用户输入用户名与密码，提交登录请求", PRIMARY),
                ("状态反馈", "表单校验失败时给出错误提示", ORANGE),
                ("进入首页", "登录成功后跳转至 HomePage", GREEN),
            ],
            "mock_login.png",
        )
    if HOME_GOLDEN.exists():
        shots["home"] = HOME_GOLDEN
    else:
        shots["home"] = _mock_page(
            "HomePage 首页",
            "主导航 · 学习入口 · 功能分发",
            ["主导航", "课程入口", "学习记录"],
            [
                (
                    "底部导航",
                    "首页 / 图谱 / 资料 / 视频 / 测验 / 进度 六大模块",
                    PRIMARY,
                ),
                ("功能入口", "把分散模块统一组织成清晰学习入口", GREEN),
                ("学习闭环", "通过进度与测验形成反馈机制", ORANGE),
            ],
            "mock_home.png",
        )
    shots["graph_list"] = _mock_page(
        "GraphListPage 图谱列表",
        "图谱主题选择 · 列表展示 · 进入详情",
        ["图谱入口", "主题列表", "数据加载"],
        [
            ("课程图谱", "展示课程整体知识结构与章节关系", PRIMARY),
            ("技术栈图谱", "展示 Flutter / Dart / SQLite 技术关系", GREEN),
            ("学习图谱", "展示学习路径、节点与实践内容", ORANGE),
        ],
        "mock_graph_list.png",
    )
    shots["graph_detail"] = _mock_page(
        "GraphDetailPage 图谱详情",
        "节点渲染 · 展开浏览 · 学习记录 · 收藏",
        ["节点交互", "缩放拖拽", "学习记录", "收藏"],
        [
            (
                "图谱可视化区域",
                "CustomPainter + InteractiveViewer 浏览图谱结构",
                PRIMARY,
            ),
            ("节点详情卡片", "点击节点后显示标题、类型、内容和操作按钮", GREEN),
            ("开始学习 / 收藏", "把浏览行为转化为学习记录和收藏数据", ORANGE),
        ],
        "mock_graph_detail.png",
    )
    return shots


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE CONTENT
# ═══════════════════════════════════════════════════════════════════════════
def build_slides(
    crops: dict[str, Path],
    shots: dict[str, Path],
    official: dict[str, Path],
) -> list[SlideSpec]:
    model_code, test_cases = load_code_summaries()
    test_report = load_test_summary()

    return [
        # 01 ─ Course intro ───────────────────────────────────────────────
        SlideSpec(
            title="课程导入",
            subtitle="先了解学习目标与整体观看思路",
            bullets=[
                "系统整体结构",
                "Flutter 与 Dart 技术栈",
                "知识图谱核心流程",
                "测试与交付成果",
            ],
            narration=(
                "欢迎进入知识图谱核心功能教学视频。"
                "本节内容围绕系统结构、Flutter 与 Dart 实战方式、知识图谱核心流程以及测试结果展开。"
                "建议先建立整体认识，再逐步深入页面、数据和交互细节。"
            ),
            image_path=crops.get("framework_full"),
        ),
        # 02 ─ Project overview ───────────────────────────────────────────
        SlideSpec(
            title="项目总览",
            subtitle="认识知识图谱学习系统的整体结构",
            bullets=[
                "Flutter 页面层",
                "Dart 业务层",
                "SQLite 数据层",
                "知识图谱驱动学习",
                "测试与交付材料",
            ],
            narration=(
                "这是整个项目的总体视图。"
                "系统以 Flutter 负责页面表现，以 Dart 负责业务逻辑与数据组织，以 SQLite 负责本地持久化。"
                "知识图谱位于系统核心，用来连接页面导航、节点学习、收藏行为和学习进度。"
            ),
            image_path=crops.get("framework_full"),
        ),
        # 03 ─ Flutter official intro ─────────────────────────────────────
        SlideSpec(
            title="Flutter 技术介绍",
            subtitle="Flutter 官方定位与本项目最佳实战",
            bullets=[
                "开源框架 — 单一代码库多平台",
                "Fast · Productive · Flexible",
                "声明式 UI，Widget 驱动渲染",
                "Hot Reload 快速迭代",
                "实战：页面层只负责渲染与交互",
            ],
            narration=(
                "Flutter 官方将自己定义为一个开源框架，用于通过单一代码库构建精美的原生编译多平台应用。"
                "官方强调三个核心特点。"
                "Fast —— 代码编译为 ARM 或 Intel 机器码，性能接近原生。"
                "Productive —— 支持 Hot Reload，修改代码几乎立即看到效果。"
                "Flexible —— 像素级控制，适配任何屏幕与平台。"
                "放到本项目中，最佳实战是让 Flutter 页面只关心渲染与交互，把数据访问和业务逻辑交给独立的 Dart 层。"
            ),
            voice_segments=[
                "Flutter 官方将自己定义为一个开源框架，",
                "用于通过单一代码库构建精美的原生编译多平台应用。",
                "官方强调三个核心特点。",
                "Fast —— 代码编译为 ARM 或 Intel 机器码，性能接近原生。",
                "Productive —— 支持 Hot Reload，修改代码几乎立即看到效果。",
                "Flexible —— 像素级控制，适配任何屏幕与平台。",
                "放到本项目中，最佳实战是让 Flutter 页面只关心渲染与交互，",
                "把数据访问和业务逻辑交给独立的 Dart 层。",
            ],
            image_path=official.get("flutter", shots.get("home")),
            image_caption="Flutter 官方品牌图 — flutter.dev",
        ),
        # 04 ─ Dart official intro ────────────────────────────────────────
        SlideSpec(
            title="Dart 技术介绍",
            subtitle="Dart 官方定位与本项目最佳实战",
            bullets=[
                "client-optimized language",
                "fast apps on any platform",
                "类型安全 + sound null safety",
                "AOT / JIT 双编译模式",
                "实战：Model / DAO / Service 分层",
            ],
            narration=(
                "Dart 官方将自己定义为一种面向客户端优化的语言，用于在任何平台上开发快速应用。"
                "它具备类型安全和 sound null safety，能在编译期消除空指针异常。"
                "支持 AOT 编译到机器码实现高性能启动，也支持 JIT 热重载提升开发效率。"
                "放到本项目中，最佳实战是用 Dart 承载模型对象、DAO 层和业务逻辑，让页面直接消费结构化结果，而不是拼接零散数据。"
            ),
            voice_segments=[
                "Dart 官方将自己定义为一种面向客户端优化的语言，",
                "用于在任何平台上开发快速应用。",
                "它具备类型安全和 sound null safety，",
                "能在编译期消除空指针异常。",
                "支持 AOT 编译到机器码实现高性能启动，",
                "也支持 JIT 热重载提升开发效率。",
                "放到本项目中，最佳实战是用 Dart 承载模型对象、DAO 层和业务逻辑，",
                "让页面直接消费结构化结果，而不是拼接零散数据。",
            ],
            image_path=official.get("dart", crops.get("class_dao")),
            image_caption="Dart 官方品牌图 — dart.dev",
        ),
        # 05 ─ Page roles ─────────────────────────────────────────────────
        SlideSpec(
            title="页面职责分析",
            subtitle="从页面分工理解系统学习闭环",
            bullets=[
                "LoginPage — 统一登录入口",
                "HomePage — 导航分发中心",
                "GraphListPage — 图谱主题选择",
                "GraphDetailPage — 核心交互页面",
                "资料 / 视频 / 测验 / 进度联动",
            ],
            narration=(
                "系统中的页面不是简单罗列，而是围绕学习闭环进行分工。"
                "登录页负责统一入口，首页负责模块导航，图谱列表页负责主题选择，图谱详情页负责节点浏览、开始学习和收藏。"
                "资料、视频、测验和进度页面则继续承接学习行为，形成完整的使用路径。"
            ),
            image_path=shots.get("home"),
        ),
        # 06 ─ System framework ───────────────────────────────────────────
        SlideSpec(
            title="系统框架总览",
            subtitle="三层架构：页面层 / 业务层 / 数据层",
            bullets=[
                "页面层负责展示与交互",
                "业务层负责规则与查询",
                "数据层负责 SQLite 读写",
                "知识图谱作为系统纽带",
            ],
            narration=(
                "从框架图可以看出，系统分为页面层、业务逻辑层和数据层。"
                "页面层承接用户交互，业务层负责认证、图谱查询、收藏管理和学习记录，数据层负责 SQLite 初始化与读写。"
                "知识图谱位于中间纽带位置，向上驱动页面，向下连接数据。"
            ),
            image_path=crops.get("framework_full"),
        ),
        # 07 ─ UI layer ───────────────────────────────────────────────────
        SlideSpec(
            title="UI 层讲解",
            subtitle="把页面名称与页面职责一一对应",
            bullets=[
                "LoginPage — 登录",
                "HomePage — 主导航",
                "GraphListPage — 图谱入口",
                "GraphDetailPage — 核心交互",
                "Document / Video / Quiz / Progress 联动",
            ],
            narration=(
                "这里聚焦 UI 层。"
                "理解 UI 层的关键不是记页面名称，而是把每个页面放回用户学习流程中。"
                "用户先登录，再通过首页进入图谱模块，在图谱详情里完成浏览、学习和收藏，最后联动其他学习资源。"
            ),
            image_path=crops.get("framework_ui"),
        ),
        # 08 ─ DAO layer ──────────────────────────────────────────────────
        SlideSpec(
            title="DAO 与数据层讲解",
            subtitle="理解页面背后的业务逻辑与数据库访问",
            bullets=[
                "AuthService — 管理登录状态",
                "GraphDao — 图谱 / 节点 / 边查询",
                "FavoriteDao — 收藏管理",
                "LearningRecordDao — 学习记录",
                "DatabaseHelper — 统一管理 SQLite",
            ],
            narration=(
                "页面背后的核心能力来自 Dart 业务层和数据层。"
                "AuthService 负责登录状态，GraphDao 负责图谱、节点和边的查询，FavoriteDao 与 LearningRecordDao 负责行为数据写入。"
                "所有这些能力最终都建立在 DatabaseHelper 对 SQLite 的统一管理之上。"
            ),
            image_path=crops.get("framework_dao"),
        ),
        # 09 ─ Core class diagram ─────────────────────────────────────────
        SlideSpec(
            title="核心类图总览",
            subtitle="从类图理解页面层、DAO 层和模型层",
            bullets=[
                "页面类 — 交互",
                "DAO 类 — 数据访问",
                "Model 类 — 结构化对象",
                "GraphDetailPage 是核心类",
            ],
            narration=(
                "类图体现了系统的职责边界。"
                "页面类负责展示与交互，DAO 类负责查询和写入，模型类负责把数据库结果组织成清晰对象。"
                "其中 GraphDetailPage 是知识图谱功能的核心页面，也是用户操作最集中的位置。"
            ),
            image_path=crops.get("class_full"),
        ),
        # 10 ─ UI class detail ────────────────────────────────────────────
        SlideSpec(
            title="UI 类职责讲解",
            subtitle="区分应用启动、导航分发与图谱交互类",
            bullets=[
                "MyApp — 应用启动",
                "HomePage — 主导航",
                "GraphListPage — 图谱入口",
                "GraphDetailPage — 图谱交互",
                "Progress / Favorites — 反馈与复习",
            ],
            narration=(
                "放大 UI 类之后，可以更清楚地看到页面类的分工。"
                "MyApp 负责应用入口，HomePage 负责导航分发，GraphListPage 负责图谱选择，GraphDetailPage 负责节点渲染与操作。"
                "其他页面则围绕学习结果和复习入口继续延展。"
            ),
            image_path=crops.get("class_ui"),
        ),
        # 11 ─ DAO model detail ───────────────────────────────────────────
        SlideSpec(
            title="DAO 与模型层讲解",
            subtitle="理解数据查询、封装与对象协同",
            bullets=[
                "GraphDao — 图谱查询",
                "FavoriteDao — 收藏",
                "LearningRecordDao — 记录",
                "GraphModel / NodeModel / EdgeModel 封装结构",
            ],
            narration=(
                "DAO 层和模型层共同决定了数据怎样被页面消费。"
                "DAO 负责向数据库取数，模型负责把原始结果整理成页面直接可用的对象结构。"
                "当这两层职责清晰时，页面就可以专注于显示，不需要承担复杂的数据拼装工作。"
            ),
            image_path=crops.get("class_dao"),
            code_title="模型代码片段",
            code_text=model_code,
        ),
        # 12 ─ Core flow overview ─────────────────────────────────────────
        SlideSpec(
            title="知识图谱核心流程总览",
            subtitle="六个环节：数据加载 → 列表显示 → 节点边生成 → 展开浏览 → 学习收藏 → 联动资源",
            bullets=[
                "① 数据库加载图谱基础数据",
                "② 图谱列表页显示主题入口",
                "③ 详情页读取节点与边",
                "④ 计算布局、展开浏览",
                "⑤ 点击节点 — 开始学习 / 收藏",
                "⑥ 写入记录 — 联动资源进度",
            ],
            narration=(
                "知识图谱的核心流程可以分成六个连续环节。"
                "首先，系统从数据库加载图谱基础数据。"
                "然后，图谱列表页把这些数据整理成主题入口，供用户选择。"
                "进入详情页后，系统继续读取节点和边，生成图谱结构对象。"
                "完成布局计算后，用户可以展开、缩放和拖拽浏览整个图谱。"
                "当用户点击节点时，可以查看详情，并执行开始学习或收藏操作。"
                "最后，这些行为被写入数据库，并继续联动资料、视频、测验和进度页面。"
            ),
            voice_segments=[
                "知识图谱的核心流程可以分成六个连续环节。",
                "首先，系统从数据库加载图谱基础数据。",
                "然后，图谱列表页把这些数据整理成主题入口，供用户选择。",
                "进入详情页后，系统继续读取节点和边，生成图谱结构对象。",
                "完成布局计算后，用户可以展开、缩放和拖拽浏览整个图谱。",
                "当用户点击节点时，可以查看详情，并执行开始学习或收藏操作。",
                "最后，这些行为被写入数据库，并联动资料、视频、测验和进度页面。",
            ],
            image_path=crops.get("sequence_full"),
        ),
        # 13 ─ Graph list loading ─────────────────────────────────────────
        SlideSpec(
            title="图谱数据加载与列表显示",
            subtitle="GraphDao 查询 → GraphListPage 展示 → 用户选择图谱",
            bullets=[
                "GraphDao.getAllGraphs() 查询数据库",
                "返回 List<GraphModel>",
                "GraphListPage 渲染图谱卡片列表",
                "用户点击卡片 → 进入详情页",
            ],
            narration=(
                "图谱功能的入口从数据加载开始。"
                "GraphDao 先查询图谱基础数据，GraphListPage 再把这些结果组织成可点击的图谱卡片列表。"
                "用户看到的是课程图谱、技术栈图谱等主题入口。"
                "点击其中一个卡片之后，系统才会进入图谱详情流程。"
            ),
            image_path=shots.get("graph_list"),
        ),
        # 14 ─ Node and edge generation ───────────────────────────────────
        SlideSpec(
            title="节点与边生成流程",
            subtitle="详情页把原始数据转换成可视化图谱结构",
            bullets=[
                "根据 graphId 读取 nodes",
                "根据 graphId 读取 edges",
                "生成 NodeModel / EdgeModel",
                "计算布局坐标 → CustomPainter 渲染",
            ],
            narration=(
                "进入图谱详情页后，系统根据当前图谱标识继续读取节点和边。"
                "这些原始数据被转换成 NodeModel 和 EdgeModel，再交给页面进行布局计算。"
                "布局完成之后，CustomPainter 将节点和边绘制成可视化图谱。"
                "完成这一层处理之后，用户看到的才是可视化图谱，而不是数据库里的原始记录。"
            ),
            image_path=crops.get("sequence_detail"),
        ),
        # 15 ─ Expand browse and favourite ────────────────────────────────
        SlideSpec(
            title="图谱展开与收藏流程",
            subtitle="节点浏览、展开、学习与收藏都在详情页完成",
            bullets=[
                "InteractiveViewer — 缩放与拖拽",
                "点击节点 → 弹出详情卡片",
                "开始学习 → 写入 learning_records",
                "收藏 → 写入 favorites",
            ],
            narration=(
                "图谱详情页不仅负责显示，还负责交互。"
                "用户通过 InteractiveViewer 可以缩放和拖拽浏览整个图谱。"
                "点击某个节点后，系统会弹出节点详情卡片，显示标题、类型和内容。"
                "如果用户点击开始学习，系统就把学习行为写入 learning_records。"
                "如果用户点击收藏，系统就把当前节点信息写入 favorites。"
            ),
            image_path=crops.get("sequence_action"),
        ),
        # 16 ─ Development early ──────────────────────────────────────────
        SlideSpec(
            title="开发过程：前期阶段",
            subtitle="需求分析、业务建模与数据库设计",
            bullets=[
                "明确图谱是系统核心",
                "定义 Graph / Node / Edge / Favorite / LearningRecord",
                "设计 SQLite 表结构",
            ],
            narration=(
                "开发前期的重点是需求分析和业务建模。"
                "首先明确知识图谱是系统核心，然后定义图谱、节点、边、收藏和学习记录等对象。"
                "最后再把这些对象映射成 SQLite 表结构，为后续编码打下基础。"
            ),
            image_path=crops.get("process_early"),
        ),
        # 17 ─ Development mid ────────────────────────────────────────────
        SlideSpec(
            title="开发过程：中期阶段",
            subtitle="DAO 开发、页面实现与图谱交互构建",
            bullets=[
                "实现 DatabaseHelper",
                "实现 GraphDao / FavoriteDao / LearningRecordDao",
                "实现 GraphListPage / GraphDetailPage",
                "构建图谱布局和节点交互",
            ],
            narration=(
                "开发中期进入真正的实现阶段。"
                "这一阶段先完成数据库和 DAO，再实现图谱列表页和图谱详情页。"
                "当数据层、页面层和交互层被连起来之后，图谱功能才真正具备使用价值。"
            ),
            image_path=crops.get("process_mid"),
        ),
        # 18 ─ Development late ───────────────────────────────────────────
        SlideSpec(
            title="开发过程：后期阶段",
            subtitle="测试验证、构建发布与教学材料交付",
            bullets=[
                "补充模型测试与组件测试",
                "生成测试用例与测试报告",
                "构建 Debug / Release APK",
                "生成 UML / PPT / 视频材料",
            ],
            narration=(
                "开发后期的关键是验证与交付。"
                "项目补充了模型测试和组件测试，整理了测试用例与测试报告，并完成 Debug 和 Release APK 构建。"
                "在此基础上，再输出 UML 图、PPT 和教学视频，形成完整交付材料。"
            ),
            image_path=crops.get("process_late"),
        ),
        # 19 ─ Test assets ────────────────────────────────────────────────
        SlideSpec(
            title="测试资产讲解",
            subtitle="模型测试 + 组件测试 + 测试用例文档",
            bullets=[
                "model_test.dart",
                "widget_test.dart",
                "home_page_widget_test.dart",
                "test_cases.md",
                "23 项测试通过",
            ],
            narration=(
                "项目不仅有页面和结构图，也有实际测试资产。"
                "模型测试覆盖 GraphModel、NodeModel、EdgeModel 等核心数据对象，组件测试覆盖登录页和首页导航。"
                "测试用例文档进一步补充了图谱、资料、视频、测验和进度等模块的验证点。"
            ),
            image_path=shots.get("login"),
            code_title="测试用例摘要",
            code_text=test_cases,
        ),
        # 20 ─ Test results ───────────────────────────────────────────────
        SlideSpec(
            title="测试结果与交付成果",
            subtitle="测试通过 · APK 构建 · 展示材料输出",
            bullets=[
                "23 项测试通过",
                "Debug APK 构建成功",
                "Release APK 构建成功",
                "测试报告已整理",
                "视频 / PPT / UML 已生成",
            ],
            narration=(
                "测试结果说明当前工程已经具备较完整的交付状态。"
                "目前模型层测试和组件测试累计通过二十三项，同时 Debug 和 Release APK 也已经构建成功。"
                "这使得项目不仅能运行，也能以测试报告、PPT 和视频的形式稳定展示。"
            ),
            image_path=shots.get("graph_list"),
            code_title="测试报告摘要",
            code_text=test_report,
        ),
        # 21 ─ Runtime UI ─────────────────────────────────────────────────
        SlideSpec(
            title="实际运行界面讲解",
            subtitle="把结构图与真实页面效果对应起来",
            bullets=[
                "登录页 — 统一入口",
                "首页 — 导航中心",
                "图谱列表页 — 主题选择",
                "图谱详情页 — 节点交互与学习动作",
            ],
            narration=(
                "最后回到运行界面。"
                "登录页负责统一入口，首页负责导航分发，图谱列表页负责主题选择，图谱详情页负责节点交互和学习动作。"
                "把结构图、类图、顺序图和真实页面放在一起理解，系统的整体逻辑就会更清楚。"
            ),
            image_path=shots.get("graph_detail"),
        ),
        # 22 ─ Review ─────────────────────────────────────────────────────
        SlideSpec(
            title="本节回顾",
            subtitle="认识 Flutter / Dart 分工 · 理解三层结构 · 看清图谱流程",
            bullets=[
                "认识 Flutter 与 Dart 分工",
                "理解三层结构与核心类",
                "看清知识图谱核心流程",
                "完成测试与交付闭环",
            ],
            narration=(
                "这一节的重点有四个。"
                "第一，理解 Flutter 负责页面，Dart 负责逻辑。"
                "第二，理解页面层、业务层和数据层的结构关系。"
                "第三，看清知识图谱从数据加载到交互记录的完整流程。"
                "第四，看到测试结果和交付材料如何支撑项目展示。"
            ),
            voice_segments=[
                "这一节的重点有四个。",
                "第一，理解 Flutter 负责页面，Dart 负责逻辑。",
                "第二，理解页面层、业务层和数据层的结构关系。",
                "第三，看清知识图谱从数据加载到交互记录的完整流程。",
                "第四，看到测试结果和交付材料如何支撑项目展示。",
            ],
            image_path=crops.get("framework_right"),
        ),
        # 23 ─ Conclusion ─────────────────────────────────────────────────
        SlideSpec(
            title="课程总结与下节预告",
            subtitle="本节回顾 · 后续学习方向说明",
            bullets=[
                "先回看系统框架与类图",
                "再回看知识图谱核心流程",
                "继续：数据库初始化源码实现",
                "继续：节点布局与交互源码阅读",
            ],
            narration=(
                "本节内容到这里结束。"
                "建议你先回看系统框架、类图和知识图谱核心流程，重新整理整体认识。"
                "下一步可以继续学习数据库初始化、图谱数据装载、节点布局计算以及交互源码实现，从展示理解逐步进入实现理解。"
            ),
            image_path=crops.get("framework_right"),
        ),
    ]


# ═══════════════════════════════════════════════════════════════════════════
# SENTENCE SPLITTING
# ═══════════════════════════════════════════════════════════════════════════
def get_sentences(slide: SlideSpec) -> list[str]:
    if slide.voice_segments:
        return [s.strip() for s in slide.voice_segments if s.strip()]
    text = slide.narration.strip()
    chunks: list[str] = []
    cur = ""
    for ch in text:
        cur += ch
        if ch in "。！？；":
            t = cur.strip()
            if t:
                chunks.append(t)
            cur = ""
    if cur.strip():
        chunks.append(cur.strip())
    # merge very short fragments (< 8 chars) with next
    merged: list[str] = []
    i = 0
    while i < len(chunks):
        if i + 1 < len(chunks) and len(chunks[i]) < 8:
            merged.append(chunks[i] + chunks[i + 1])
            i += 2
        else:
            merged.append(chunks[i])
            i += 1
    return merged if merged else [text]


# ═══════════════════════════════════════════════════════════════════════════
# TTS ENGINE  (edge_tts primary → pyttsx3 fallback)
# ═══════════════════════════════════════════════════════════════════════════
async def _edge_async(text: str, path: Path) -> None:
    comm = edge_tts.Communicate(text, TTS_VOICE, rate=TTS_RATE)
    await comm.save(str(path))


_pyttsx3_engine = None


def _init_pyttsx3() -> None:
    global _pyttsx3_engine
    if _pyttsx3_engine is not None or not _PYTTSX3_OK:
        return
    try:
        eng = _pyttsx3_mod.init()
        eng.setProperty("rate", 155)
        voices = eng.getProperty("voices") or []
        for v in voices:
            combined = f"{getattr(v, 'name', '')} {getattr(v, 'id', '')}".lower()
            if any(k in combined for k in ["zh", "chinese", "huihui", "xiaoxiao"]):
                eng.setProperty("voice", v.id)
                break
        _pyttsx3_engine = eng
    except Exception as e:
        print(f"  pyttsx3 init error: {e}")


def _pyttsx3_gen(text: str, wav_path: Path) -> bool:
    _init_pyttsx3()
    if _pyttsx3_engine is None:
        return False
    try:
        _pyttsx3_engine.save_to_file(text, str(wav_path))
        _pyttsx3_engine.runAndWait()
        time.sleep(0.3)
        return wav_path.exists() and wav_path.stat().st_size > 200
    except Exception:
        return False


def generate_tts(text: str, stem: Path) -> Path | None:
    """
    Generate TTS for `text`. Returns path to audio file (mp3 or wav), or None.
    Caches: if file already exists and is valid, returns immediately.
    """
    mp3 = stem.with_suffix(".mp3")
    wav = stem.with_suffix(".wav")

    if mp3.exists() and mp3.stat().st_size > 500:
        return mp3
    if wav.exists() and wav.stat().st_size > 500:
        return wav

    if _EDGE_OK:
        for attempt in range(3):
            try:
                asyncio.run(_edge_async(text, mp3))
                if mp3.exists() and mp3.stat().st_size > 500:
                    return mp3
                if mp3.exists():
                    mp3.unlink(missing_ok=True)
            except Exception as e:
                print(f"    edge_tts attempt {attempt + 1}: {e}")
                mp3.unlink(missing_ok=True)
                if attempt < 2:
                    time.sleep(1.5 * (attempt + 1))
        print("    edge_tts failed, switching to pyttsx3 fallback")

    if _pyttsx3_gen(text, wav):
        return wav

    return None


def get_duration(audio_path: Path) -> float:
    """Get duration of audio file via moviepy."""
    try:
        clip = AudioFileClip(str(audio_path))
        dur = clip.duration
        clip.close()
        return dur
    except Exception:
        return 0.0


def make_silent_wav(duration: float, path: Path) -> None:
    sr = 22050
    n = max(1, int(duration * sr))
    with wave.open(str(path), "wb") as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(sr)
        wf.writeframes(b"\x00\x00" * n)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE RENDERING
# ═══════════════════════════════════════════════════════════════════════════
def _draw_bullets(
    draw: ImageDraw.ImageDraw,
    bullets: list[str],
    x1: int,
    y1: int,
    x2: int,
    y2: int,
) -> None:
    y = y1
    for bullet in bullets:
        if y + 28 > y2:
            break
        draw.ellipse((x1, y + 7, x1 + 13, y + 20), fill=ACCENT)
        y = _draw_wrapped(draw, bullet, x1 + 22, y, FN, TEXT, x2 - x1 - 24, 5) + 10


def _draw_code_box(
    draw: ImageDraw.ImageDraw,
    title: str,
    code: str,
    x1: int,
    y1: int,
    x2: int,
    y2: int,
) -> None:
    draw.rounded_rectangle(
        (x1, y1, x2, y2), radius=12, fill=CODE_BG, outline=BORDER, width=2
    )
    draw.text((x1 + 14, y1 + 10), title, font=FKB, fill=PRIMARY)
    cy = y1 + 38
    for line in code.splitlines():
        if cy > y2 - 16:
            break
        draw.text((x1 + 14, cy), line.rstrip(), font=FK, fill=TEXT)
        cy += _tsz(draw, line or "M", FK)[1] + 4


def render_base_slide(slide: SlideSpec, idx: int, total: int) -> Image.Image:
    """
    Renders the full slide EXCEPT the subtitle bar.
    The subtitle bar area (CONTENT_H..H) is filled with SUB_BG but left blank —
    add_subtitle_bar() fills it per sentence.
    """
    img = Image.new("RGB", (W, H), BG_TOP)
    _gradient(img)
    draw = ImageDraw.Draw(img)

    # ── title bar ────────────────────────────────────────────────────────
    draw.rectangle((0, 0, W, TITLE_H), fill=TITLE_BG)
    draw.line((0, TITLE_H, W, TITLE_H), fill=ACCENT, width=3)
    draw.text((48, 12), slide.title, font=FT, fill=ACCENT)
    draw.text((50, 58), slide.subtitle, font=FS, fill=TEXT_MUTE)

    badge = f"{idx} / {total}"
    bw, _ = _tsz(draw, badge, FKB)
    bx = W - bw - 50
    draw.rounded_rectangle(
        (bx - 12, 20, W - 20, 64), radius=10, fill=PRIMARY, outline=PRIMARY, width=1
    )
    draw.text((bx, 28), badge, font=FKB, fill=(8, 16, 36))

    # ── layout constants ─────────────────────────────────────────────────
    LEFT_W = 570
    MARGIN = 42
    RSTART = LEFT_W + MARGIN * 2
    CTOP = TITLE_H + 18
    CBOT = CONTENT_H - 8

    # code box splits right panel if needed
    has_code = bool(slide.code_title and slide.code_text)
    IMG_BOT = CBOT - 145 if has_code else CBOT

    # ── left panel (bullets) ──────────────────────────────────────────────
    draw.rounded_rectangle(
        (MARGIN, CTOP, LEFT_W + MARGIN, CBOT),
        radius=16,
        fill=PANEL,
        outline=BORDER,
        width=2,
    )
    _draw_bullets(
        draw, slide.bullets, MARGIN + 18, CTOP + 18, LEFT_W + MARGIN - 10, CBOT - 10
    )

    # ── right panel (image) ───────────────────────────────────────────────
    draw.rounded_rectangle(
        (RSTART, CTOP, W - MARGIN, IMG_BOT),
        radius=16,
        fill=PANEL,
        outline=BORDER,
        width=2,
    )
    _paste_fit(
        img, slide.image_path, (RSTART + 8, CTOP + 8, W - MARGIN - 8, IMG_BOT - 28)
    )
    if slide.image_caption:
        draw.text(
            (RSTART + 12, IMG_BOT - 24),
            slide.image_caption[:90],
            font=FK,
            fill=TEXT_MUTE,
        )

    # ── code / summary box ───────────────────────────────────────────────
    if has_code:
        _draw_code_box(
            draw,
            slide.code_title,
            slide.code_text,  # type: ignore[arg-type]
            RSTART,
            IMG_BOT + 8,
            W - MARGIN,
            CBOT,
        )

    # ── subtitle area placeholder ─────────────────────────────────────────
    draw.rectangle((0, CONTENT_H, W, H), fill=SUB_BG)
    draw.line((0, CONTENT_H, W, CONTENT_H), fill=(40, 44, 68), width=1)

    return img


def add_subtitle_bar(base: Image.Image, sentence: str) -> Image.Image:
    """
    Returns a copy of `base` with `sentence` burned into the subtitle bar.
    This is the per-sentence frame shown while that sentence is spoken.
    """
    result = base.copy()
    draw = ImageDraw.Draw(result)

    # clear and redraw subtitle area
    draw.rectangle((0, CONTENT_H, W, H), fill=SUB_BG)
    draw.line((0, CONTENT_H, W, CONTENT_H), fill=(48, 52, 78), width=1)

    if not sentence:
        return result

    font = FSU
    max_w = W - 180
    lines = _wrap(draw, sentence, font, max_w)[:3]  # max 3 lines
    lh = _tsz(draw, "M", font)[1] + 7
    total = len(lines) * lh
    y = CONTENT_H + (SUB_H - total) // 2

    for line in lines:
        lw, _ = _tsz(draw, line, font)
        draw.text(((W - lw) // 2, y), line, font=font, fill=WHITE)
        y += lh

    return result


# ═══════════════════════════════════════════════════════════════════════════
# SRT
# ═══════════════════════════════════════════════════════════════════════════
def _ts(s: float) -> str:
    h = int(s // 3600)
    m = int((s % 3600) // 60)
    sc = int(s % 60)
    ms = int(round((s - int(s)) * 1000))
    return f"{h:02d}:{m:02d}:{sc:02d},{ms:03d}"


def write_srt(entries: list[tuple[str, float, float]]) -> None:
    lines: list[str] = []
    for i, (text, start, end) in enumerate(entries, 1):
        lines += [str(i), f"{_ts(start)} --> {_ts(end)}", text, ""]
    SRT_PATH.write_text("\n".join(lines), encoding="utf-8")
    print(f"SRT  → {SRT_PATH}  ({len(entries)} entries)")


# ═══════════════════════════════════════════════════════════════════════════
# PPTX  (uses base slide images, no subtitle bar)
# ═══════════════════════════════════════════════════════════════════════════
def build_pptx(slides: list[SlideSpec], base_paths: list[Path]) -> None:
    if not _PPTX_OK or Presentation is None:
        print("python-pptx not available, skipping PPTX.")
        return
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for img_path, spec in zip(base_paths, slides):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        tb = sl.shapes.add_textbox(
            Inches(0.38), Inches(0.12), Inches(10.5), Inches(0.52)
        )
        p = tb.text_frame.paragraphs[0]
        p.text = spec.title
        p.font.size = Pt(28)
        p.font.bold = True
        tb2 = sl.shapes.add_textbox(
            Inches(0.40), Inches(0.62), Inches(12.2), Inches(0.34)
        )
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = spec.subtitle
        p2.font.size = Pt(14)
        sl.shapes.add_picture(
            str(img_path),
            Inches(0.28),
            Inches(1.02),
            width=Inches(12.78),
            height=Inches(6.18),
        )
    prs.save(str(PPTX_PATH))
    print(f"PPTX → {PPTX_PATH}")


# ═══════════════════════════════════════════════════════════════════════════
# SCRIPT MARKDOWN
# ═══════════════════════════════════════════════════════════════════════════
def build_script(slides: list[SlideSpec]) -> None:
    lines = [
        "# 知识图谱核心功能教学视频脚本 v6",
        "",
        "## 技术特性",
        "",
        "- **edge_tts** (zh-CN-XiaoxiaoNeural) 高质量中文语音",
        "- **moviepy** 精确帧级同步：每句独立帧，一次性编码，彻底消除漂移",
        "- 字幕逐句烧录到画面；SRT 时间戳来自真实 MP3 时长",
        "",
        "## 成片结构",
        "",
    ]
    for i, s in enumerate(slides, 1):
        lines.append(f"{i}. **{s.title}** — {s.subtitle}")
    lines += ["", "## 讲解稿", ""]
    for i, s in enumerate(slides, 1):
        lines += [
            f"### {i}. {s.title}",
            f"> {s.subtitle}",
            "",
            f"**旁白：** {s.narration}",
            "",
        ]
    SCRIPT_PATH.write_text("\n".join(lines), encoding="utf-8")
    print(f"Script → {SCRIPT_PATH}")


# ═══════════════════════════════════════════════════════════════════════════
# MAIN VIDEO PIPELINE
# ═══════════════════════════════════════════════════════════════════════════
def build_video(slides: list[SlideSpec]) -> tuple[list[Path], bool]:
    """
    Core sync pipeline:
      For each slide → for each sentence:
        1. Generate TTS audio via edge_tts (or pyttsx3 fallback)
        2. Render base slide image once per slide
        3. Overlay subtitle text → per-sentence PNG
        4. moviepy clip: ImageClip(png).set_duration(real_dur + tail).set_audio(AudioFileClip(mp3))
      concatenate_videoclips → write_videofile
        → single AAC encode pass, no per-clip encoding delay, zero drift
      SRT timestamps = cumulative real audio durations, exact match with burned subtitle.
    """
    if not _MOVIEPY_OK:
        print("ERROR: moviepy is required. pip install moviepy==1.0.3")
        return [], False

    all_clips: list = []
    srt_entries: list[tuple[str, float, float]] = []
    base_paths: list[Path] = []
    t = 0.0  # running video timestamp (seconds)
    total = len(slides)

    for si, slide in enumerate(slides, 1):
        print(f"\n[{si:02d}/{total}] {slide.title}")

        # render base slide image (once per slide)
        base_img = render_base_slide(slide, si, total)
        base_path = SLIDES_DIR / f"{si:02d}_{_safe(slide.title)}.png"
        base_img.save(base_path, "PNG")
        base_paths.append(base_path)

        sentences = get_sentences(slide)
        nsent = len(sentences)

        for ki, sentence in enumerate(sentences):
            is_last = ki == nsent - 1
            tail = SLIDE_TAIL if is_last else CLIP_TAIL

            audio_stem = AUDIO_DIR / f"{si:02d}_{ki + 1:02d}"
            print(f"  [{ki + 1}/{nsent}] TTS → {sentence[:55]}")

            audio_path = generate_tts(sentence, audio_stem)

            # silent fallback
            if audio_path is None or not audio_path.exists():
                audio_path = audio_stem.with_suffix(".wav")
                estimated = max(1.8, len(sentence) * 0.14)
                make_silent_wav(estimated, audio_path)

            real_dur = get_duration(audio_path)
            if real_dur < 0.05:
                real_dur = max(1.8, len(sentence) * 0.14)

            # per-sentence frame
            frame_img = add_subtitle_bar(base_img, sentence)
            frame_path = SENT_DIR / f"{si:02d}_{ki + 1:02d}.png"
            frame_img.save(frame_path, "PNG")

            # moviepy clip  ← THE KEY SYNC MECHANISM
            clip_dur = real_dur + tail
            try:
                audio_clip = AudioFileClip(str(audio_path))
                vid_clip = (
                    ImageClip(str(frame_path))
                    .set_duration(clip_dur)
                    .set_audio(audio_clip)
                )
            except Exception as e:
                print(f"    clip error: {e}  (silent fallback)")
                vid_clip = ImageClip(str(frame_path)).set_duration(clip_dur)

            all_clips.append(vid_clip)

            # SRT entry covers the speech portion (not the tail silence)
            srt_entries.append((sentence, t, t + real_dur))
            t += clip_dur

    # ── write SRT ───────────────────────────────────────────────────────
    write_srt(srt_entries)

    # ── concatenate and encode ───────────────────────────────────────────
    print(f"\nConcatenating {len(all_clips)} clips  (total ≈ {t:.1f} s)...")
    try:
        final = concatenate_videoclips(all_clips, method="compose")
        print(f"Encoding video → {VIDEO_PATH}")
        final.write_videofile(
            str(VIDEO_PATH),
            fps=FPS,
            codec="libx264",
            audio_codec="aac",
            audio_bitrate="192k",
            temp_audiofile=str(TEMP_DIR / "tmp_audio_v6.m4a"),
            remove_temp=True,
            verbose=True,
            write_logfile=False,
            threads=4,
        )
        final.close()
        print(f"Video → {VIDEO_PATH}")
        return base_paths, True
    except Exception as e:
        print(f"ERROR during video write: {e}")
        import traceback

        traceback.print_exc()
        return base_paths, False


# ═══════════════════════════════════════════════════════════════════════════
# ENTRY POINT
# ═══════════════════════════════════════════════════════════════════════════
def main() -> None:
    ensure_dirs()

    # Check UML source images
    required = [FRAMEWORK_IMG, CLASS_IMG, SEQUENCE_IMG, PROCESS_IMG]
    missing = [p for p in required if not p.exists()]
    if missing:
        print("Missing UML images:")
        for p in missing:
            print(f"  {p}")
        raise SystemExit(1)

    print("=" * 60)
    print("  知识图谱教学视频生成器 v6")
    print("  edge_tts + moviepy  严格帧级同步")
    print("=" * 60)

    print("\n[1/6] Downloading official brand assets...")
    official = ensure_official_assets()
    for k, v in official.items():
        print(f"  {k}: {v.name}")

    print("\n[2/6] Generating UML crops...")
    crops = generate_uml_crops()
    print(f"  {len(crops)} crops generated")

    print("\n[3/6] Generating runtime screenshots...")
    shots = generate_runtime_shots()
    print(f"  {len(shots)} screenshots generated")

    print("\n[4/6] Building slides...")
    slides = build_slides(crops, shots, official)
    print(f"  {len(slides)} slides defined")

    print("\n[5/6] Building video (TTS + moviepy)...")
    base_paths, ok = build_video(slides)

    print("\n[6/6] Building PPTX and script...")
    build_pptx(slides, base_paths)
    build_script(slides)

    print("\n" + "=" * 60)
    print("v6 generation complete")
    print(f"  Video  : {VIDEO_PATH}")
    print(f"  PPTX   : {PPTX_PATH}")
    print(f"  SRT    : {SRT_PATH}")
    print(f"  Script : {SCRIPT_PATH}")
    if not ok:
        print("  !! Video generation FAILED – see errors above.")
    print("=" * 60)


if __name__ == "__main__":
    main()
