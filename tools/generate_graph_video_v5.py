from __future__ import annotations

import os
import shutil
import struct
import subprocess
import urllib.request
import wave
from dataclasses import dataclass
from pathlib import Path
from typing import Protocol, cast

from PIL import Image, ImageDraw, ImageFont

PillowFont = ImageFont.ImageFont | ImageFont.FreeTypeFont | ImageFont.TransposedFont

try:
    import pyttsx3
except Exception:
    pyttsx3 = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None
    Inches = None
    Pt = None


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
DIAGRAMS_V3_DIR = DOCS_DIR / "diagrams" / "v3"
VIDEO_V5_DIR = DOCS_DIR / "video" / "v5"
TESTING_DIR = DOCS_DIR / "testing"
VIDEO_OUTPUT_DIR = ROOT / "video_output"

SLIDES_DIR = VIDEO_V5_DIR / "slides"
AUDIO_DIR = VIDEO_V5_DIR / "audio"
CLIPS_DIR = VIDEO_V5_DIR / "clips"
TEMP_DIR = VIDEO_V5_DIR / "temp"
CROPS_DIR = VIDEO_V5_DIR / "crops"
ASSETS_DIR = VIDEO_V5_DIR / "assets"
SEGMENT_AUDIO_DIR = AUDIO_DIR / "segments"

FRAMEWORK_IMG = DIAGRAMS_V3_DIR / "flutter_dart_framework_architecture.png"
CLASS_IMG = DIAGRAMS_V3_DIR / "flutter_dart_core_class_diagram.png"
SEQUENCE_IMG = DIAGRAMS_V3_DIR / "graph_feature_sequence_diagram.png"
PROCESS_IMG = DIAGRAMS_V3_DIR / "knowledge_graph_development_process.png"

SCRIPT_PATH = VIDEO_V5_DIR / "video_script_v5.md"
PPTX_PATH = VIDEO_OUTPUT_DIR / "知识图谱核心功能_图谱功能教程_v5.pptx"
VIDEO_PATH = VIDEO_OUTPUT_DIR / "知识图谱核心功能_图谱功能教程_v5.mp4"
SRT_PATH = VIDEO_V5_DIR / "video_subtitles_v5.srt"
FLUTTER_OFFICIAL_IMG = ASSETS_DIR / "flutter_official.png"
DART_OFFICIAL_IMG = ASSETS_DIR / "dart_official.png"

MODEL_TEST_PATH = ROOT / "test" / "models" / "model_test.dart"
HOME_WIDGET_TEST_PATH = ROOT / "test" / "widgets" / "home_page_widget_test.dart"
LOGIN_WIDGET_TEST_PATH = ROOT / "test" / "widget_test.dart"
TEST_CASES_PATH = TESTING_DIR / "test_cases.md"
TEST_REPORT_PATH = TESTING_DIR / "test_report.md"

LOGIN_GOLDEN_PATH = ROOT / "test" / "screenshots" / "goldens" / "login_page.png"
HOME_GOLDEN_PATH = ROOT / "test" / "screenshots" / "goldens" / "home_page.png"

WIDTH = 1920
HEIGHT = 1080

BG_TOP = (244, 248, 255)
BG_BOTTOM = (230, 238, 252)
WHITE = (255, 255, 255)
TEXT = (34, 44, 66)
TEXT_LIGHT = (92, 104, 128)
PRIMARY = (76, 110, 245)
PRIMARY_DARK = (48, 77, 201)
ACCENT = (18, 184, 134)
ORANGE = (255, 146, 43)
RED = (235, 87, 87)
PURPLE = (125, 92, 255)
BORDER = (210, 220, 240)
PANEL_BG = (250, 252, 255)
CODE_BG = (245, 247, 251)

TITLE_FONT_SIZE = 52
SUBTITLE_FONT_SIZE = 28
SECTION_FONT_SIZE = 28
BODY_FONT_SIZE = 25
SMALL_FONT_SIZE = 20
CODE_FONT_SIZE = 20


class VoiceLike(Protocol):
    id: str
    name: str


class TTSEngineLike(Protocol):
    def getProperty(self, name: str) -> object: ...
    def setProperty(self, name: str, value: str | float) -> None: ...
    def save_to_file(
        self, text: str, filename: str, name: str | None = None
    ) -> None: ...
    def runAndWait(self) -> None: ...
    def stop(self) -> None: ...


@dataclass
class SlideSpec:
    title: str
    subtitle: str
    bullets: list[str]
    narration: str
    image_path: Path | None = None
    image_caption: str | None = None
    code_title: str | None = None
    code_text: str | None = None
    voice_segments: list[str] | None = None
    summary_text: str | None = None


def ensure_dirs() -> None:
    for path in [
        VIDEO_V5_DIR,
        SLIDES_DIR,
        AUDIO_DIR,
        CLIPS_DIR,
        TEMP_DIR,
        CROPS_DIR,
        ASSETS_DIR,
        SEGMENT_AUDIO_DIR,
        VIDEO_OUTPUT_DIR,
    ]:
        path.mkdir(parents=True, exist_ok=True)


def download_if_missing(url: str, out_path: Path) -> Path | None:
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    try:
        urllib.request.urlretrieve(url, str(out_path))
        return out_path if out_path.exists() else None
    except Exception:
        return None


def ensure_official_assets() -> dict[str, Path]:
    assets: dict[str, Path] = {}

    flutter = download_if_missing(
        "https://flutter.dev/assets/lockup_flutter_horizontal.d0515092173211776ceed19b39c2a041.png",
        FLUTTER_OFFICIAL_IMG,
    )
    dart = download_if_missing(
        "https://dart.dev/assets/img/logo/dart-logo-for-shares.png",
        DART_OFFICIAL_IMG,
    )

    if flutter is not None:
        assets["flutter"] = flutter
    if dart is not None:
        assets["dart"] = dart

    return assets


def load_font(size: int, bold: bool = False):
    candidates: list[str] = []
    if os.name == "nt":
        if bold:
            candidates.extend(
                [
                    r"C:\Windows\Fonts\msyhbd.ttc",
                    r"C:\Windows\Fonts\simhei.ttf",
                    r"C:\Windows\Fonts\arialbd.ttf",
                    r"C:\Windows\Fonts\calibrib.ttf",
                ]
            )
        else:
            candidates.extend(
                [
                    r"C:\Windows\Fonts\msyh.ttc",
                    r"C:\Windows\Fonts\simsun.ttc",
                    r"C:\Windows\Fonts\arial.ttf",
                    r"C:\Windows\Fonts\calibri.ttf",
                ]
            )
    else:
        candidates.extend(
            [
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            ]
        )

    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size=size)
        except Exception:
            continue
    return ImageFont.load_default()


FONT_TITLE = load_font(TITLE_FONT_SIZE, True)
FONT_SUBTITLE = load_font(SUBTITLE_FONT_SIZE, False)
FONT_SECTION = load_font(SECTION_FONT_SIZE, True)
FONT_BODY = load_font(BODY_FONT_SIZE, False)
FONT_BODY_BOLD = load_font(BODY_FONT_SIZE, True)
FONT_SMALL = load_font(SMALL_FONT_SIZE, False)
FONT_SMALL_BOLD = load_font(SMALL_FONT_SIZE, True)
FONT_CODE = load_font(CODE_FONT_SIZE, False)


def draw_vertical_gradient(
    img: Image.Image, top: tuple[int, int, int], bottom: tuple[int, int, int]
) -> None:
    draw = ImageDraw.Draw(img)
    for y in range(img.height):
        ratio = y / max(1, img.height - 1)
        color = tuple(int(top[i] * (1 - ratio) + bottom[i] * ratio) for i in range(3))
        draw.line([(0, y), (img.width, y)], fill=color)


def rounded_rect(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int, int, int],
    fill: tuple[int, int, int],
    outline: tuple[int, int, int] | None = None,
    radius: int = 24,
    width: int = 2,
) -> None:
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def text_size(
    draw: ImageDraw.ImageDraw, text: str, font: PillowFont
) -> tuple[int, int]:
    bbox = draw.textbbox((0, 0), text, font=font)
    return int(bbox[2] - bbox[0]), int(bbox[3] - bbox[1])


def center_text(
    draw: ImageDraw.ImageDraw,
    x: int,
    y: int,
    text: str,
    font: PillowFont,
    fill: tuple[int, int, int],
) -> None:
    w, h = text_size(draw, text, font)
    draw.text((x - w / 2, y - h / 2), text, font=font, fill=fill)


def wrap_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    font: PillowFont,
    max_width: int,
) -> list[str]:
    lines: list[str] = []
    for paragraph in text.splitlines():
        p = paragraph.strip()
        if not p:
            lines.append("")
            continue
        current = ""
        for ch in p:
            candidate = current + ch
            w, _ = text_size(draw, candidate, font)
            if w <= max_width or not current:
                current = candidate
            else:
                lines.append(current)
                current = ch
        if current:
            lines.append(current)
    return lines


def draw_wrapped_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    box: tuple[int, int, int, int],
    font: PillowFont,
    fill: tuple[int, int, int],
    line_spacing: int = 6,
) -> int:
    x1, y1, x2, y2 = box
    y = y1
    max_width = x2 - x1
    for line in wrap_text(draw, text, font, max_width):
        if y > y2:
            break
        draw.text((x1, y), line, font=font, fill=fill)
        _, h = text_size(draw, line or "中", font)
        y += h + line_spacing
    return y


def safe_name(text: str) -> str:
    keep = []
    for ch in text:
        if ch.isalnum() or ch in ("_", "-", " "):
            keep.append(ch)
        elif "\u4e00" <= ch <= "\u9fff":
            keep.append(ch)
        else:
            keep.append("_")
    return "".join(keep).strip().replace(" ", "_")


def read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8") if path.exists() else ""


def extract_block_by_marker(path: Path, marker: str, lines_after: int = 25) -> str:
    text = read_text(path)
    if not text:
        return ""
    lines = text.splitlines()
    for i, line in enumerate(lines):
        if marker in line:
            return "\n".join(lines[i : i + lines_after])
    return "\n".join(lines[:lines_after])


def select_lines(text: str, keywords: list[str], max_lines: int = 10) -> str:
    selected: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        lowered = stripped.lower()
        if any(key.lower() in lowered for key in keywords):
            selected.append(stripped)
        if len(selected) >= max_lines:
            break
    return "\n".join(selected) if selected else "\n".join(text.splitlines()[:max_lines])


def load_test_summary() -> str:
    report = read_text(TEST_REPORT_PATH)
    return select_lines(
        report,
        ["23", "pass", "通过", "apk", "debug", "release", "测试", "构建"],
        10,
    )


def load_test_case_summary() -> str:
    cases = read_text(TEST_CASES_PATH)
    return select_lines(
        cases,
        [
            "tc-",
            "图谱",
            "登录",
            "测验",
            "资料",
            "视频",
            "progress",
            "model",
        ],
        16,
    )


def load_code_summaries() -> tuple[str, str, str]:
    model_code = extract_block_by_marker(MODEL_TEST_PATH, "group('GraphModel'", 30)
    home_code = extract_block_by_marker(
        HOME_WIDGET_TEST_PATH, "group('HomePage widget tests'", 30
    )
    login_code = extract_block_by_marker(
        LOGIN_WIDGET_TEST_PATH, "testWidgets('Login page shows core UI elements'", 24
    )
    return model_code, home_code, login_code


def crop_and_save(
    source_path: Path,
    crop_box: tuple[int, int, int, int],
    output_name: str,
    border_color: tuple[int, int, int] = PRIMARY,
) -> Path:
    img = Image.open(source_path).convert("RGB")
    cropped = img.crop(crop_box)

    canvas = Image.new("RGB", (1600, 900), WHITE)
    draw_vertical_gradient(canvas, BG_TOP, BG_BOTTOM)
    draw = ImageDraw.Draw(canvas)

    rounded_rect(draw, (18, 18, 1582, 882), PANEL_BG, BORDER, 22, 3)

    target_w = 1530
    target_h = 840
    ratio = min(target_w / cropped.width, target_h / cropped.height)
    new_size = (max(1, int(cropped.width * ratio)), max(1, int(cropped.height * ratio)))
    resized = cropped.resize(new_size)
    x = 35 + (target_w - new_size[0]) // 2
    y = 30 + (target_h - new_size[1]) // 2
    canvas.paste(resized, (x, y))

    draw.rounded_rectangle(
        (18, 18, 1582, 882), radius=22, outline=border_color, width=4
    )
    out = CROPS_DIR / output_name
    canvas.save(out)
    return out


def generate_uml_crops() -> dict[str, Path]:
    crops: dict[str, Path] = {}

    if FRAMEWORK_IMG.exists():
        img = Image.open(FRAMEWORK_IMG)
        w, h = img.size
        crops["framework_full"] = FRAMEWORK_IMG
        crops["framework_ui"] = crop_and_save(
            FRAMEWORK_IMG,
            (int(w * 0.02), int(h * 0.02), int(w * 0.40), int(h * 0.56)),
            "framework_ui_crop.png",
            PRIMARY,
        )
        crops["framework_dao"] = crop_and_save(
            FRAMEWORK_IMG,
            (int(w * 0.26), int(h * 0.30), int(w * 0.64), int(h * 0.78)),
            "framework_dao_crop.png",
            ACCENT,
        )
        crops["framework_right"] = crop_and_save(
            FRAMEWORK_IMG,
            (int(w * 0.58), int(h * 0.20), int(w * 0.95), int(h * 0.66)),
            "framework_right_crop.png",
            ORANGE,
        )

    if CLASS_IMG.exists():
        img = Image.open(CLASS_IMG)
        w, h = img.size
        crops["class_full"] = CLASS_IMG
        crops["class_ui"] = crop_and_save(
            CLASS_IMG,
            (int(w * 0.02), int(h * 0.03), int(w * 0.42), int(h * 0.58)),
            "class_ui_crop.png",
            PRIMARY,
        )
        crops["class_dao"] = crop_and_save(
            CLASS_IMG,
            (int(w * 0.44), int(h * 0.22), int(w * 0.76), int(h * 0.64)),
            "class_dao_crop.png",
            ACCENT,
        )
        crops["class_model"] = crop_and_save(
            CLASS_IMG,
            (int(w * 0.56), int(h * 0.56), int(w * 0.94), int(h * 0.92)),
            "class_model_crop.png",
            PURPLE,
        )

    if SEQUENCE_IMG.exists():
        img = Image.open(SEQUENCE_IMG)
        w, h = img.size
        crops["sequence_full"] = SEQUENCE_IMG
        crops["sequence_enter"] = crop_and_save(
            SEQUENCE_IMG,
            (int(w * 0.04), 0, int(w * 0.96), int(h * 0.36)),
            "sequence_enter_crop.png",
            PRIMARY,
        )
        crops["sequence_detail"] = crop_and_save(
            SEQUENCE_IMG,
            (int(w * 0.04), int(h * 0.24), int(w * 0.96), int(h * 0.62)),
            "sequence_detail_crop.png",
            ACCENT,
        )
        crops["sequence_action"] = crop_and_save(
            SEQUENCE_IMG,
            (int(w * 0.04), int(h * 0.56), int(w * 0.96), int(h * 0.92)),
            "sequence_action_crop.png",
            ORANGE,
        )

    if PROCESS_IMG.exists():
        img = Image.open(PROCESS_IMG)
        w, h = img.size
        crops["process_full"] = PROCESS_IMG
        crops["process_early"] = crop_and_save(
            PROCESS_IMG,
            (int(w * 0.02), int(h * 0.06), int(w * 0.34), int(h * 0.46)),
            "process_early_crop.png",
            PRIMARY,
        )
        crops["process_middle"] = crop_and_save(
            PROCESS_IMG,
            (int(w * 0.30), int(h * 0.28), int(w * 0.68), int(h * 0.68)),
            "process_middle_crop.png",
            ACCENT,
        )
        crops["process_late"] = crop_and_save(
            PROCESS_IMG,
            (int(w * 0.60), int(h * 0.50), int(w * 0.97), int(h * 0.92)),
            "process_late_crop.png",
            PURPLE,
        )

    return crops


def generate_mock_runtime_screenshots() -> dict[str, Path]:
    out: dict[str, Path] = {}

    def base_page(
        title: str,
        subtitle: str,
        tags: list[str],
        cards: list[tuple[str, str, tuple[int, int, int]]],
        filename: str,
    ) -> Path:
        img = Image.new("RGB", (1440, 900), WHITE)
        draw_vertical_gradient(img, BG_TOP, BG_BOTTOM)
        draw = ImageDraw.Draw(img)

        rounded_rect(draw, (28, 28, 1412, 872), WHITE, BORDER, 24, 3)
        rounded_rect(draw, (28, 28, 1412, 120), PRIMARY, PRIMARY, 24, 1)

        draw.text((60, 52), title, font=FONT_TITLE, fill=WHITE)
        draw.text((62, 98), subtitle, font=FONT_SMALL, fill=(230, 236, 255))

        x = 60
        for tag in tags:
            w, h = text_size(draw, tag, FONT_SMALL_BOLD)
            rounded_rect(
                draw, (x, 150, x + w + 34, 194), (238, 243, 255), BORDER, 16, 1
            )
            draw.text((x + 16, 160), tag, font=FONT_SMALL_BOLD, fill=PRIMARY_DARK)
            x += w + 52

        y = 240
        for card_title, card_desc, color in cards:
            rounded_rect(draw, (60, y, 1380, y + 150), WHITE, BORDER, 24, 3)
            rounded_rect(draw, (84, y + 24, 110, y + 50), color, color, 10, 1)
            draw.text((128, y + 18), card_title, font=FONT_BODY_BOLD, fill=TEXT)
            draw_wrapped_text(
                draw, card_desc, (128, y + 60, 1335, y + 128), FONT_BODY, TEXT_LIGHT
            )
            y += 175

        out_path = CROPS_DIR / filename
        img.save(out_path)
        return out_path

    if LOGIN_GOLDEN_PATH.exists():
        out["login"] = LOGIN_GOLDEN_PATH
    else:
        out["login"] = base_page(
            "LoginPage 登录页",
            "登录页运行效果（页面示意图）",
            ["统一入口", "表单输入", "登录校验"],
            [
                ("用户名输入", "用户输入用户名和密码并提交登录请求", PRIMARY),
                ("状态反馈", "表单校验失败时给出提示信息", ORANGE),
                ("进入首页", "登录成功后跳转至 HomePage", ACCENT),
            ],
            "runtime_login_mock.png",
        )

    if HOME_GOLDEN_PATH.exists():
        out["home"] = HOME_GOLDEN_PATH
    else:
        out["home"] = base_page(
            "HomePage 首页",
            "首页运行效果（页面示意图）",
            ["主导航", "课程入口", "学习记录"],
            [
                ("底部导航", "在首页、图谱、资料、视频、测验与进度之间切换", PRIMARY),
                ("功能入口", "把分散模块统一组织成清晰的学习入口", ACCENT),
                ("学习闭环", "通过进度与测验形成反馈机制", ORANGE),
            ],
            "runtime_home_mock.png",
        )

    out["graph_list"] = base_page(
        "GraphListPage 图谱列表",
        "图谱页运行效果（页面示意图）",
        ["图谱入口", "列表加载", "图谱选择"],
        [
            ("01 课程图谱详细图谱", "展示课程整体知识结构与章节关系", PRIMARY),
            ("02 技术栈图谱详细图谱", "展示 Flutter、Dart、SQLite 等技术关系", ACCENT),
            ("06 学习图谱详细图谱", "展示学习路径、复习节点与实践内容", ORANGE),
        ],
        "runtime_graph_list_mock.png",
    )

    out["graph_detail"] = base_page(
        "GraphDetailPage 图谱详情",
        "图谱详情运行效果（页面示意图）",
        ["节点交互", "缩放拖拽", "学习记录", "收藏"],
        [
            (
                "图谱可视化区域",
                "通过 CustomPainter 与 InteractiveViewer 浏览结构",
                PRIMARY,
            ),
            ("节点详情卡片", "点击节点后显示标题、类型、内容和操作按钮", ACCENT),
            ("开始学习 / 收藏", "把浏览行为转化为学习记录和收藏数据", RED),
        ],
        "runtime_graph_detail_mock.png",
    )

    return out


def build_script_markdown(slides: list[SlideSpec]) -> None:
    lines = [
        "# 知识图谱核心功能教学视频脚本 v5",
        "",
        "## 说明",
        "",
        "- v5 重点解决语音、字幕与页面信息同步问题。",
        "- 成片中去除教学提示等说明性标签，改为更简洁的课程展示风格。",
        "- 新增 Flutter 与 Dart 介绍页，并强化知识图谱核心流程说明。",
        "",
        "## 成片结构",
        "",
    ]
    for i, slide in enumerate(slides, start=1):
        lines.append(f"{i}. {slide.title} - {slide.subtitle}")
    lines += ["", "## 讲解稿", ""]
    for i, slide in enumerate(slides, start=1):
        lines += [
            f"### {i}. {slide.title}",
            f"- 副标题：{slide.subtitle}",
            f"- 要点：{', '.join(slide.bullets)}",
            f"- 讲解：{slide.narration}",
            "",
        ]
    _ = SCRIPT_PATH.write_text("\n".join(lines), encoding="utf-8")


def build_slides(
    crops: dict[str, Path],
    runtime_shots: dict[str, Path],
    official_assets: dict[str, Path],
) -> list[SlideSpec]:
    model_code, _home_code, _login_code = load_code_summaries()
    test_case_summary = load_test_case_summary()
    test_report_summary = load_test_summary()

    return [
        SlideSpec(
            title="课程导入",
            subtitle="先了解学习目标与观看重点",
            bullets=[
                "整体结构",
                "技术栈",
                "核心流程",
                "测试结果",
            ],
            narration=(
                "欢迎进入知识图谱核心功能教学视频。"
                "这一节会围绕系统结构、Flutter 与 Dart 的实战使用方式、知识图谱核心流程以及测试结果展开。"
                "请先建立整体认识，再逐步进入页面、数据和交互细节。"
            ),
            image_path=crops.get("framework_full"),
        ),
        SlideSpec(
            title="项目总览",
            subtitle="认识知识图谱学习系统的整体结构",
            bullets=[
                "Flutter 页面层",
                "Dart 业务层",
                "SQLite 数据层",
                "知识图谱驱动学习",
                "测试与交付材料",
            ],
            narration=(
                "这是整个项目的总体视图。"
                "系统以 Flutter 负责页面表现，以 Dart 负责业务逻辑与数据组织，以 SQLite 负责本地持久化。"
                "知识图谱位于系统核心位置，用来连接页面导航、节点学习、收藏行为和学习进度。"
            ),
            image_path=crops.get("framework_full"),
        ),
        SlideSpec(
            title="Flutter 技术介绍",
            subtitle="Flutter 官方介绍与本项目最佳实战",
            bullets=[
                "开源框架",
                "单一代码库",
                "多平台应用",
                "Fast / Productive / Flexible",
                "界面层与业务层分离",
            ],
            narration=(
                "Flutter 官方将自己定义为一个开源框架，用于通过单一代码库构建精美的原生编译多平台应用。"
                "官方还强调它具有 Fast、Productive 和 Flexible 三个特点。"
                "放到本项目中，最佳实战是让 Flutter 专注页面渲染、导航和交互反馈，而把数据访问与业务逻辑放到独立的 Dart 层。"
            ),
            image_path=official_assets.get("flutter", runtime_shots.get("home")),
            image_caption="Flutter 官方品牌图与官方定位：单一代码库构建多平台应用",
            voice_segments=[
                "Flutter 官方将自己定义为一个开源框架，",
                "用于通过单一代码库构建精美的原生编译多平台应用。",
                "官方还强调它具有 Fast、Productive 和 Flexible 三个特点。",
                "放到本项目中，最佳实战是让 Flutter 专注页面渲染、导航和交互反馈，",
                "而把数据访问与业务逻辑放到独立的 Dart 层。",
            ],
            summary_text=(
                "Flutter 官方定位：\n"
                "开源框架、单一代码库、多平台应用。\n"
                "项目实战：页面层负责渲染与交互，业务逻辑不直接堆在 Widget 内。"
            ),
        ),
        SlideSpec(
            title="Dart 技术介绍",
            subtitle="Dart 官方介绍与本项目最佳实战",
            bullets=[
                "client-optimized language",
                "fast apps on any platform",
                "type safe",
                "sound null safety",
                "模型与 DAO 分层",
            ],
            narration=(
                "Dart 官方将自己定义为一种面向客户端优化的语言，用于在任何平台上开发快速应用。"
                "官方还强调它具备类型安全和 sound null safety，并能够支撑高效开发体验。"
                "放到本项目中，最佳实战是用 Dart 承载模型对象、DAO 层和业务逻辑，让页面直接消费结构化结果，而不是拼接零散数据。"
            ),
            image_path=official_assets.get("dart", crops.get("class_dao")),
            image_caption="Dart 官方品牌图与官方定位：client-optimized language for fast apps",
            voice_segments=[
                "Dart 官方将自己定义为一种面向客户端优化的语言，",
                "用于在任何平台上开发快速应用。",
                "官方还强调它具备类型安全和 sound null safety，",
                "并能够支撑高效开发体验。",
                "放到本项目中，最佳实战是用 Dart 承载模型对象、DAO 层和业务逻辑，",
                "让页面直接消费结构化结果，而不是拼接零散数据。",
            ],
            summary_text=(
                "Dart 官方定位：\n"
                "client-optimized language for fast apps on any platform。\n"
                "项目实战：模型、DAO、业务逻辑分层，页面直接使用结构化数据。"
            ),
        ),
        SlideSpec(
            title="页面职责分析",
            subtitle="从页面分工理解系统学习闭环",
            bullets=[
                "LoginPage 统一登录入口",
                "HomePage 负责导航分发",
                "GraphListPage 负责图谱选择",
                "GraphDetailPage 负责核心交互",
                "资源与进度页形成学习闭环",
            ],
            narration=(
                "系统中的页面不是简单罗列，而是围绕学习闭环进行分工。"
                "登录页负责统一入口，首页负责模块导航，图谱列表页负责主题选择，图谱详情页负责节点浏览、开始学习和收藏。"
                "资料、视频、测验和进度页面则继续承接学习行为，形成完整的使用路径。"
            ),
            image_path=runtime_shots.get("home"),
        ),
        SlideSpec(
            title="系统框架总览",
            subtitle="先建立页面层、逻辑层和数据层的整体认识",
            bullets=[
                "页面层负责展示",
                "业务层负责规则",
                "数据层负责持久化",
                "图谱连接学习模块",
            ],
            narration=(
                "从框架图可以看出，系统分为页面层、业务逻辑层和数据层。"
                "页面层承接用户交互，业务层负责认证、图谱查询、收藏管理和学习记录，数据层负责 SQLite 初始化与读写。"
                "知识图谱位于中间纽带位置，向上驱动页面，向下连接数据。"
            ),
            image_path=crops.get("framework_full"),
        ),
        SlideSpec(
            title="UI 层讲解",
            subtitle="把页面名称与页面职责一一对应起来",
            bullets=[
                "LoginPage 登录",
                "HomePage 导航",
                "GraphListPage 入口",
                "GraphDetailPage 核心交互",
                "资料 / 视频 / 测验 / 进度联动",
            ],
            narration=(
                "这里聚焦 UI 层。"
                "理解 UI 层的关键不是记页面名称，而是把每个页面放回用户学习流程中。"
                "用户先登录，再通过首页导航进入图谱模块，随后在图谱详情里完成节点浏览、学习和收藏，最后再联动其他学习资源。"
            ),
            image_path=crops.get("framework_ui"),
        ),
        SlideSpec(
            title="DAO 与数据层讲解",
            subtitle="理解页面背后的业务逻辑与数据库访问",
            bullets=[
                "AuthService 管理登录状态",
                "GraphDao 读取图谱数据",
                "FavoriteDao 管理收藏",
                "LearningRecordDao 管理学习记录",
                "DatabaseHelper 统一管理 SQLite",
            ],
            narration=(
                "页面背后的核心能力来自 Dart 业务层和数据层。"
                "AuthService 负责登录状态，GraphDao 负责图谱、节点和边的查询，FavoriteDao 与 LearningRecordDao 负责行为数据写入。"
                "所有这些能力最终都建立在 DatabaseHelper 对 SQLite 的统一管理之上。"
            ),
            image_path=crops.get("framework_dao"),
        ),
        SlideSpec(
            title="核心类图总览",
            subtitle="从类图理解页面层、DAO 层和模型层",
            bullets=[
                "页面类负责交互",
                "DAO 类负责访问",
                "Model 类负责结构化对象",
                "GraphDetailPage 是核心类",
            ],
            narration=(
                "类图体现了系统的职责边界。"
                "页面类负责展示与交互，DAO 类负责查询和写入，模型类负责把数据库结果组织成清晰对象。"
                "其中 GraphDetailPage 是知识图谱功能的核心页面，也是用户操作最集中的位置。"
            ),
            image_path=crops.get("class_full"),
        ),
        SlideSpec(
            title="UI 类职责讲解",
            subtitle="区分应用启动、导航分发与图谱交互类",
            bullets=[
                "MyApp 应用启动",
                "HomePage 主导航",
                "GraphListPage 图谱入口",
                "GraphDetailPage 图谱交互",
                "Progress / Favorites 提供反馈与复习",
            ],
            narration=(
                "放大 UI 类之后，可以更清楚地看到页面类的分工。"
                "MyApp 负责应用入口，HomePage 负责导航分发，GraphListPage 负责图谱选择，GraphDetailPage 负责节点渲染与操作。"
                "其他页面则围绕学习结果和复习入口继续延展。"
            ),
            image_path=crops.get("class_ui"),
        ),
        SlideSpec(
            title="DAO 与模型层讲解",
            subtitle="理解数据查询、封装与对象协同",
            bullets=[
                "GraphDao 查询图谱",
                "FavoriteDao 管理收藏",
                "LearningRecordDao 管理记录",
                "GraphModel / NodeModel / EdgeModel 封装结构",
            ],
            narration=(
                "DAO 层和模型层共同决定了数据怎样被页面消费。"
                "DAO 负责向数据库取数，模型负责把原始结果整理成页面直接可用的对象结构。"
                "当这两层职责清晰时，页面就可以专注于显示，不需要承担复杂的数据拼装工作。"
            ),
            image_path=crops.get("class_dao"),
            code_title="模型测试代码片段",
            code_text=model_code,
        ),
        SlideSpec(
            title="知识图谱核心流程总览",
            subtitle="先看清数据从加载到交互落库的完整链路",
            bullets=[
                "图谱数据加载",
                "列表显示",
                "节点与边生成",
                "图谱展开浏览",
                "收藏与学习记录",
                "联动资源与进度",
            ],
            narration=(
                "知识图谱的核心流程可以分成六个连续环节。"
                "首先，系统从数据库加载图谱基础数据。"
                "然后，在图谱列表页显示不同主题，供用户选择进入。"
                "接着，详情页继续读取节点和边，生成图谱结构对象。"
                "随后，页面计算布局并支持展开、缩放和拖拽浏览。"
                "当用户点击节点后，可以查看详情、开始学习或加入收藏。"
                "最后，学习记录、收藏结果以及相关资源联动都会写回系统。"
            ),
            image_path=crops.get("sequence_full"),
            voice_segments=[
                "知识图谱的核心流程可以分成六个连续环节。",
                "首先，系统从数据库加载图谱基础数据。",
                "然后，在图谱列表页显示不同主题，供用户选择进入。",
                "接着，详情页继续读取节点和边，生成图谱结构对象。",
                "随后，页面计算布局并支持展开、缩放和拖拽浏览。",
                "当用户点击节点后，可以查看详情、开始学习或加入收藏。",
                "最后，学习记录、收藏结果以及相关资源联动都会写回系统。",
            ],
            summary_text=(
                "核心链路：\n"
                "数据加载 → 列表显示 → 节点与边生成 → 图谱展开浏览 → 收藏 / 学习记录 → 联动资源"
            ),
        ),
        SlideSpec(
            title="图谱数据加载与列表显示",
            subtitle="先把图谱主题取出来，再形成列表入口",
            bullets=[
                "GraphDao 查询图谱基础数据",
                "GraphListPage 接收图谱列表",
                "页面按主题显示图谱卡片",
                "用户点击卡片进入详情页",
            ],
            narration=(
                "图谱功能的入口从数据加载开始。"
                "GraphDao 先查询图谱基础数据，GraphListPage 再把这些结果组织成可点击的图谱卡片列表。"
                "用户看到的是课程图谱、技术栈图谱等主题入口，点击其中一个卡片之后，系统才会进入图谱详情流程。"
            ),
            image_path=runtime_shots.get("graph_list"),
            summary_text=(
                "入口阶段：\n"
                "数据库读取图谱基础数据。\n"
                "GraphListPage 负责把结果显示成图谱列表入口。"
            ),
        ),
        SlideSpec(
            title="节点与边生成流程",
            subtitle="图谱详情页把原始数据转换成可视化结构",
            bullets=[
                "根据 graphId 读取 nodes",
                "根据 graphId 读取 edges",
                "生成 NodeModel / EdgeModel",
                "计算布局后渲染图谱",
            ],
            narration=(
                "进入图谱详情页后，系统会根据当前图谱标识继续读取节点和边。"
                "这些原始数据会被转换成 NodeModel 和 EdgeModel，再交给页面进行布局计算。"
                "完成这一层处理之后，用户看到的才是可视化图谱，而不是数据库里的原始记录。"
            ),
            image_path=crops.get("sequence_detail"),
            summary_text=(
                "详情阶段：\n"
                "读取 nodes 和 edges。\n"
                "生成模型对象。\n"
                "布局完成后形成可视化图谱。"
            ),
        ),
        SlideSpec(
            title="图谱展开与收藏流程",
            subtitle="节点浏览、展开、学习与收藏都发生在详情页",
            bullets=[
                "点击节点查看详情",
                "支持拖拽、缩放与展开浏览",
                "开始学习写入 learning_records",
                "收藏写入 favorites",
            ],
            narration=(
                "图谱详情页不仅负责显示，还负责交互。"
                "用户可以拖拽、缩放并展开浏览节点关系，点击某个节点后会看到对应详情。"
                "如果用户点击开始学习，系统就把学习行为写入 learning_records。"
                "如果用户点击收藏，系统就把当前节点信息写入 favorites。"
            ),
            image_path=crops.get("sequence_action"),
            summary_text=(
                "交互阶段：\n"
                "展开浏览节点关系。\n"
                "开始学习写入 learning_records。\n"
                "收藏写入 favorites。"
            ),
        ),
        SlideSpec(
            title="开发过程：前期阶段",
            subtitle="需求分析、业务建模与数据库设计",
            bullets=[
                "明确图谱是系统核心",
                "定义 Graph / Node / Edge / Favorite / LearningRecord",
                "设计 SQLite 表结构",
            ],
            narration=(
                "开发前期的重点是需求分析和业务建模。"
                "首先明确知识图谱是系统核心，然后定义图谱、节点、边、收藏和学习记录等对象。"
                "最后再把这些对象映射成 SQLite 表结构，为后续编码打下基础。"
            ),
            image_path=crops.get("process_early"),
        ),
        SlideSpec(
            title="开发过程：中期阶段",
            subtitle="DAO 开发、页面实现与图谱交互构建",
            bullets=[
                "实现 DatabaseHelper",
                "实现 GraphDao / FavoriteDao / LearningRecordDao",
                "实现 GraphListPage / GraphDetailPage",
                "构建图谱布局和节点交互",
            ],
            narration=(
                "开发中期进入真正的实现阶段。"
                "这一阶段先完成数据库和 DAO，再实现图谱列表页和图谱详情页。"
                "当数据层、页面层和交互层被连起来之后，图谱功能才真正具备使用价值。"
            ),
            image_path=crops.get("process_middle"),
        ),
        SlideSpec(
            title="开发过程：后期阶段",
            subtitle="测试验证、构建发布与教学材料交付",
            bullets=[
                "补充模型测试与组件测试",
                "生成测试用例与测试报告",
                "构建 Debug / Release APK",
                "生成 UML、PPT、视频材料",
            ],
            narration=(
                "开发后期的关键是验证与交付。"
                "项目补充了模型测试和组件测试，整理了测试用例与测试报告，并完成 Debug 和 Release APK 构建。"
                "在此基础上，再输出 UML 图、PPT 和教学视频，形成完整交付材料。"
            ),
            image_path=crops.get("process_late"),
        ),
        SlideSpec(
            title="测试资产讲解",
            subtitle="从测试代码与测试用例验证系统质量",
            bullets=[
                "model_test.dart",
                "widget_test.dart",
                "home_page_widget_test.dart",
                "test_cases.md",
                "23 项测试通过",
            ],
            narration=(
                "项目不仅有页面和结构图，也有实际测试资产。"
                "模型测试覆盖 GraphModel、NodeModel、EdgeModel 等核心数据对象，组件测试覆盖登录页和首页导航。"
                "测试用例文档则进一步补充了图谱、资料、视频、测验和进度等模块的验证点。"
            ),
            image_path=runtime_shots.get("login"),
            code_title="测试用例摘要",
            code_text=test_case_summary,
        ),
        SlideSpec(
            title="测试结果与交付成果",
            subtitle="理解测试、构建与展示材料的闭环关系",
            bullets=[
                "23 项测试通过",
                "Debug APK 构建成功",
                "Release APK 构建成功",
                "测试报告已整理",
                "视频 / PPT / UML 已生成",
            ],
            narration=(
                "测试结果说明当前工程已经具备较完整的交付状态。"
                "目前模型层测试和组件测试累计通过二十三项，同时 Debug 和 Release APK 也已经构建成功。"
                "这使得项目不仅能运行，也能以测试报告、PPT 和视频的形式稳定展示。"
            ),
            image_path=runtime_shots.get("graph_list"),
            code_title="测试报告摘要",
            code_text=test_report_summary,
        ),
        SlideSpec(
            title="实际运行界面讲解",
            subtitle="把结构图理解与真实页面效果对应起来",
            bullets=[
                "登录页：统一入口",
                "首页：导航中心",
                "图谱列表页：图谱主题选择",
                "图谱详情页：节点交互与学习动作",
            ],
            narration=(
                "最后回到运行界面。"
                "登录页负责统一入口，首页负责导航分发，图谱列表页负责主题选择，图谱详情页负责节点交互和学习动作。"
                "把结构图、类图、顺序图和真实页面放在一起理解，系统的整体逻辑就会更清楚。"
            ),
            image_path=runtime_shots.get("graph_detail"),
        ),
        SlideSpec(
            title="本节回顾",
            subtitle="回顾本节知识点与整体实现逻辑",
            bullets=[
                "认识 Flutter 与 Dart 分工",
                "理解三层结构与核心类",
                "看清知识图谱核心流程",
                "完成测试与交付闭环",
            ],
            narration=(
                "这一节的重点有四个。"
                "第一，理解 Flutter 负责页面，Dart 负责逻辑。"
                "第二，理解页面层、业务层和数据层的结构关系。"
                "第三，看清知识图谱从加载到交互再到记录写入的完整流程。"
                "第四，看到测试结果和交付材料如何支撑项目展示。"
            ),
            image_path=crops.get("framework_right"),
            voice_segments=[
                "这一节的重点有四个。",
                "第一，理解 Flutter 负责页面，Dart 负责逻辑。",
                "第二，理解页面层、业务层和数据层的结构关系。",
                "第三，看清知识图谱从加载到交互再到记录写入的完整流程。",
                "第四，看到测试结果和交付材料如何支撑项目展示。",
            ],
        ),
        SlideSpec(
            title="课程总结与下节预告",
            subtitle="本节回顾与后续学习方向说明",
            bullets=[
                "先回看系统框架与类图",
                "再回看核心流程链路",
                "继续进入数据库初始化实现",
                "继续进入节点布局与交互源码阅读",
            ],
            narration=(
                "本节内容到这里结束。"
                "建议你先回看系统框架、类图和知识图谱核心流程，重新整理整体认识。"
                "下一步可以继续学习数据库初始化、图谱数据装载、节点布局计算以及交互源码实现，从展示理解逐步进入实现理解。"
            ),
            image_path=crops.get("framework_right"),
        ),
    ]


def paste_image(
    base: Image.Image, source: Path | None, box: tuple[int, int, int, int]
) -> None:
    if source is None or not source.exists():
        return
    img = Image.open(source).convert("RGB")
    target_w = box[2] - box[0]
    target_h = box[3] - box[1]
    ratio = min(target_w / img.width, target_h / img.height)
    new_size = (max(1, int(img.width * ratio)), max(1, int(img.height * ratio)))
    img = img.resize(new_size)
    x = box[0] + (target_w - new_size[0]) // 2
    y = box[1] + (target_h - new_size[1]) // 2
    base.paste(img, (x, y))


def draw_bullets(
    draw: ImageDraw.ImageDraw, bullets: list[str], box: tuple[int, int, int, int]
) -> None:
    x1, y1, x2, y2 = box
    y = y1
    for bullet in bullets:
        rounded_rect(draw, (x1, y + 8, x1 + 18, y + 26), PRIMARY, PRIMARY, 6, 1)
        y = (
            draw_wrapped_text(
                draw, bullet, (x1 + 32, y, x2, y + 66), FONT_BODY, TEXT, 8
            )
            + 12
        )
        if y > y2:
            break


def draw_code_block(
    draw: ImageDraw.ImageDraw, title: str, code: str, box: tuple[int, int, int, int]
) -> None:
    rounded_rect(draw, box, CODE_BG, (224, 230, 242), 18, 2)
    x1, y1, _x2, y2 = box
    draw.text((x1 + 18, y1 + 14), title, font=FONT_SMALL_BOLD, fill=PRIMARY_DARK)
    y = y1 + 52
    for line in code.splitlines():
        if y > y2 - 24:
            break
        draw.text((x1 + 18, y), line.rstrip(), font=FONT_CODE, fill=TEXT)
        _, h = text_size(draw, line or "中", FONT_CODE)
        y += h + 5


def split_text_segments(text: str, max_chars: int = 22) -> list[str]:
    source = text.strip()
    chunks: list[str] = []
    current = ""

    for ch in source:
        current += ch
        if ch in "，。！？；":
            trimmed = current.strip()
            if trimmed:
                chunks.append(trimmed)
            current = ""

    tail = current.strip()
    if tail:
        chunks.append(tail)

    cleaned_chunks = [
        chunk
        for chunk in chunks
        if any(ch.isalnum() or "\u4e00" <= ch <= "\u9fff" for ch in chunk)
    ]
    if not cleaned_chunks:
        cleaned_chunks = [source]

    merged: list[str] = []
    current_line = ""
    for chunk in cleaned_chunks:
        candidate = current_line + chunk
        if len(candidate) <= max_chars or not current_line:
            current_line = candidate
        else:
            merged.append(current_line)
            current_line = chunk
    if current_line:
        merged.append(current_line)

    return merged


def get_voice_segments(slide: SlideSpec) -> list[str]:
    if slide.voice_segments:
        return [seg.strip() for seg in slide.voice_segments if seg.strip()]
    return split_text_segments(slide.narration.strip(), 22)


def summary_text_for_slide(slide: SlideSpec) -> str:
    if slide.summary_text:
        return slide.summary_text
    return "\n".join(get_voice_segments(slide)[:5])


def render_slide(slide: SlideSpec, index: int, total: int) -> Path:
    img = Image.new("RGB", (WIDTH, HEIGHT), WHITE)
    draw_vertical_gradient(img, BG_TOP, BG_BOTTOM)
    draw = ImageDraw.Draw(img)

    rounded_rect(draw, (48, 34, 1870, 150), WHITE, BORDER, 28, 2)
    draw.text((84, 50), slide.title, font=FONT_TITLE, fill=TEXT)
    draw.text((88, 108), slide.subtitle, font=FONT_SUBTITLE, fill=TEXT_LIGHT)

    rounded_rect(draw, (1660, 54, 1835, 108), PRIMARY, PRIMARY, 18, 1)
    center_text(draw, 1748, 82, f"{index}/{total}", FONT_SMALL_BOLD, WHITE)

    left_box = (56, 180, 650, 930)
    right_box = (680, 180, 1862, 785)
    bottom_box = (680, 810, 1862, 948)

    rounded_rect(draw, left_box, WHITE, BORDER, 30, 3)
    rounded_rect(draw, right_box, WHITE, BORDER, 30, 3)
    rounded_rect(draw, bottom_box, WHITE, BORDER, 30, 3)

    draw_bullets(draw, slide.bullets, (96, 216, 595, 900))

    _ = paste_image(img, slide.image_path, (700, 205, 1840, 760))
    if slide.image_caption:
        _ = draw_wrapped_text(
            draw,
            slide.image_caption,
            (720, 742, 1820, 778),
            FONT_SMALL,
            TEXT_LIGHT,
            2,
        )

    _ = draw_wrapped_text(
        draw,
        summary_text_for_slide(slide),
        (710, 836, 1834, 930),
        FONT_BODY,
        TEXT,
        6,
    )

    out_path = SLIDES_DIR / f"{index:02d}_{safe_name(slide.title)}.png"
    img.save(out_path)
    return out_path


def render_all_slides(slides: list[SlideSpec]) -> list[Path]:
    return [
        render_slide(slide, i, len(slides)) for i, slide in enumerate(slides, start=1)
    ]


def choose_voice(engine: TTSEngineLike, preferred_keywords: list[str]) -> str | None:
    voices_obj = engine.getProperty("voices")
    voices = cast(list[VoiceLike], voices_obj) if isinstance(voices_obj, list) else []
    lowered = [key.lower() for key in preferred_keywords]
    for voice in voices:
        combined = f"{voice.id} {voice.name}".lower()
        if any(key in combined for key in lowered):
            return voice.id
    return voices[0].id if voices else None


def clean_tts_text(text: str) -> str:
    replacements = {
        "Flutter": "Flutter",
        "Dart": "Dart",
        "SQLite": "S Q Lite",
        "GraphDao": "Graph Dao",
        "FavoriteDao": "Favorite Dao",
        "LearningRecordDao": "Learning Record Dao",
        "DatabaseHelper": "Database Helper",
        "GraphListPage": "Graph List Page",
        "GraphDetailPage": "Graph Detail Page",
        "HomePage": "Home Page",
        "QuizPage": "Quiz Page",
        "ProgressPage": "Progress Page",
        "LoginPage": "Login Page",
        "UML": "U M L",
        "APK": "A P K",
        "favorites": "favorites",
        "learning_records": "learning records",
    }
    result = text
    for old, new in replacements.items():
        result = result.replace(old, new)
    return result


def save_tts_audio(text: str, out_path: Path, rate: int = 158) -> bool:
    if pyttsx3 is None:
        return False
    try:
        engine = cast(TTSEngineLike, pyttsx3.init())
        voice_id = choose_voice(
            engine, ["zh", "chinese", "huihui", "hanhan", "xiaoxiao"]
        )
        if voice_id:
            _ = engine.setProperty("voice", voice_id)
        _ = engine.setProperty("rate", rate)
        _ = engine.save_to_file(clean_tts_text(text), str(out_path))
        _ = engine.runAndWait()
        _ = engine.stop()
        return out_path.exists() and out_path.stat().st_size > 0
    except Exception:
        return False


def create_silent_wav(duration_seconds: float, out_path: Path) -> None:
    sample_rate = 22050
    frames = max(1, int(duration_seconds * sample_rate))
    silence = struct.pack("<h", 0)
    with wave.open(str(out_path), "wb") as wav_file:
        wav_file.setnchannels(1)
        wav_file.setsampwidth(2)
        wav_file.setframerate(sample_rate)
        for _ in range(frames):
            wav_file.writeframesraw(silence)


def wav_duration(path: Path) -> float:
    with wave.open(str(path), "rb") as wav_file:
        return wav_file.getnframes() / float(wav_file.getframerate())


def combine_wavs(inputs: list[Path], output_path: Path) -> bool:
    if not inputs:
        return False

    first_params = None
    with wave.open(str(output_path), "wb") as out_wav:
        for idx, wav_path in enumerate(inputs):
            if not wav_path.exists():
                return False
            with wave.open(str(wav_path), "rb") as in_wav:
                params = (
                    in_wav.getnchannels(),
                    in_wav.getsampwidth(),
                    in_wav.getframerate(),
                )
                if idx == 0:
                    first_params = params
                    out_wav.setnchannels(params[0])
                    out_wav.setsampwidth(params[1])
                    out_wav.setframerate(params[2])
                elif params != first_params:
                    return False
                out_wav.writeframes(in_wav.readframes(in_wav.getnframes()))
    return output_path.exists() and output_path.stat().st_size > 0


def which(names: list[str]) -> str | None:
    for name in names:
        found = shutil.which(name)
        if found:
            return found
    return None


def ffmpeg_path() -> str | None:
    return which(
        ["ffmpeg", r"C:\Users\ldl\AppData\Local\Microsoft\WinGet\Links\ffmpeg.exe"]
    )


def run_ffmpeg(args: list[str]) -> bool:
    ffmpeg = ffmpeg_path()
    if not ffmpeg:
        return False
    process = subprocess.run(
        [ffmpeg] + args, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False
    )
    return process.returncode == 0


def build_clip(image_path: Path, audio_path: Path, output_path: Path) -> bool:
    return run_ffmpeg(
        [
            "-y",
            "-loop",
            "1",
            "-i",
            str(image_path),
            "-i",
            str(audio_path),
            "-c:v",
            "libx264",
            "-tune",
            "stillimage",
            "-c:a",
            "pcm_s16le",
            "-pix_fmt",
            "yuv420p",
            "-shortest",
            "-vf",
            "scale=1920:1080,fps=30",
            str(output_path),
        ]
    )


def concat_clips(clips: list[Path], output_path: Path) -> bool:
    ffmpeg = ffmpeg_path()
    if not ffmpeg:
        return False
    manifest = TEMP_DIR / "clips_manifest_v5.txt"
    _ = manifest.write_text(
        "\n".join(f"file '{clip.as_posix()}'" for clip in clips), encoding="utf-8"
    )
    process = subprocess.run(
        [
            ffmpeg,
            "-y",
            "-f",
            "concat",
            "-safe",
            "0",
            "-i",
            str(manifest),
            "-c",
            "copy",
            str(output_path),
        ],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
    )
    return process.returncode == 0 and output_path.exists()


def seconds_to_srt(ts: float) -> str:
    hours = int(ts // 3600)
    minutes = int((ts % 3600) // 60)
    seconds = int(ts % 60)
    millis = int((ts - int(ts)) * 1000)
    return f"{hours:02d}:{minutes:02d}:{seconds:02d},{millis:03d}"


def build_slide_audio(
    slide: SlideSpec, slide_index: int
) -> tuple[Path, list[tuple[str, float]]]:
    segments = get_voice_segments(slide)
    segment_paths: list[Path] = []
    segment_info: list[tuple[str, float]] = []

    for seg_index, segment in enumerate(segments, start=1):
        seg_path = SEGMENT_AUDIO_DIR / f"{slide_index:02d}_{seg_index:02d}.wav"
        ok = save_tts_audio(segment, seg_path)
        if not ok:
            create_silent_wav(1.5, seg_path)
        duration = max(0.1, wav_duration(seg_path))
        segment_paths.append(seg_path)
        segment_info.append((segment, duration))

    slide_audio_path = AUDIO_DIR / f"{slide_index:02d}_{safe_name(slide.title)}.wav"
    if not combine_wavs(segment_paths, slide_audio_path):
        create_silent_wav(
            max(3.0, sum(duration for _, duration in segment_info)), slide_audio_path
        )

    return slide_audio_path, segment_info


def generate_srt_from_entries(entries: list[tuple[str, float, float]]) -> None:
    lines: list[str] = []
    for i, (text, start, end) in enumerate(entries, start=1):
        lines.append(str(i))
        lines.append(f"{seconds_to_srt(start)} --> {seconds_to_srt(end)}")
        lines.append(text)
        lines.append("")
    _ = SRT_PATH.write_text("\n".join(lines), encoding="utf-8")


def build_video(slides: list[SlideSpec], images: list[Path]) -> bool:
    clips: list[Path] = []
    srt_entries: list[tuple[str, float, float]] = []
    current_time = 0.0

    for idx, (slide, image_path) in enumerate(zip(slides, images), start=1):
        audio_path, segment_info = build_slide_audio(slide, idx)
        clip_path = CLIPS_DIR / f"{idx:02d}_{safe_name(slide.title)}.mkv"

        segment_start = current_time
        slide_duration = 0.0
        for segment_text, segment_duration in segment_info:
            segment_end = segment_start + segment_duration
            srt_entries.append((segment_text, segment_start, segment_end))
            segment_start = segment_end
            slide_duration += segment_duration

        if slide_duration <= 0:
            slide_duration = wav_duration(audio_path)

        if not build_clip(image_path, audio_path, clip_path):
            return False

        clips.append(clip_path)
        current_time += slide_duration

    generate_srt_from_entries(srt_entries)

    plain_video_path = TEMP_DIR / "knowledge_graph_video_v5_plain.mkv"
    if not concat_clips(clips, plain_video_path):
        return False

    subtitle_path = (
        SRT_PATH.resolve().as_posix().replace(":", "\\:").replace(",", "\\,")
    )
    return run_ffmpeg(
        [
            "-y",
            "-i",
            str(plain_video_path),
            "-vf",
            f"subtitles='{subtitle_path}':force_style='FontName=Microsoft YaHei,FontSize=20,PrimaryColour=&H00FFFFFF,OutlineColour=&H00000000,BorderStyle=3,Outline=1,Shadow=0,MarginV=24'",
            "-c:v",
            "libx264",
            "-c:a",
            "aac",
            "-b:a",
            "192k",
            str(VIDEO_PATH),
        ]
    )


def build_pptx(images: list[Path], slides: list[SlideSpec]) -> None:
    if Presentation is None or Inches is None or Pt is None:
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for img_path, slide_spec in zip(images, slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(
            Inches(0.45), Inches(0.18), Inches(9.8), Inches(0.5)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = slide_spec.title
        p.font.size = Pt(28)
        p.font.bold = True

        subtitle_box = slide.shapes.add_textbox(
            Inches(0.48), Inches(0.70), Inches(10.5), Inches(0.35)
        )
        p2 = subtitle_box.text_frame.paragraphs[0]
        p2.text = slide_spec.subtitle
        p2.font.size = Pt(15)

        _ = slide.shapes.add_picture(
            str(img_path),
            Inches(0.36),
            Inches(1.12),
            width=Inches(12.55),
            height=Inches(6.0),
        )

    prs.save(str(PPTX_PATH))


def print_summary(
    images: list[Path],
    video_ok: bool,
    crops: dict[str, Path],
    runtime_shots: dict[str, Path],
) -> None:
    print("已生成 v5 视频素材：")
    print(f"- {SCRIPT_PATH}")
    print(f"- {SRT_PATH}")
    print(f"- {PPTX_PATH}")
    for key, value in sorted(crops.items()):
        print(f"- {key}: {value}")
    for key, value in sorted(runtime_shots.items()):
        print(f"- {key}: {value}")
    for img in images:
        print(f"- {img}")
    if video_ok:
        print(f"- {VIDEO_PATH}")
    else:
        print("- 视频未成功生成，请检查 ffmpeg 或音频环境。")


def main() -> None:
    ensure_dirs()

    required = [FRAMEWORK_IMG, CLASS_IMG, SEQUENCE_IMG, PROCESS_IMG]
    missing = [path for path in required if not path.exists()]
    if missing:
        print("缺少以下 UML 图片文件：")
        for path in missing:
            print(f"- {path}")
        raise SystemExit(1)

    official_assets = ensure_official_assets()
    crops = generate_uml_crops()
    runtime_shots = generate_mock_runtime_screenshots()
    slides = build_slides(crops, runtime_shots, official_assets)
    build_script_markdown(slides)
    images = render_all_slides(slides)
    build_pptx(images, slides)
    video_ok = build_video(slides, images)
    print_summary(images, video_ok, crops, runtime_shots)


if __name__ == "__main__":
    main()
