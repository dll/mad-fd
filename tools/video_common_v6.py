#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
video_common_v6.py — 共享视频生成引擎

所有功能教学视频共用此模块，包含：
  - 颜色/字体常量
  - 绘图工具函数
  - edge_tts + pyttsx3 双引擎 TTS
  - moviepy 帧级同步视频生成
  - UML 裁切图生成
  - Mock 运行截图生成
  - SRT / PPTX / 脚本输出
"""

from __future__ import annotations

import asyncio
import os
import struct
import sys
import time
import urllib.request
import wave
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont

# ── edge_tts ──────────────────────────────────────────────────────────────
try:
    import edge_tts

    _EDGE_OK = True
except ImportError:
    _EDGE_OK = False
    print("Warning: edge-tts not found.  pip install edge-tts")

# ── moviepy ───────────────────────────────────────────────────────────────
try:
    from moviepy.editor import AudioFileClip, ImageClip, concatenate_videoclips

    _MOVIEPY_OK = True
except ImportError:
    _MOVIEPY_OK = False
    print("Error: moviepy not found.  pip install moviepy==1.0.3")

# ── pyttsx3 fallback ──────────────────────────────────────────────────────
try:
    import pyttsx3 as _pyttsx3_mod  # type: ignore[import]

    _PYTTSX3_OK = True
except Exception:
    _pyttsx3_mod = None
    _PYTTSX3_OK = False

# ── python-pptx ───────────────────────────────────────────────────────────
try:
    from pptx import Presentation  # type: ignore[import]
    from pptx.util import Inches, Pt  # type: ignore[import]

    _PPTX_OK = True
except Exception:
    Presentation = Inches = Pt = None
    _PPTX_OK = False


# ═══════════════════════════════════════════════════════════════════════════
# PROJECT ROOT
# ═══════════════════════════════════════════════════════════════════════════
ROOT = Path(__file__).resolve().parents[1]
DIAGRAMS_DIR = ROOT / "docs" / "diagrams" / "v3"
TESTING_DIR = ROOT / "docs" / "testing"
OUT_DIR = ROOT / "video_output"

FRAMEWORK_IMG = DIAGRAMS_DIR / "flutter_dart_framework_architecture.png"
CLASS_IMG = DIAGRAMS_DIR / "flutter_dart_core_class_diagram.png"
SEQUENCE_IMG = DIAGRAMS_DIR / "graph_feature_sequence_diagram.png"
PROCESS_IMG = DIAGRAMS_DIR / "knowledge_graph_development_process.png"

LOGIN_GOLDEN = ROOT / "test" / "screenshots" / "goldens" / "login_page.png"
HOME_GOLDEN = ROOT / "test" / "screenshots" / "goldens" / "home_page.png"


# ═══════════════════════════════════════════════════════════════════════════
# VIDEO PARAMETERS
# ═══════════════════════════════════════════════════════════════════════════
W, H = 1920, 1080
FPS = 30
TTS_VOICE = "zh-CN-XiaoxiaoNeural"
TTS_RATE = "-5%"
CLIP_TAIL = 0.40  # seconds held after speech for mid-slide sentences
SLIDE_TAIL = 1.20  # extra tail for last sentence of each slide
SUB_H = 115  # subtitle bar height (px)
CONTENT_H = H - SUB_H  # 965 px – main content area
TITLE_H = 84  # title bar height


# ═══════════════════════════════════════════════════════════════════════════
# COLOURS
# ═══════════════════════════════════════════════════════════════════════════
BG_TOP = (20, 36, 68)
BG_BOT = (32, 54, 96)
TITLE_BG = (13, 24, 50)
ACCENT = (255, 198, 55)
TEXT = (236, 241, 252)
TEXT_MUTE = (165, 185, 220)
PRIMARY = (92, 172, 255)
GREEN = (72, 208, 126)
ORANGE = (255, 152, 55)
PURPLE = (168, 110, 255)
RED_C = (235, 85, 85)
BORDER = (52, 84, 140)
CODE_BG = (24, 30, 44)
PANEL = (26, 44, 76)
SUB_BG = (6, 6, 12)
WHITE = (255, 255, 255)


# ═══════════════════════════════════════════════════════════════════════════
# FONTS
# ═══════════════════════════════════════════════════════════════════════════
def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    paths: list[str] = (
        [r"C:\Windows\Fonts\msyhbd.ttc", r"C:\Windows\Fonts\simhei.ttf"]
        if bold
        else [r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\simsun.ttc"]
    ) + [r"C:\Windows\Fonts\arial.ttf"]
    for p in paths:
        try:
            return ImageFont.truetype(p, size)
        except Exception:
            continue
    return ImageFont.load_default()  # type: ignore[return-value]


FT = _font(48, True)  # slide title
FS = _font(24, False)  # slide subtitle
FB = _font(22, True)  # section label
FN = _font(22, False)  # body / bullets
FK = _font(18, False)  # code / small
FKB = _font(18, True)  # code bold / badge
FSU = _font(27, False)  # subtitle bar text


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE SPEC
# ═══════════════════════════════════════════════════════════════════════════
@dataclass
class SlideSpec:
    title: str
    subtitle: str
    bullets: list[str]
    narration: str
    image_path: Path | None = None
    image_caption: str | None = None
    code_title: str | None = None
    code_text: str | None = None
    voice_segments: list[str] | None = None


# ═══════════════════════════════════════════════════════════════════════════
# DRAWING UTILITIES
# ═══════════════════════════════════════════════════════════════════════════
def tsz(draw: ImageDraw.ImageDraw, text: str, font: Any) -> tuple[int, int]:
    bb = draw.textbbox((0, 0), text, font=font)
    return int(bb[2] - bb[0]), int(bb[3] - bb[1])


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: Any, max_w: int) -> list[str]:
    lines: list[str] = []
    for para in text.splitlines():
        cur = ""
        for ch in para.strip():
            cand = cur + ch
            if tsz(draw, cand, font)[0] > max_w and cur:
                lines.append(cur)
                cur = ch
            else:
                cur = cand
        if cur:
            lines.append(cur)
    return lines


def draw_wrapped(
    draw: ImageDraw.ImageDraw,
    text: str,
    x: int,
    y: int,
    font: Any,
    fill: tuple[int, int, int],
    max_w: int,
    gap: int = 5,
) -> int:
    for line in wrap_text(draw, text, font, max_w):
        draw.text((x, y), line, font=font, fill=fill)
        y += tsz(draw, line or "M", font)[1] + gap
    return y


def gradient_fill(img: Image.Image) -> None:
    draw = ImageDraw.Draw(img)
    for y in range(img.height):
        r = BG_TOP[0] + int((BG_BOT[0] - BG_TOP[0]) * y / img.height)
        g = BG_TOP[1] + int((BG_BOT[1] - BG_TOP[1]) * y / img.height)
        b = BG_TOP[2] + int((BG_BOT[2] - BG_TOP[2]) * y / img.height)
        draw.line([(0, y), (img.width, y)], fill=(r, g, b))


def paste_fit(
    base: Image.Image,
    src: Path | None,
    box: tuple[int, int, int, int],
) -> None:
    if not src or not src.exists():
        return
    try:
        im = Image.open(src).convert("RGB")
    except Exception:
        return
    bw, bh = box[2] - box[0], box[3] - box[1]
    ratio = min(bw / im.width, bh / im.height)
    nw = max(1, int(im.width * ratio))
    nh = max(1, int(im.height * ratio))
    im = im.resize((nw, nh), Image.LANCZOS)  # type: ignore[attr-defined]
    base.paste(im, (box[0] + (bw - nw) // 2, box[1] + (bh - nh) // 2))


def draw_bullets(
    draw: ImageDraw.ImageDraw,
    bullets: list[str],
    x1: int,
    y1: int,
    x2: int,
    y2: int,
) -> None:
    y = y1
    for bullet in bullets:
        if y + 28 > y2:
            break
        draw.ellipse((x1, y + 7, x1 + 13, y + 20), fill=ACCENT)
        y = draw_wrapped(draw, bullet, x1 + 22, y, FN, TEXT, x2 - x1 - 26, 5) + 10


def draw_code_box(
    draw: ImageDraw.ImageDraw,
    title: str,
    code: str,
    x1: int,
    y1: int,
    x2: int,
    y2: int,
) -> None:
    draw.rounded_rectangle(
        (x1, y1, x2, y2), radius=12, fill=CODE_BG, outline=BORDER, width=2
    )
    draw.text((x1 + 14, y1 + 10), title, font=FKB, fill=PRIMARY)
    cy = y1 + 38
    for line in code.splitlines():
        if cy > y2 - 16:
            break
        draw.text((x1 + 14, cy), line.rstrip(), font=FK, fill=TEXT)
        cy += tsz(draw, line or "M", FK)[1] + 4


# ═══════════════════════════════════════════════════════════════════════════
# SAFE FILENAME HELPER
# ═══════════════════════════════════════════════════════════════════════════
def safe_name(text: str) -> str:
    keep: list[str] = []
    for ch in text:
        if ch.isalnum() or ch in ("-", "_") or "\u4e00" <= ch <= "\u9fff":
            keep.append(ch)
        else:
            keep.append("_")
    return "".join(keep).strip("_") or "slide"


# ═══════════════════════════════════════════════════════════════════════════
# TEXT FILE UTILITIES
# ═══════════════════════════════════════════════════════════════════════════
def read_file(path: Path) -> str:
    return path.read_text(encoding="utf-8") if path.exists() else ""


def extract_block(path: Path, marker: str, n: int = 20) -> str:
    text = read_file(path)
    for i, line in enumerate(text.splitlines()):
        if marker in line:
            return "\n".join(text.splitlines()[i : i + n])
    return "\n".join(text.splitlines()[:n])


def select_lines(text: str, keywords: list[str], n: int = 10) -> str:
    picked: list[str] = []
    for line in text.splitlines():
        s = line.strip()
        if s and any(k.lower() in s.lower() for k in keywords):
            picked.append(s)
            if len(picked) >= n:
                break
    return "\n".join(picked) if picked else "\n".join(text.splitlines()[:n])


# ═══════════════════════════════════════════════════════════════════════════
# SENTENCE SPLITTING
# ═══════════════════════════════════════════════════════════════════════════
def get_sentences(slide: SlideSpec) -> list[str]:
    if slide.voice_segments:
        return [s.strip() for s in slide.voice_segments if s.strip()]
    text = slide.narration.strip()
    chunks: list[str] = []
    cur = ""
    for ch in text:
        cur += ch
        if ch in "。！？；":
            t = cur.strip()
            if t:
                chunks.append(t)
            cur = ""
    if cur.strip():
        chunks.append(cur.strip())
    # merge very short fragments with next
    merged: list[str] = []
    i = 0
    while i < len(chunks):
        if i + 1 < len(chunks) and len(chunks[i]) < 8:
            merged.append(chunks[i] + chunks[i + 1])
            i += 2
        else:
            merged.append(chunks[i])
            i += 1
    return merged if merged else [text]


# ═══════════════════════════════════════════════════════════════════════════
# TTS ENGINE  (edge_tts primary → pyttsx3 fallback)
# ═══════════════════════════════════════════════════════════════════════════
async def _edge_async(text: str, path: Path) -> None:
    comm = edge_tts.Communicate(text, TTS_VOICE, rate=TTS_RATE)
    await comm.save(str(path))


_pyttsx3_engine: Any = None


def _init_pyttsx3() -> None:
    global _pyttsx3_engine
    if _pyttsx3_engine is not None or not _PYTTSX3_OK:
        return
    try:
        eng = _pyttsx3_mod.init()
        eng.setProperty("rate", 155)
        voices = eng.getProperty("voices") or []
        for v in voices:
            combined = f"{getattr(v, 'name', '')} {getattr(v, 'id', '')}".lower()
            if any(k in combined for k in ["zh", "chinese", "huihui", "xiaoxiao"]):
                eng.setProperty("voice", v.id)
                break
        _pyttsx3_engine = eng
    except Exception as e:
        print(f"  pyttsx3 init error: {e}")


def _pyttsx3_gen(text: str, wav_path: Path) -> bool:
    _init_pyttsx3()
    if _pyttsx3_engine is None:
        return False
    try:
        _pyttsx3_engine.save_to_file(text, str(wav_path))
        _pyttsx3_engine.runAndWait()
        time.sleep(0.3)
        return wav_path.exists() and wav_path.stat().st_size > 200
    except Exception:
        return False


def generate_tts(text: str, stem: Path) -> Path | None:
    """Generate TTS; caches by file existence. Returns audio path or None."""
    mp3 = stem.with_suffix(".mp3")
    wav = stem.with_suffix(".wav")

    if mp3.exists() and mp3.stat().st_size > 500:
        return mp3
    if wav.exists() and wav.stat().st_size > 500:
        return wav

    if _EDGE_OK:
        for attempt in range(3):
            try:
                asyncio.run(_edge_async(text, mp3))
                if mp3.exists() and mp3.stat().st_size > 500:
                    return mp3
                mp3.unlink(missing_ok=True)
            except Exception as e:
                print(f"    edge_tts attempt {attempt + 1}: {e}")
                mp3.unlink(missing_ok=True)
                if attempt < 2:
                    time.sleep(1.5 * (attempt + 1))
        print("    edge_tts failed, switching to pyttsx3")

    if _pyttsx3_gen(text, wav):
        return wav
    return None


def get_audio_duration(audio_path: Path) -> float:
    try:
        clip = AudioFileClip(str(audio_path))
        dur = clip.duration
        clip.close()
        return dur
    except Exception:
        return 0.0


def make_silent_wav(duration: float, path: Path) -> None:
    sr = 22050
    n = max(1, int(duration * sr))
    with wave.open(str(path), "wb") as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(sr)
        wf.writeframes(b"\x00\x00" * n)


# ═══════════════════════════════════════════════════════════════════════════
# UML CROPS
# ═══════════════════════════════════════════════════════════════════════════
def _crop_save(
    src: Path,
    box: tuple[float, float, float, float],
    out_path: Path,
    border: tuple[int, int, int] = PRIMARY,
) -> Path:
    img = Image.open(src).convert("RGB")
    iw, ih = img.size
    cropped = img.crop(
        (
            int(iw * box[0]),
            int(ih * box[1]),
            int(iw * box[2]),
            int(ih * box[3]),
        )
    )
    canvas = Image.new("RGB", (1600, 900), BG_TOP)
    tw, th = 1560, 860
    ratio = min(tw / cropped.width, th / cropped.height)
    nw = max(1, int(cropped.width * ratio))
    nh = max(1, int(cropped.height * ratio))
    resized = cropped.resize((nw, nh), Image.LANCZOS)  # type: ignore[attr-defined]
    canvas.paste(resized, (20 + (tw - nw) // 2, 20 + (th - nh) // 2))
    draw = ImageDraw.Draw(canvas)
    draw.rounded_rectangle((10, 10, 1590, 890), radius=16, outline=border, width=4)
    canvas.save(out_path)
    return out_path


def generate_uml_crops(crops_dir: Path) -> dict[str, Path]:
    crops: dict[str, Path] = {}
    crops_dir.mkdir(parents=True, exist_ok=True)

    if FRAMEWORK_IMG.exists():
        crops["framework_full"] = FRAMEWORK_IMG
        crops["framework_ui"] = _crop_save(
            FRAMEWORK_IMG, (0.01, 0.01, 0.40, 0.55), crops_dir / "fw_ui.png", PRIMARY
        )
        crops["framework_dao"] = _crop_save(
            FRAMEWORK_IMG, (0.26, 0.30, 0.64, 0.78), crops_dir / "fw_dao.png", GREEN
        )
        crops["framework_right"] = _crop_save(
            FRAMEWORK_IMG, (0.58, 0.20, 0.95, 0.66), crops_dir / "fw_right.png", ORANGE
        )

    if CLASS_IMG.exists():
        crops["class_full"] = CLASS_IMG
        crops["class_ui"] = _crop_save(
            CLASS_IMG, (0.01, 0.02, 0.42, 0.58), crops_dir / "cls_ui.png", PRIMARY
        )
        crops["class_dao"] = _crop_save(
            CLASS_IMG, (0.44, 0.22, 0.76, 0.64), crops_dir / "cls_dao.png", GREEN
        )

    if SEQUENCE_IMG.exists():
        crops["sequence_full"] = SEQUENCE_IMG
        crops["sequence_enter"] = _crop_save(
            SEQUENCE_IMG, (0.04, 0.00, 0.96, 0.36), crops_dir / "seq_enter.png", PRIMARY
        )
        crops["sequence_detail"] = _crop_save(
            SEQUENCE_IMG, (0.04, 0.24, 0.96, 0.62), crops_dir / "seq_detail.png", GREEN
        )
        crops["sequence_action"] = _crop_save(
            SEQUENCE_IMG, (0.04, 0.56, 0.96, 0.92), crops_dir / "seq_action.png", ORANGE
        )

    if PROCESS_IMG.exists():
        crops["process_full"] = PROCESS_IMG
        crops["process_early"] = _crop_save(
            PROCESS_IMG, (0.02, 0.06, 0.34, 0.46), crops_dir / "proc_early.png", PRIMARY
        )
        crops["process_mid"] = _crop_save(
            PROCESS_IMG, (0.30, 0.28, 0.68, 0.68), crops_dir / "proc_mid.png", GREEN
        )
        crops["process_late"] = _crop_save(
            PROCESS_IMG, (0.60, 0.50, 0.97, 0.92), crops_dir / "proc_late.png", PURPLE
        )

    return crops


# ═══════════════════════════════════════════════════════════════════════════
# MOCK SCREENSHOT BUILDER
# ═══════════════════════════════════════════════════════════════════════════
def mock_page(
    title: str,
    subtitle: str,
    tags: list[str],
    cards: list[tuple[str, str, tuple[int, int, int]]],
    out_path: Path,
) -> Path:
    img = Image.new("RGB", (1440, 900), BG_TOP)
    gradient_fill(img)
    draw = ImageDraw.Draw(img)

    # header bar
    draw.rectangle((0, 0, 1440, 88), fill=TITLE_BG)
    draw.line((0, 88, 1440, 88), fill=ACCENT, width=3)
    draw.text((40, 10), title, font=_font(36, True), fill=ACCENT)
    draw.text((42, 56), subtitle, font=_font(16), fill=TEXT_MUTE)

    # tag chips
    x = 40
    for tag in tags:
        tw, _ = tsz(draw, tag, _font(16, True))
        draw.rounded_rectangle(
            (x, 102, x + tw + 26, 132),
            radius=10,
            fill=(36, 58, 102),
            outline=BORDER,
            width=1,
        )
        draw.text((x + 13, 105), tag, font=_font(16, True), fill=PRIMARY)
        x += tw + 38

    # content cards
    y = 152
    for ctitle, cdesc, color in cards:
        draw.rounded_rectangle(
            (40, y, 1400, y + 126),
            radius=14,
            fill=(28, 46, 78),
            outline=BORDER,
            width=2,
        )
        draw.rounded_rectangle(
            (58, y + 16, 86, y + 44), radius=8, fill=color, outline=color, width=1
        )
        draw.text((98, y + 10), ctitle, font=_font(22, True), fill=TEXT)
        draw_wrapped(draw, cdesc, 98, y + 44, _font(18), TEXT_MUTE, 1270, 4)
        y += 146

    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)
    return out_path


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE RENDERING
# ═══════════════════════════════════════════════════════════════════════════
def render_base_slide(
    slide: SlideSpec, idx: int, total: int, out_path: Path
) -> Image.Image:
    """
    Renders full slide except subtitle bar (left blank / SUB_BG).
    Saves PNG to out_path and returns the Image for subtitle overlay.
    """
    img = Image.new("RGB", (W, H), BG_TOP)
    gradient_fill(img)
    draw = ImageDraw.Draw(img)

    # ── title bar ─────────────────────────────────────────────────────────
    draw.rectangle((0, 0, W, TITLE_H), fill=TITLE_BG)
    draw.line((0, TITLE_H, W, TITLE_H), fill=ACCENT, width=3)
    draw.text((48, 12), slide.title, font=FT, fill=ACCENT)
    draw.text((50, 58), slide.subtitle, font=FS, fill=TEXT_MUTE)

    badge = f"{idx} / {total}"
    bw, _ = tsz(draw, badge, FKB)
    bx = W - bw - 50
    draw.rounded_rectangle(
        (bx - 12, 20, W - 20, 64), radius=10, fill=PRIMARY, outline=PRIMARY, width=1
    )
    draw.text((bx, 28), badge, font=FKB, fill=(8, 16, 36))

    # ── layout ────────────────────────────────────────────────────────────
    LEFT_W = 570
    MG = 42
    RSTART = LEFT_W + MG * 2
    CTOP = TITLE_H + 18
    CBOT = CONTENT_H - 8

    has_code = bool(slide.code_title and slide.code_text)
    IMG_BOT = CBOT - 148 if has_code else CBOT

    # left panel
    draw.rounded_rectangle(
        (MG, CTOP, LEFT_W + MG, CBOT), radius=16, fill=PANEL, outline=BORDER, width=2
    )
    draw_bullets(draw, slide.bullets, MG + 18, CTOP + 18, LEFT_W + MG - 10, CBOT - 10)

    # right panel (image)
    draw.rounded_rectangle(
        (RSTART, CTOP, W - MG, IMG_BOT), radius=16, fill=PANEL, outline=BORDER, width=2
    )
    paste_fit(img, slide.image_path, (RSTART + 8, CTOP + 8, W - MG - 8, IMG_BOT - 28))
    if slide.image_caption:
        draw.text(
            (RSTART + 12, IMG_BOT - 24),
            slide.image_caption[:92],
            font=FK,
            fill=TEXT_MUTE,
        )

    # code / summary box
    if has_code:
        draw_code_box(
            draw,
            slide.code_title,  # type: ignore[arg-type]
            slide.code_text,  # type: ignore[arg-type]
            RSTART,
            IMG_BOT + 8,
            W - MG,
            CBOT,
        )

    # subtitle placeholder
    draw.rectangle((0, CONTENT_H, W, H), fill=SUB_BG)
    draw.line((0, CONTENT_H, W, CONTENT_H), fill=(40, 44, 68), width=1)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path, "PNG")
    return img


def add_subtitle_bar(base: Image.Image, sentence: str, out_path: Path) -> Path:
    """Overlay sentence text into subtitle bar area; save to out_path."""
    result = base.copy()
    draw = ImageDraw.Draw(result)

    draw.rectangle((0, CONTENT_H, W, H), fill=SUB_BG)
    draw.line((0, CONTENT_H, W, CONTENT_H), fill=(48, 52, 78), width=1)

    if sentence:
        max_w = W - 180
        lines = wrap_text(draw, sentence, FSU, max_w)[:3]
        lh = tsz(draw, "M", FSU)[1] + 7
        total = len(lines) * lh
        y = CONTENT_H + (SUB_H - total) // 2
        for line in lines:
            lw, _ = tsz(draw, line, FSU)
            draw.text(((W - lw) // 2, y), line, font=FSU, fill=WHITE)
            y += lh

    out_path.parent.mkdir(parents=True, exist_ok=True)
    result.save(out_path, "PNG")
    return out_path


# ═══════════════════════════════════════════════════════════════════════════
# SRT
# ═══════════════════════════════════════════════════════════════════════════
def ts_fmt(s: float) -> str:
    h = int(s // 3600)
    m = int((s % 3600) // 60)
    sc = int(s % 60)
    ms = int(round((s - int(s)) * 1000))
    return f"{h:02d}:{m:02d}:{sc:02d},{ms:03d}"


def write_srt(entries: list[tuple[str, float, float]], srt_path: Path) -> None:
    lines: list[str] = []
    for i, (text, start, end) in enumerate(entries, 1):
        lines += [str(i), f"{ts_fmt(start)} --> {ts_fmt(end)}", text, ""]
    srt_path.parent.mkdir(parents=True, exist_ok=True)
    srt_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"SRT  → {srt_path}  ({len(entries)} entries)")


# ═══════════════════════════════════════════════════════════════════════════
# PPTX
# ═══════════════════════════════════════════════════════════════════════════
def build_pptx(
    slides: list[SlideSpec],
    base_paths: list[Path],
    pptx_path: Path,
) -> None:
    if not _PPTX_OK or Presentation is None:
        print("python-pptx not available, skipping PPTX.")
        return
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for img_path, spec in zip(base_paths, slides):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        tb = sl.shapes.add_textbox(
            Inches(0.38), Inches(0.12), Inches(10.5), Inches(0.52)
        )
        p = tb.text_frame.paragraphs[0]
        p.text = spec.title
        p.font.size = Pt(28)
        p.font.bold = True
        tb2 = sl.shapes.add_textbox(
            Inches(0.40), Inches(0.62), Inches(12.2), Inches(0.34)
        )
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = spec.subtitle
        p2.font.size = Pt(13)
        sl.shapes.add_picture(
            str(img_path),
            Inches(0.28),
            Inches(1.02),
            width=Inches(12.78),
            height=Inches(6.18),
        )
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    print(f"PPTX → {pptx_path}")


# ═══════════════════════════════════════════════════════════════════════════
# SCRIPT MARKDOWN
# ═══════════════════════════════════════════════════════════════════════════
def build_script(
    slides: list[SlideSpec],
    script_path: Path,
    feature_name: str,
) -> None:
    lines = [
        f"# {feature_name}教学视频脚本 v6",
        "",
        "## 技术特性",
        "",
        "- **edge_tts** (zh-CN-XiaoxiaoNeural) 高质量中文语音",
        "- **moviepy** 精确帧级同步：每句独立帧，一次性编码，彻底消除漂移",
        "- 字幕逐句烧录到画面；SRT 时间戳来自真实 MP3 时长",
        "",
        "## 成片结构",
        "",
    ]
    for i, s in enumerate(slides, 1):
        lines.append(f"{i}. **{s.title}** — {s.subtitle}")
    lines += ["", "## 讲解稿", ""]
    for i, s in enumerate(slides, 1):
        lines += [
            f"### {i}. {s.title}",
            f"> {s.subtitle}",
            "",
            f"**旁白：** {s.narration}",
            "",
        ]
    script_path.parent.mkdir(parents=True, exist_ok=True)
    script_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"Script → {script_path}")


# ═══════════════════════════════════════════════════════════════════════════
# MAIN VIDEO PIPELINE
# ═══════════════════════════════════════════════════════════════════════════
def build_video(
    slides: list[SlideSpec],
    work_dir: Path,
    video_path: Path,
    srt_path: Path,
) -> tuple[list[Path], bool]:
    """
    Core sync pipeline (same as v6 knowledge-graph video):
      Per slide → per sentence:
        1. edge_tts MP3  (real duration)
        2. render base PNG (once per slide)
        3. overlay subtitle → per-sentence PNG
        4. moviepy: ImageClip.set_duration(real_dur + tail).set_audio(AudioFileClip)
      concatenate_videoclips → write_videofile (single AAC encode, zero drift)
      SRT timestamps = cumulative real durations.
    """
    if not _MOVIEPY_OK:
        print("ERROR: moviepy required.  pip install moviepy==1.0.3")
        return [], False

    slides_dir = work_dir / "slides"
    sent_dir = work_dir / "sent"
    audio_dir = work_dir / "audio"
    temp_dir = work_dir / "temp"
    for d in [slides_dir, sent_dir, audio_dir, temp_dir]:
        d.mkdir(parents=True, exist_ok=True)

    all_clips: list[Any] = []
    srt_entries: list[tuple[str, float, float]] = []
    base_paths: list[Path] = []
    t = 0.0
    total = len(slides)

    for si, slide in enumerate(slides, 1):
        print(f"\n[{si:02d}/{total}] {slide.title}")

        base_path = slides_dir / f"{si:02d}_{safe_name(slide.title)}.png"
        base_img = render_base_slide(slide, si, total, base_path)
        base_paths.append(base_path)

        sentences = get_sentences(slide)
        nsent = len(sentences)

        for ki, sentence in enumerate(sentences):
            is_last = ki == nsent - 1
            tail = SLIDE_TAIL if is_last else CLIP_TAIL
            stem = audio_dir / f"{si:02d}_{ki + 1:02d}"

            print(f"  [{ki + 1}/{nsent}] TTS → {sentence[:58]}")

            audio_path = generate_tts(sentence, stem)

            if audio_path is None or not audio_path.exists():
                audio_path = stem.with_suffix(".wav")
                make_silent_wav(max(1.8, len(sentence) * 0.14), audio_path)

            real_dur = get_audio_duration(audio_path)
            if real_dur < 0.05:
                real_dur = max(1.8, len(sentence) * 0.14)

            frame_path = sent_dir / f"{si:02d}_{ki + 1:02d}.png"
            add_subtitle_bar(base_img, sentence, frame_path)

            clip_dur = real_dur + tail
            try:
                audio_clip = AudioFileClip(str(audio_path))
                vid_clip = (
                    ImageClip(str(frame_path))
                    .set_duration(clip_dur)
                    .set_audio(audio_clip)
                )
            except Exception as e:
                print(f"    clip error: {e}  (silent fallback)")
                vid_clip = ImageClip(str(frame_path)).set_duration(clip_dur)

            all_clips.append(vid_clip)
            srt_entries.append((sentence, t, t + real_dur))
            t += clip_dur

    write_srt(srt_entries, srt_path)

    print(f"\nConcatenating {len(all_clips)} clips  (≈ {t:.1f} s)...")
    try:
        final = concatenate_videoclips(all_clips, method="compose")
        print(f"Encoding → {video_path}")
        video_path.parent.mkdir(parents=True, exist_ok=True)
        final.write_videofile(
            str(video_path),
            fps=FPS,
            codec="libx264",
            audio_codec="aac",
            audio_bitrate="192k",
            temp_audiofile=str(temp_dir / "tmp_audio.m4a"),
            remove_temp=True,
            verbose=True,
            write_logfile=False,
            threads=4,
        )
        final.close()
        print(f"Video → {video_path}")
        return base_paths, True
    except Exception as e:
        print(f"ERROR during video write: {e}")
        import traceback

        traceback.print_exc()
        return base_paths, False


# ═══════════════════════════════════════════════════════════════════════════
# ASSET DOWNLOAD
# ═══════════════════════════════════════════════════════════════════════════
def download_asset(url: str, out_path: Path) -> Path | None:
    if out_path.exists() and out_path.stat().st_size > 500:
        return out_path
    try:
        urllib.request.urlretrieve(url, str(out_path))
        return out_path if out_path.exists() and out_path.stat().st_size > 500 else None
    except Exception as e:
        print(f"  Download failed ({out_path.name}): {e}")
        return None


def ensure_official_assets(assets_dir: Path) -> dict[str, Path]:
    assets_dir.mkdir(parents=True, exist_ok=True)
    result: dict[str, Path] = {}
    f = download_asset(
        "https://flutter.dev/assets/lockup_flutter_horizontal.d0515092173211776ceed19b39c2a041.png",
        assets_dir / "flutter_official.png",
    )
    if f:
        result["flutter"] = f
    d = download_asset(
        "https://dart.dev/assets/img/logo/dart-logo-for-shares.png",
        assets_dir / "dart_official.png",
    )
    if d:
        result["dart"] = d
    return result
